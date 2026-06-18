import time
import threading
import queue
import wave
import re
import os
import json
import sys
import glob
import tempfile
import shutil
import random
import hashlib
import base64
import subprocess
import traceback
import requests
import platform
from datetime import datetime, timedelta
from email.utils import parsedate_to_datetime
from io import BytesIO
from pathlib import Path
from urllib.parse import parse_qs, parse_qsl, urlencode, urlparse, quote, urlunparse

from portable_runtime import configure_portable_environment, get_models_dir


configure_portable_environment()

def _configure_ffmpeg_path() -> str | None:
    """跨平台注册 imageio-ffmpeg 提供的二进制，兼容 Windows 和 Linux 部署环境。"""
    bundled_ffmpeg = str(os.environ.get("IMAGEIO_FFMPEG_EXE", "") or "").strip()
    if bundled_ffmpeg and os.path.exists(bundled_ffmpeg):
        ffmpeg_dir = os.path.dirname(bundled_ffmpeg)
        if ffmpeg_dir not in os.environ.get("PATH", ""):
            os.environ["PATH"] = ffmpeg_dir + os.pathsep + os.environ.get("PATH", "")
            print(f"Added bundled ffmpeg to PATH: {ffmpeg_dir}")
        return bundled_ffmpeg

    try:
        import imageio_ffmpeg

        ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
        if not ffmpeg_path:
            return None

        ffmpeg_dir = os.path.dirname(ffmpeg_path)
        ffmpeg_name = "ffmpeg.exe" if os.name == "nt" else "ffmpeg"
        target_ffmpeg = os.path.join(ffmpeg_dir, ffmpeg_name)

        if not os.path.exists(target_ffmpeg):
            try:
                shutil.copy2(ffmpeg_path, target_ffmpeg)
                if os.name != "nt":
                    os.chmod(target_ffmpeg, 0o755)
            except Exception as e:
                print(f"Failed to create ffmpeg alias {target_ffmpeg}: {e}")

        if ffmpeg_dir not in os.environ.get("PATH", ""):
            os.environ["PATH"] = ffmpeg_dir + os.pathsep + os.environ.get("PATH", "")
            print(f"Added ffmpeg to PATH: {ffmpeg_dir}")

        # 直接返回可执行文件路径，避免 Linux 环境误用 ffmpeg.exe 导致 yt-dlp 识别失败。
        return ffmpeg_path
    except ImportError:
        print("imageio-ffmpeg not found. Assuming ffmpeg is in PATH.")
        return None
    except Exception as e:
        print(f"Failed to auto-configure ffmpeg: {e}")
        return None


# 自动配置 ffmpeg 环境 (从 imageio-ffmpeg 获取)
ffmpeg_binary_path = _configure_ffmpeg_path()

from requests.adapters import HTTPAdapter
from urllib3 import Retry
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api._errors import (
    AgeRestricted,
    InvalidVideoId,
    IpBlocked,
    NoTranscriptFound,
    PoTokenRequired,
    RequestBlocked,
    TranscriptsDisabled,
    VideoUnavailable,
    VideoUnplayable,
)

# 自动配置 nodejs 环境
try:
    # 自动检测并添加 node.exe 到 PATH (仅本地 Windows 调试时可能用到)
    possible_node_paths = [
        r"d:\Program Files",
        r"D:\Program Files",
        r"C:\Program Files\nodejs",
        r"C:\Program Files (x86)\nodejs",
    ]
    for p in possible_node_paths:
        node_exe = os.path.join(p, "node.exe")
        if os.path.exists(node_exe):
            if p not in os.environ.get("PATH", ""):
                os.environ["PATH"] = p + os.pathsep + os.environ.get("PATH", "")
                print(f"已自动添加 nodejs 到 PATH: {p}")
            break

except ImportError:
    print("未安装 imageio-ffmpeg，无法自动配置 ffmpeg 环境。")
except Exception as e:
    print(f"自动配置环境失败: {e}")


def resolve_cookie_file(
    cookies_file: str = "",
    cookies_content: str = "",
    cookies_content_b64: str = "",
) -> str:
    """
    解析可供 yt-dlp 使用的 cookies 文件路径。

    支持三种输入方式：
    1. 直接传入现成的 cookies 文件路径
    2. 传入 Netscape cookies 文本内容
    3. 传入 Base64 编码后的 cookies 文本
    """
    file_path = (cookies_file or "").strip()
    if file_path:
        return file_path

    raw_content = (cookies_content or "").strip()
    if not raw_content and (cookies_content_b64 or "").strip():
        try:
            raw_content = base64.b64decode(cookies_content_b64).decode("utf-8")
        except Exception as e:
            raise RuntimeError(f"Cookies Base64 内容无效：{e}") from e

    if not raw_content:
        return ""

    runtime_dir = Path(tempfile.gettempdir()) / "youtube_summarizer_cookies"
    runtime_dir.mkdir(parents=True, exist_ok=True)
    content_hash = hashlib.sha256(raw_content.encode("utf-8")).hexdigest()[:16]
    cookie_path = runtime_dir / f"cookies_{content_hash}.txt"
    if not cookie_path.exists():
        cookie_path.write_text(raw_content, encoding="utf-8")
    return str(cookie_path)


def build_cookie_runtime_diagnostics(
    api=None,
    *,
    cookies_file: str = "",
    cookies_from_browser: str = "",
    resolved_cookie_file: str = "",
    resolve_error: str = "",
) -> str:
    """
    生成 cookies 运行时诊断信息。

    只输出来源、长度、文件存在性等元数据，不泄露真实 cookie 内容。
    """
    raw_file = str(cookies_file or getattr(api, "_cookies_file", "") or "").strip()
    raw_content = str(getattr(api, "_cookies_content", "") or "").strip()
    raw_b64 = str(getattr(api, "_cookies_content_b64", "") or "").strip()
    browser = str(cookies_from_browser or getattr(api, "_cookies_from_browser", "") or "").strip().lower()

    parts = [
        f"input_file={'yes' if raw_file else 'no'}",
        f"input_content={'yes' if raw_content else 'no'}",
        f"input_b64={'yes' if raw_b64 else 'no'}",
        f"browser={browser or 'none'}",
    ]
    if raw_b64:
        parts.append(f"b64_len={len(raw_b64)}")
    if resolve_error:
        parts.append(f"resolve_error={resolve_error}")

    resolved_path = str(resolved_cookie_file or "").strip()
    if resolved_path:
        try:
            cookie_path = Path(resolved_path)
            exists = cookie_path.exists()
            parts.append(f"resolved_file={'yes' if exists else 'no'}")
            parts.append(f"resolved_name={cookie_path.name}")
            if exists:
                parts.append(f"resolved_size={cookie_path.stat().st_size}")
                first_line = cookie_path.read_text(encoding="utf-8", errors="ignore").splitlines()
                first_line_text = first_line[0].strip() if first_line else ""
                parts.append(f"netscape_header={'yes' if first_line_text.startswith('# Netscape HTTP Cookie File') else 'no'}")
        except Exception as e:
            parts.append(f"resolved_check_error={type(e).__name__}:{e}")
    else:
        parts.append("resolved_file=no")

    return "; ".join(parts)


def build_runtime_version_diagnostics() -> str:
    """输出当前运行实例的版本信息，便于确认 Render 是否部署到最新提交。"""
    render_commit = str(os.environ.get("RENDER_GIT_COMMIT", "") or "").strip()
    render_service = str(os.environ.get("RENDER_SERVICE_ID", "") or "").strip()
    render_instance = str(os.environ.get("RENDER_INSTANCE_ID", "") or "").strip()
    app_expected_commit = "latest-local"
    parts = [f"expected_commit={app_expected_commit}"]
    parts.append(f"render_commit={render_commit or 'unknown'}")
    parts.append(f"render_service={'set' if render_service else 'unknown'}")
    parts.append(f"render_instance={'set' if render_instance else 'unknown'}")
    return "; ".join(parts)


def get_remote_worker_status(timeout_seconds: float = 4.0) -> dict:
    """检查远程抓取节点配置、健康状态和 Render ASR 兜底策略。"""
    remote_enabled = str(
        os.environ.get("REMOTE_TRANSCRIBE_ENABLED", "0") or "0"
    ).strip().lower() in {"1", "true", "yes"}
    worker_url = str(os.environ.get("REMOTE_TRANSCRIBE_URL", "") or "").strip()
    remote_mode = str(os.environ.get("REMOTE_TRANSCRIBE_MODE", "") or "").strip().lower()
    worker_token = str(os.environ.get("REMOTE_TRANSCRIBE_TOKEN", "") or "").strip()
    running_on_render = bool(str(os.environ.get("RENDER_SERVICE_ID", "") or "").strip())
    disable_render_asr_fallback = running_on_render and str(
        os.environ.get("REMOTE_TRANSCRIBE_DISABLE_RENDER_ASR_FALLBACK", "0") or "0"
    ).strip().lower() not in {"0", "false", "no"}

    status = {
        "configured": remote_enabled and bool(worker_url),
        "remote_enabled": remote_enabled,
        "remote_mode": remote_mode if remote_enabled else "disabled",
        "worker_url": worker_url,
        "worker_host": "",
        "worker_health_url": "",
        "worker_token_configured": bool(worker_token),
        "running_on_render": running_on_render,
        "disable_render_asr_fallback": disable_render_asr_fallback,
        "health_ok": False,
        "health_status_code": None,
        "health_error": "",
        "health_payload": {},
    }

    if not remote_enabled:
        status["health_error"] = "REMOTE_TRANSCRIBE_ENABLED=0"
        return status

    if not worker_url:
        status["health_error"] = "未配置 REMOTE_TRANSCRIBE_URL"
        return status

    parsed = urlparse(worker_url)
    base = worker_url.rsplit("/fetch-transcript", 1)[0] if "/fetch-transcript" in worker_url else worker_url.rstrip("/")
    health_url = f"{base}/health"
    status["worker_host"] = parsed.netloc or worker_url
    status["worker_health_url"] = health_url

    try:
        resp = requests.get(health_url, timeout=max(1.0, float(timeout_seconds)), verify=False)
        status["health_status_code"] = resp.status_code
        if resp.ok:
            try:
                payload = resp.json()
            except Exception:
                payload = {"raw_text": resp.text[:500]}
            status["health_payload"] = payload if isinstance(payload, dict) else {"raw": payload}
            status["health_ok"] = bool(status["health_payload"].get("ok", True))
            if not status["health_ok"]:
                status["health_error"] = str(status["health_payload"].get("error") or "健康检查返回 ok=false")
        else:
            status["health_error"] = f"HTTP {resp.status_code}"
    except requests.exceptions.RequestException as e:
        status["health_error"] = f"{type(e).__name__}: {e}"

    return status


def try_fetch_transcript_via_remote_worker(
    video_id: str,
    video_url: str,
    languages: list[str],
    api=None,
) -> str:
    """
    通过外部抓取节点拉取 transcript 文本。

    仅返回 transcript 文本，不负责总结。
    """
    remote_enabled = str(
        os.environ.get("REMOTE_TRANSCRIBE_ENABLED", "0") or "0"
    ).strip().lower() in {"1", "true", "yes"}
    if not remote_enabled:
        raise RuntimeError("远程抓取节点已禁用（REMOTE_TRANSCRIBE_ENABLED=0）")

    worker_url = str(os.environ.get("REMOTE_TRANSCRIBE_URL", "") or "").strip()
    worker_token = str(os.environ.get("REMOTE_TRANSCRIBE_TOKEN", "") or "").strip()
    if not worker_url:
        raise RuntimeError("未配置 REMOTE_TRANSCRIBE_URL")

    timeout_seconds = float(getattr(api, "_timeout_seconds", 60.0) or 60.0) if api else 60.0
    worker_timeout_seconds = float(os.environ.get("REMOTE_TRANSCRIBE_TIMEOUT_SECONDS", "180") or "180")
    processing_extension_seconds = float(os.environ.get("REMOTE_TRANSCRIBE_PROCESSING_EXTENSION_SECONDS", "240") or "240")
    poll_interval_seconds = float(os.environ.get("REMOTE_TRANSCRIBE_POLL_INTERVAL_SECONDS", "2.0") or "2.0")
    payload = {
        "video_id": video_id,
        "video_url": video_url,
        "languages": list(languages or []),
        "cookies_file": str(getattr(api, "_cookies_file", "") or "").strip() if api else "",
        "cookies_content": str(getattr(api, "_cookies_content", "") or "") if api else "",
        "cookies_content_b64": str(getattr(api, "_cookies_content_b64", "") or "").strip() if api else "",
        "cookies_from_browser": str(getattr(api, "_cookies_from_browser", "") or "").strip().lower() if api else "",
        "asr_enabled": bool(getattr(api, "_asr_enabled", False)) if api else False,
        "asr_model": str(getattr(api, "_asr_model", "") or "") if api else "",
        "asr_language": str(getattr(api, "_asr_language", "") or "") if api else "",
        "asr_fast_mode": bool(getattr(api, "_asr_fast_mode", False)) if api else False,
        "asr_force_cpu": bool(getattr(api, "_asr_force_cpu", False)) if api else False,
    }
    headers = {"Content-Type": "application/json"}
    if worker_token:
        headers["X-Worker-Token"] = worker_token

    def _raise_remote_connectivity_error(exc: Exception) -> None:
        host = urlparse(worker_url).netloc or worker_url
        detail = f"{type(exc).__name__}: {exc}"
        if "trycloudflare.com" in host:
            raise RuntimeError(
                "本地抓取节点不可达：当前配置的是临时 Cloudflare Tunnel 地址，且该地址已失效或无法解析。"
                f" host={host}；detail={detail}。"
                " 请在本地重启 `cloudflared tunnel --url http://127.0.0.1:8787`，"
                "拿到新的 trycloudflare 地址后，立即更新 Render 的 `REMOTE_TRANSCRIBE_URL`。"
                " 如果需要给朋友持续测试，建议改用固定域名的 Cloudflare Tunnel、Tailscale Funnel 或 ngrok。"
            )
        raise RuntimeError(
            f"本地抓取节点不可达：host={host}；detail={detail}。"
            " 请检查本地抓取节点、隧道进程以及 Render 中的 `REMOTE_TRANSCRIBE_URL` 是否仍然有效。"
        )

    try:
        submit_resp = requests.post(
            worker_url,
            headers=headers,
            json=payload,
            timeout=(10.0, 20.0),
        )
    except requests.exceptions.ConnectionError as e:
        _raise_remote_connectivity_error(e)
    except requests.exceptions.RequestException as e:
        raise RuntimeError(f"提交本地抓取任务失败：{type(e).__name__}: {e}")
    if submit_resp.status_code == 401:
        raise RuntimeError(
            "本地抓取节点拒绝了请求（401 Unauthorized）。"
            " 这通常表示 Render 端的 `REMOTE_TRANSCRIBE_TOKEN` 与本地抓取节点的 `REMOTE_TRANSCRIBE_TOKEN` 不一致，"
            "或本地节点要求 token 但 Render 未正确携带。"
        )
    submit_resp.raise_for_status()
    data = submit_resp.json()
    if not isinstance(data, dict):
        raise RuntimeError("本地抓取节点返回格式无效")
    if not data.get("ok"):
        raise RuntimeError(str(data.get("error") or "本地抓取节点执行失败"))
    task_id = str(data.get("task_id") or "").strip()
    if not task_id:
        transcript_text = str(data.get("transcript_text") or "").strip()
        transcript_label = str(data.get("transcript_label") or "remote-worker").strip()
        if not transcript_text:
            raise RuntimeError("本地抓取节点未返回 transcript 文本")
        return f"[{transcript_label}]\n\n{transcript_text}"

    status_base = worker_url.rsplit("/fetch-transcript", 1)[0]
    status_url = f"{status_base}/task/{quote(task_id)}"
    deadline = time.monotonic() + max(30.0, timeout_seconds, worker_timeout_seconds)
    last_stage = ""
    last_stage_detail = ""
    last_status = ""
    processing_deadline_extended = False

    while time.monotonic() < deadline:
        time.sleep(max(0.5, poll_interval_seconds))
        try:
            poll_resp = requests.get(status_url, headers=headers, timeout=(10.0, 20.0))
        except requests.exceptions.ConnectionError as e:
            _raise_remote_connectivity_error(e)
        except requests.exceptions.RequestException as e:
            raise RuntimeError(f"查询本地抓取节点状态失败：{type(e).__name__}: {e}")
        if poll_resp.status_code == 401:
            raise RuntimeError(
                "查询本地抓取节点状态被拒绝（401 Unauthorized）。"
                " 请检查 Render 端与本地节点的 `REMOTE_TRANSCRIBE_TOKEN` 是否完全一致。"
            )
        poll_resp.raise_for_status()
        poll_data = poll_resp.json()
        if not isinstance(poll_data, dict) or not poll_data.get("ok"):
            raise RuntimeError("本地抓取节点状态查询失败")
        status = str(poll_data.get("status") or "").strip().lower()
        last_status = status
        last_stage = str(poll_data.get("stage") or "").strip()
        last_stage_detail = str(poll_data.get("stage_detail") or "").strip()
        if (
            not processing_deadline_extended
            and status in {"queued", "running"}
            and last_stage == "processing"
            and any(token in last_stage_detail.lower() for token in ["whisper", "转写", "audio", "音频"])
        ):
            deadline = max(deadline, time.monotonic() + max(60.0, processing_extension_seconds))
            processing_deadline_extended = True
        if status in {"queued", "running"}:
            continue
        if status == "failed":
            raise RuntimeError(str(poll_data.get("error") or "本地抓取节点执行失败"))
        if status == "success":
            transcript_text = str(poll_data.get("transcript_text") or "").strip()
            transcript_label = str(poll_data.get("transcript_label") or "remote-worker").strip()
            if not transcript_text:
                raise RuntimeError("本地抓取节点未返回 transcript 文本")
            return f"[{transcript_label}]\n\n{transcript_text}"
        raise RuntimeError(f"本地抓取节点返回未知状态: {status or 'unknown'}")

    timeout_hint = f"本地抓取节点处理超时（>{int(max(30.0, timeout_seconds, worker_timeout_seconds))}s）"
    if last_status or last_stage or last_stage_detail:
        timeout_hint += (
            f"；last_status={last_status or 'unknown'}"
            f"；stage={last_stage or 'unknown'}"
            f"；detail={last_stage_detail or 'none'}"
        )
    raise RuntimeError(timeout_hint)


class TimeoutSession(requests.Session):
    def __init__(self, timeout_seconds: float):
        super().__init__()
        self.verify = False  # 全局禁用 SSL 验证，解决代理证书问题
        self._timeout_seconds = timeout_seconds
        self._exception_retries = 2
        self._deadline_s: float | None = None
        self.headers.update(
            {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                "Accept-Encoding": "identity",
            }
        )

    def request(self, method, url, **kwargs):
        if self._deadline_s is not None and time.monotonic() > self._deadline_s:
            raise requests.exceptions.Timeout("已超过整体超时时间")
        if "timeout" not in kwargs:
            kwargs["timeout"] = (min(10.0, self._timeout_seconds), self._timeout_seconds)
        last_exc: Exception | None = None
        retries = max(0, int(getattr(self, "_exception_retries", 0)))
        for attempt in range(retries + 1):
            try:
                return super().request(method, url, **kwargs)
            except (requests.exceptions.ChunkedEncodingError, requests.exceptions.ConnectionError, requests.exceptions.ReadTimeout) as e:
                last_exc = e
                if attempt >= retries:
                    break
                time.sleep(0.6 * (2**attempt))
        assert last_exc is not None
        raise last_exc


def _strip_trailing_punct(text: str) -> str:
    return re.sub(r"[!！。,.，?？;；:：\)\]）】]+$", "", text or "")


def extract_video_id(url_or_id: str) -> str:
    candidate = _strip_trailing_punct(url_or_id.strip())
    
    # YouTube ID (11 chars)
    if re.fullmatch(r"[A-Za-z0-9_-]{11}", candidate):
        return candidate

    if not re.match(r"^[a-zA-Z][a-zA-Z0-9+.-]*://", candidate) and (
        candidate.startswith("www.") or "youtube.com" in candidate or "youtu.be" in candidate
    ):
        candidate = "https://" + candidate

    parsed = urlparse(candidate)
    host = (parsed.netloc or "").lower()

    if "youtu.be" in host:
        video_id = _strip_trailing_punct(parsed.path.strip("/").split("/")[0])
        if re.fullmatch(r"[A-Za-z0-9_-]{11}", video_id):
            return video_id

    if "youtube.com" in host:
        qs = parse_qs(parsed.query)
        if "v" in qs and qs["v"]:
            video_id = _strip_trailing_punct(qs["v"][0])
            if re.fullmatch(r"[A-Za-z0-9_-]{11}", video_id):
                return video_id

        m = re.search(r"/(shorts|embed)/([A-Za-z0-9_-]{11})", parsed.path)
        if m:
            return m.group(2)

    raise ValueError("无法从输入解析出视频 ID（支持 YouTube 11位 ID）")


def normalize_video_url(url_or_id: str) -> str:
    s = _strip_trailing_punct(url_or_id.strip())
    
    # YouTube
    if re.fullmatch(r"[A-Za-z0-9_-]{11}", s):
        return f"https://www.youtube.com/watch?v={s}"
    if not re.match(r"^[a-zA-Z][a-zA-Z0-9+.-]*://", s) and (
        s.startswith("www.") or "youtube.com" in s or "youtu.be" in s
    ):
        s = "https://" + s
    return s


class CookieManager:
    """集中管理 Cookie 的获取、验证与错误处理"""
    
    @staticmethod
    def is_cookie_error(msg: str | Exception) -> bool:
        """统一识别 yt-dlp 报告的 Cookie 相关错误"""
        s = strip_ansi(str(msg or "")).lower()
        # 使用正则进行更鲁棒的匹配
        patterns = [
            r"could not copy.*cookie", # 数据库被锁定
            r"cookie database",
            r"database is locked",
            r"permission denied.*cookie",
            r"access is denied",
            r"used by another process",
            r"winerror 32", # 文件占用
            r"winerror 5",  # 拒绝访问
            r"sqlite3.*locked",
            r"dpapi.*decrypt",
            r"dpapi.*failed",
            r"unable to extract.*cookie",
        ]
        return any(re.search(p, s) for p in patterns)

    @staticmethod
    def get_fatal_msg(err_msg: str, browser: str) -> str:
        """生成统一且带有具体指引的致命 Cookie 错误信息"""
        browser_name = browser or "浏览器"
        details_raw = strip_ansi(err_msg).replace("ERROR: ", "").strip()
        detail_lines = [x.strip() for x in re.split(r"[\r\n]+", details_raw) if x and x.strip()]
        detail_seen: set[str] = set()
        detail_dedup: list[str] = []
        for x in detail_lines:
            k = x.lower()
            if k in detail_seen:
                continue
            detail_seen.add(k)
            detail_dedup.append(x)
        details = " | ".join(detail_dedup) if detail_dedup else details_raw
        d_lower = details.lower()
        
        # 1. 针对新版 Chromium 的解密限制 (DPAPI)
        if "dpapi" in d_lower:
            return (
                f"❌ 浏览器 Cookie 解密失败 (DPAPI Error)。\n"
                f"原因：{browser_name} (新版基于 Chromium 127+) 引入了增强加密，外部程序已无法直接解密 Cookie。\n\n"
                f"💡 解决方法：\n"
                f"1. **使用 Firefox**：Firefox 暂不受此限制影响，请在 Firefox 中登录 YouTube 后，在设置中将浏览器改为 firefox (或 auto)。\n"
                f"2. **手动导出 Cookie**：安装浏览器插件（如 'Get cookies.txt LOCALLY'）导出 cookies.txt，并在设置中提供文件路径。\n\n"
                f"原始细节: {details}"
            )
            
        # 2. 针对数据库锁定/无法复制 (常见于浏览器未彻底关闭)
        if any(x in d_lower for x in ["locked", "another process", "could not copy", "32"]):
             return (
                f"❌ 浏览器 Cookie 数据库被锁定或无法访问。\n"
                f"原因：{browser_name} 正在运行或其数据库文件被占用。\n\n"
                f"💡 解决方法：\n"
                f"1. **彻底关闭浏览器**：请确保所有 {browser_name} 窗口已关闭（包括后台进程），然后重试。\n"
                f"2. **使用 Firefox 或导出 Cookie**：参考上述方案。\n\n"
                f"原始细节: {details}"
            )
        
        return f"❌ 无法从 {browser_name} 获取 Cookie。\n细节: {details}"

    @staticmethod
    def get_sources(cookies_file: str, cookies_from_browser: str, force_browser_cookie: bool = False) -> list[tuple[str, str]]:
        """返回 (cookiefile, browser) 元组列表，支持智能回退"""
        file_path = (cookies_file or "").strip()
        browser_name = (cookies_from_browser or "").strip().lower()
        
        sources: list[tuple[str, str]] = []
        if file_path:
            # 只要上层已提供 cookies 文件路径，就优先尝试文件模式。
            # 不在这里强依赖 exists 判断，避免云端时序/挂载差异导致错误退化为 cookie=none。
            sources.append((file_path, ""))
            return sources
        
        # 确定浏览器候选列表
        candidates = []
        if browser_name and browser_name != "auto":
            candidates = [browser_name]
            # 后备方案：如果显式指定了容易受限的 Chromium，则将 firefox 作为最后保单
            if browser_name in ["chrome", "edge", "brave", "chromium"] and "firefox" not in candidates:
                candidates.append("firefox")
        else:
            # 自动模式或默认：Firefox 优先级最高（最稳），其次是主流 Chromium
            candidates = ["firefox", "edge", "chrome", "brave", "chromium"]

        for b in candidates:
            if CookieManager.is_browser_available(b):
                sources.append(("", b))
        
        if not sources and force_browser_cookie:
            sources.append(("", "chrome"))
        if ("", "") not in sources:
            sources.append(("", ""))
            
        return sources

    @staticmethod
    def is_browser_available(browser: str) -> bool:
        """检测系统中是否安装了指定浏览器并存在 Cookie 数据库"""
        b = (browser or "").strip().lower()
        if not b: return False
        if os.name != "nt": return True # 非 Windows 暂不严格检查路径
        
        local = os.environ.get("LOCALAPPDATA", "")
        roaming = os.environ.get("APPDATA", "")
        
        def has_chromium_cookie(base_dir: str) -> bool:
            if not base_dir or not os.path.isdir(base_dir): return False
            try:
                # 检查 Default 或 Profile x 目录
                entries = os.listdir(base_dir)
                profile_dirs = [d for d in entries if d == "Default" or d.startswith("Profile")]
                if not profile_dirs: profile_dirs = ["Default"]
                for d in profile_dirs:
                    if os.path.exists(os.path.join(base_dir, d, "Cookies")): return True
                    if os.path.exists(os.path.join(base_dir, d, "Network", "Cookies")): return True
            except: pass
            return False

        if b == "chrome":
            return has_chromium_cookie(os.path.join(local, "Google/Chrome/User Data"))
        if b == "edge":
            return has_chromium_cookie(os.path.join(local, "Microsoft/Edge/User Data"))
        if b == "brave":
            return has_chromium_cookie(os.path.join(local, "BraveSoftware/Brave-Browser/User Data"))
        if b == "chromium":
            return has_chromium_cookie(os.path.join(local, "Chromium/User Data"))
        if b == "firefox":
            # 检查 Roaming 下的 Profiles
            base = os.path.join(roaming, "Mozilla/Firefox/Profiles")
            if os.path.isdir(base):
                try:
                    for name in os.listdir(base):
                        if os.path.exists(os.path.join(base, name, "cookies.sqlite")): return True
                except: pass
            return False
        return False




def expand_languages(languages: list[str]) -> list[str]:
    out: list[str] = []
    seen: set[str] = set()

    def add(x: str) -> None:
        x = x.strip()
        if not x or x in seen:
            return
        seen.add(x)
        out.append(x)

    for lang in languages:
        add(lang)
        if lang == "zh-Hans":
            add("zh-Hant")
            add("zh")
            add("zh-CN")
            add("zh-TW")
        elif lang == "zh-Hant":
            add("zh-Hans")
            add("zh")
            add("zh-TW")
            add("zh-CN")
        elif lang == "zh":
            add("zh-Hans")
            add("zh-Hant")
            add("zh-CN")
            add("zh-TW")
        elif lang.lower() == "en":
            add("en-US")
            add("en-GB")

    return out


def detect_windows_proxy() -> tuple[str, str]:
    # 1. 尝试从注册表读取
    try:
        import winreg  # type: ignore
        key_path = r"Software\Microsoft\Windows\CurrentVersion\Internet Settings"
        with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path) as key:
            proxy_enable = int(winreg.QueryValueEx(key, "ProxyEnable")[0] or 0)
            proxy_server = str(winreg.QueryValueEx(key, "ProxyServer")[0] or "")
            auto_config_url = str(winreg.QueryValueEx(key, "AutoConfigURL")[0] or "")
        
        if auto_config_url:
            return "", f"检测到 PAC: {auto_config_url}"

        if proxy_enable == 1 and proxy_server:
            server = proxy_server.strip()
            if ";" in server or "=" in server:
                parts = [p.strip() for p in server.split(";") if p.strip()]
                mapping: dict[str, str] = {}
                for p in parts:
                    if "=" in p:
                        k, v = p.split("=", 1)
                        mapping[k.strip().lower()] = v.strip()
                    else:
                        mapping["http"] = p
                        mapping["https"] = p
                candidate = mapping.get("https") or mapping.get("http") or ""
                if candidate:
                    if not re.match(r"^[a-zA-Z][a-zA-Z0-9+.-]*://", candidate):
                        candidate = "http://" + candidate
                    return candidate, ""
            else:
                if not re.match(r"^[a-zA-Z][a-zA-Z0-9+.-]*://", server):
                    server = "http://" + server
                return server, ""
    except Exception:
        pass

    # 2. 尝试常见代理端口兜底检测
    common_ports = [7890, 7897, 10809, 1080, 8080]
    for port in common_ports:
        candidate = f"http://127.0.0.1:{port}"
        try:
            # 快速检测端口是否开放
            import socket
            with socket.create_connection(("127.0.0.1", port), timeout=0.5):
                return candidate, f"自动扫描到本地代理: {candidate}"
        except Exception:
            continue

    return "", ""


def strip_ansi(s: str) -> str:
    return re.sub(r"\x1B\[[0-?]*[ -/]*[@-~]", "", s)

def has_js_challenge_failure(lines: list[str]) -> bool:
    tail = "\n".join(lines[-120:]).lower()
    return (
        "challenge solving failed" in tail
        or "error solving n challenge" in tail
        or "only images are available" in tail
    )

def has_po_token_required(lines: list[str]) -> bool:
    tail = "\n".join(lines[-160:]).lower()
    return (
        ("po token" in tail and "required" in tail)
        or ("po token" in tail and "not provided" in tail)
        or ("po token" in tail and "skipped" in tail)
        or ("gvs po token" in tail)
    )


def detect_js_runtime() -> tuple[bool, str]:
    runtimes = []
    for name in ("node", "deno", "bun", "quickjs", "qjs"):
        if shutil.which(name):
            runtimes.append(name)
    return bool(runtimes), ",".join(runtimes) if runtimes else "none"

def has_login_required(lines: list[str], msg: str = "") -> bool:
    tail = "\n".join(lines[-160:]).lower()
    m = (msg or "").lower()
    return (
        "login_required" in tail
        or "sign in to confirm you’re not a bot" in tail
        or "sign in to confirm you're not a bot" in tail
        or "use --cookies-from-browser" in tail
        or "use --cookies" in tail
        or "login_required" in m
        or "sign in to confirm you’re not a bot" in m
        or "sign in to confirm you're not a bot" in m
        or "use --cookies-from-browser" in m
        or "use --cookies" in m
    )


def has_premium_only_warning(lines: list[str], msg: str = "") -> bool:
    """检查是否为会员专享或高码率受限。"""
    tail = "\n".join(lines[-160:]).lower()
    m = (msg or "").lower()
    premium_markers = [
        "become a premium member",
        "formats) 1080p",
        "formats are missing",
        "format(s)",
        "高码率 are missing",
        "高码率",
    ]
    combined = tail + "\n" + m
    return any(marker in combined for marker in premium_markers)

def is_html_like_text(text: str | None) -> bool:
    if not text:
        return False
    sample = text.strip().lower()
    if sample.startswith("<!doctype") or sample.startswith("<html"):
        return True
    if "<html" in sample[:2000]:
        return True
    head_idx = sample.find("<head")
    body_idx = sample.find("<body")
    if head_idx != -1 and body_idx != -1 and abs(head_idx - body_idx) < 5000:
        return True
    return False

def _audio_cache_root() -> Path | None:
    try:
        root = Path(__file__).resolve().parent / ".cache" / "audio"
        root.mkdir(parents=True, exist_ok=True)
        return root
    except Exception:
        return None

def _audio_cache_key(video_url: str) -> str:
    try:
        vid = extract_video_id(video_url)
        if vid:
            return vid
    except Exception:
        pass
    h = hashlib.sha1(video_url.encode("utf-8", errors="ignore")).hexdigest()
    return h

def _audio_cache_path(video_url: str) -> Path | None:
    root = _audio_cache_root()
    if not root:
        return None
    key = _audio_cache_key(video_url)
    return root / f"{key}.wav"


def get_effective_proxy(proxy_url: str, use_system_proxy: bool) -> tuple[str, str]:
    proxy_url = (proxy_url or "").strip()
    if proxy_url:
        return proxy_url, ""
    if use_system_proxy:
        return detect_windows_proxy()
    return "", ""


def is_fatal_network_error(msg: str) -> bool:
    """
    检测是否为不可恢复的全局网络错误（代理失效、连接拒绝、DNS 解析失败等）。
    出现此类错误时，继续轮换客户端/Cookie 策略毫无意义，应立即中断所有重试。
    """
    s = str(msg or "").lower()
    return any([
        "proxyerror" in s,
        "proxy error" in s,
        "cannot connect to proxy" in s,
        "tunnel connection failed" in s,
        "failed to establish a new connection" in s,
        "name or service not known" in s,
        "getaddrinfo failed" in s,
        "nodename nor servname provided" in s,
        "connection refused" in s,
        "connection reset by peer" in s,
        "network is unreachable" in s,
        "no route to host" in s,
        # 整体超时（不同于单次请求超时）
        "已超过整体超时时间" in s,
        "timed out" in s and ("connect" in s or "proxy" in s),
    ])



def build_api(proxy_url: str, timeout_seconds: float, use_system_proxy: bool, retries: int) -> YouTubeTranscriptApi:
    session = TimeoutSession(timeout_seconds=max(1.0, float(timeout_seconds)))
    session.trust_env = True
    effective_proxy, pac_note = get_effective_proxy(proxy_url, use_system_proxy)
    if effective_proxy:
        session.trust_env = False
        session.proxies = {"http": effective_proxy, "https": effective_proxy}
    session._deadline_s = time.monotonic() + max(2.0, float(timeout_seconds))  # type: ignore[attr-defined]
    retry_cfg = Retry(total=max(0, int(retries)), backoff_factor=0.6, status_forcelist=[429, 500, 502, 503, 504])
    adapter = HTTPAdapter(max_retries=retry_cfg)
    session.mount("http://", adapter)
    session.mount("https://", adapter)
    session._exception_retries = max(1, int(retries))
    api = YouTubeTranscriptApi(http_client=session)
    api._pac_note = pac_note  # type: ignore[attr-defined]
    api._effective_proxy = effective_proxy  # type: ignore[attr-defined]
    api._timeout_seconds = float(timeout_seconds)  # type: ignore[attr-defined]
    api._retries = int(retries)  # type: ignore[attr-defined]
    return api


def vtt_to_text(vtt: str) -> str:
    lines: list[str] = []
    for raw in vtt.splitlines():
        s = raw.strip()
        if not s:
            continue
        if s.startswith("WEBVTT"):
            continue
        if "-->" in s:
            continue
        if s.isdigit():
            continue
        if s.startswith("NOTE"):
            continue
        lines.append(s)

    out: list[str] = []
    prev = None
    for s in lines:
        if s == prev:
            continue
        out.append(s)
        prev = s
    return "\n".join(out)


def fetch_subtitles_with_ytdlp(
    video_url: str,
    preferred_langs: list[str],
    proxy_url: str,
    timeout_seconds: float,
    retries: int,
    cookies_file: str,
    cookies_from_browser: str,
) -> tuple[str, str]:
    try:
        from yt_dlp import YoutubeDL  # type: ignore
        from yt_dlp.utils import DownloadError  # type: ignore
    except Exception as e:
        raise RuntimeError(
            "未安装 yt-dlp，无法启用兜底抓取。可执行：python -m pip install yt-dlp -i https://pypi.tuna.tsinghua.edu.cn/simple"
        ) from e

    # 移除 detect_js_runtime 的定义，改用 yt-dlp 原生解释器或 yt_dlp_ejs
    langs = preferred_langs[:] if preferred_langs else []
    with tempfile.TemporaryDirectory() as tmp:
        outtmpl = os.path.join(tmp, "%(id)s.%(ext)s")
        last_err: Exception | None = None
        last_cookie_error: RuntimeError | None = None
        disabled_browsers: set[str] = set()
        last_video_id = ""
        disabled_web = False
        force_browser_cookie = False

        import random

        user_agents = [
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 14_0) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.0 Safari/605.1.15",
            "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/121.0.0.0 Safari/537.36",
        ]

        class YdlLogger:
            def __init__(self) -> None:
                self.lines: list[str] = []

            def _add(self, level: str, msg: str) -> None:
                s = strip_ansi(str(msg or "")).strip()
                if not s:
                    return
                self.lines.append(f"[{level}] {s}")
                if len(self.lines) > 400:
                    self.lines = self.lines[-250:]

            def debug(self, msg: str) -> None:
                self._add("debug", msg)

            def info(self, msg: str) -> None:
                self._add("info", msg)

            def warning(self, msg: str) -> None:
                self._add("warning", msg)

            def error(self, msg: str) -> None:
                self._add("error", msg)

        def extract_info_with_format_rotation(ydl_opts: dict) -> dict:
            format_candidates = [None, "bestaudio/best", "worstaudio/worst"]
            last_format_err: Exception | None = None
            for fmt in format_candidates:
                if fmt:
                    ydl_opts["format"] = fmt
                else:
                    ydl_opts.pop("format", None)
                try:
                    with YoutubeDL(ydl_opts) as ydl:
                        return ydl.extract_info(video_url, download=False, process=False)
                except DownloadError as e:
                    msg = strip_ansi(str(e))
                    if "Requested format is not available" in msg:
                        last_format_err = e
                        continue
                    raise
            if last_format_err is not None:
                raise last_format_err
            return {}

        # 移除此处的本地 is_cookie_error 定义，改用顶层的 _is_cookie_error

        def build_download_session() -> requests.Session:
            sess = TimeoutSession(timeout_seconds=max(1.0, float(timeout_seconds)))
            sess.trust_env = True
            if proxy_url:
                sess.trust_env = False
                sess.proxies = {"http": proxy_url, "https": proxy_url}
            sess._deadline_s = time.monotonic() + max(2.0, float(timeout_seconds))  # type: ignore[attr-defined]
            retry_cfg = Retry(
                total=max(0, int(retries)),
                backoff_factor=0.8,
                status_forcelist=[500, 502, 503, 504],
            )
            adapter = HTTPAdapter(max_retries=retry_cfg)
            sess.mount("http://", adapter)
            sess.mount("https://", adapter)
            sess._exception_retries = max(1, int(retries))
            return sess

        def maybe_force_vtt(url: str) -> str:
            if "timedtext" not in url:
                return url
            if "fmt=" in url:
                return url
            return url + ("&" if "?" in url else "?") + "fmt=vtt"

        def try_extract_and_download_by_url() -> tuple[str, str] | None:
            nonlocal force_browser_cookie, last_cookie_error, disabled_browsers
            sess = build_download_session()
            client_sets: list[list[str]] = [["android"], ["ios"], ["web_safari"], ["web"], ["tv"]]
            attempt_retries = max(1, int(retries))
            
            # 用于记录 "无 Cookie" 模式下的错误
            no_cookie_error: Exception | None = None
            
            for client_set in client_sets:
                if disabled_web and "web" in client_set:
                    continue
                for attempt in range(attempt_retries):
                    ua = random.choice(user_agents)
                    for cookiefile, cfb in CookieManager.get_sources(cookies_file, cookies_from_browser, force_browser_cookie):
                        if cfb and cfb in disabled_browsers:
                            continue
                        info = None
                        logger = YdlLogger()
                        opts: dict = {
                            "skip_download": True,
                            "quiet": True,
                            "no_warnings": True,
                            "verbose": True,
                            "nocheckcertificate": True,
                            "socket_timeout": float(timeout_seconds),
                            "geo_bypass": True,
                            "ignoreerrors": False,
                            "ignore_config": True,  # 忽略全局配置文件，避免意外读取 Cookie
                            "impersonate": "chrome",
                            # "format": "bestaudio/best", # 移除强制格式，允许自动选择最佳可用格式（仅获取信息）
                            "http_headers": {"User-Agent": ua, "Accept-Language": "en-US,en;q=0.9"},
                            "extractor_args": {"youtube": {"player_client": client_set}},
                            "logger": logger,
                        }
                        if proxy_url:
                            opts["proxy"] = proxy_url
                        if cookiefile:
                            opts["cookiefile"] = cookiefile
                        elif cfb:
                            opts["cookiesfrombrowser"] = (cfb,)

                        try:
                            info = extract_info_with_format_rotation(opts)
                        except DownloadError as e:
                            last_msg = strip_ansi(str(e))
                            has_cookie_in_log = any(CookieManager.is_cookie_error(line) for line in logger.lines)
                            
                            # ⚡ 快速中断：致命网络错误（代理/DNS/连接拒绝）— 继续重试毫无意义
                            if is_fatal_network_error(last_msg):
                                raise RuntimeError(f"网络连接失败（代理/DNS 错误），已中断所有重试: {last_msg}") from e
                            
                            if CookieManager.is_cookie_error(last_msg) or has_cookie_in_log:
                                custom_err = RuntimeError(CookieManager.get_fatal_msg(last_msg, cfb))
                                if not cookiefile and not cfb:
                                    no_cookie_error = custom_err
                                    last_err = custom_err
                                else:
                                    last_cookie_error = custom_err
                                    if cfb:
                                        disabled_browsers.add(cfb)
                                    last_err = no_cookie_error if no_cookie_error else custom_err
                                continue
                                
                            if has_login_required([], last_msg) or "Sign in to confirm" in last_msg:
                                force_browser_cookie = True
                                last_err = RuntimeError("需要登录或验证（可能触发人机校验）。已尝试自动读取浏览器 Cookie；如仍失败请手动开启、提供 cookies 文件或改用 Firefox。")
                                continue
                            
                            # 错误处理优化
                            if not cookiefile and not cfb:
                                no_cookie_error = e
                                last_err = e
                            else:
                                last_err = no_cookie_error if no_cookie_error else e

                            if "HTTP Error 429" in last_msg or "429" in last_msg:
                                time.sleep(2.5 * (attempt + 1))
                                continue
                            if cfb and CookieManager.is_cookie_error(last_msg):
                                if cfb:
                                    disabled_browsers.add(cfb)
                                continue
                            continue
                        except Exception as e:
                            if not cookiefile and not cfb:
                                no_cookie_error = e
                                last_err = e
                            else:
                                if no_cookie_error:
                                    last_err = no_cookie_error
                                elif CookieManager.is_cookie_error(e):
                                    if no_cookie_error:
                                        last_err = no_cookie_error
                                    else:
                                        last_err = e
                                else:
                                    last_err = e
                                    
                            if cfb and CookieManager.is_cookie_error(e):
                                if cfb:
                                    disabled_browsers.add(cfb)
                                continue
                            continue
                    
                        if info is not None:
                            if has_login_required(logger.lines):
                                force_browser_cookie = True
                                disabled_web = disabled_web or ("web" in client_set)
                                last_err = RuntimeError("需要登录或验证（可能触发人机校验）。已尝试自动读取浏览器 Cookie；如仍失败请手动开启或提供 cookies 文件。")
                                continue
                            if has_po_token_required(logger.lines) and any(c in {"android", "ios", "mweb"} for c in client_set):
                                last_err = RuntimeError("该客户端需要 PO Token，已降级到其他客户端。")
                                disabled_web = disabled_web or ("web" in client_set)
                                continue
                            if has_js_challenge_failure(logger.lines):
                                disabled_web = True
                                last_err = RuntimeError("JS challenge 失败，可能导致格式缺失。建议升级 yt-dlp 或更换网络。")
                                continue
                            break

                    if not isinstance(info, dict):
                        continue

                    video_id = str(info.get("id") or "")
                    subtitles = info.get("subtitles") or {}
                    auto = info.get("automatic_captions") or {}
                    merged: dict = {}
                    if isinstance(subtitles, dict):
                        merged.update(subtitles)
                    if isinstance(auto, dict):
                        for k, v in auto.items():
                            if k not in merged:
                                merged[k] = v

                    preferred = [l for l in langs if l and l != "all"]
                    try_langs = preferred or [str(k) for k in merged.keys()]
                    for lang in try_langs:
                        tracks = merged.get(lang)
                        if not isinstance(tracks, list) or not tracks:
                            continue
                        track = None
                        for t in tracks:
                            if isinstance(t, dict) and (t.get("ext") == "vtt" or str(t.get("ext") or "") == "vtt"):
                                track = t
                                break
                        if track is None:
                            for t in tracks:
                                if isinstance(t, dict):
                                    track = t
                                    break
                        if not isinstance(track, dict):
                            continue
                        url = str(track.get("url") or "")
                        if not url:
                            continue
                        url = maybe_force_vtt(url)

                        try:
                            r = sess.get(url, headers={"User-Agent": ua, "Accept-Encoding": "identity"})
                            if r.status_code == 429:
                                retry_after = r.headers.get("Retry-After", "")
                                try:
                                    wait_s = int(retry_after)
                                except Exception:
                                    wait_s = 20
                                last_err = RuntimeError(f"timedtext 触发 429（Retry-After={retry_after or 'n/a'}）")
                                time.sleep(max(8, min(90, wait_s)))
                                continue
                            r.raise_for_status()
                        except requests.exceptions.HTTPError as e:
                            last_err = e
                            continue
                        except requests.exceptions.RequestException:
                            last_err = None
                            continue

                        raw = r.text or ""
                        text = vtt_to_text(raw) if raw else ""
                        if text.strip():
                            label = f"{video_id or ''}{' ' if video_id else ''}{lang}".strip()
                            return label, text

            return None

        def download_with_lang(sub_langs: list[str]) -> tuple[str, list[Path]]:
            nonlocal force_browser_cookie
            nonlocal last_err, last_video_id
            attempt_retries = max(1, int(retries))
            
            # 记录无 Cookie 时的错误
            no_cookie_error: Exception | None = None
            
            client_strategies = [["android"], ["ios"]] + ([["web_safari"], ["web"]] if allow_web_client else []) + [["tv"]]

            for attempt in range(attempt_retries + 1): # 多尝试一次以覆盖更多策略
                # 轮换客户端策略
                client_set = client_strategies[attempt % len(client_strategies)]
                if disabled_web and "web" in client_set:
                    continue
                
                # ua = random.choice(user_agents) # 不再手动指定 UA，让 yt-dlp 根据 client 自动匹配
                
                for cookiefile, cfb in CookieManager.get_sources(cookies_file, cookies_from_browser, force_browser_cookie):
                    if cfb and cfb in disabled_browsers:
                        continue
                    logger = YdlLogger()
                    opts: dict = {
                        "skip_download": True,
                        "writesubtitles": True,
                        "writeautomaticsub": True,
                        "subtitleslangs": sub_langs,
                        "subtitlesformat": "vtt",
                        "outtmpl": outtmpl,
                        "quiet": True,
                        "no_warnings": True,
                        "verbose": True,
                        "nocheckcertificate": True,
                        "sleep_interval": 2,
                        "max_sleep_interval": 6,
                        "retries": 1,
                        "fragment_retries": 1,
                        "extractor_retries": 1,
                        "socket_timeout": float(timeout_seconds),
                        "geo_bypass": True,
                        "ignoreerrors": False,
                        "ignore_config": True,
                        "impersonate": "chrome",
                        "http_headers": {"Accept-Language": "en-US,en;q=0.9"},
                        "extractor_args": {"youtube": {"player_client": client_set}},
                        "logger": logger,
                    }
                    if proxy_url:
                        opts["proxy"] = proxy_url
                    if cookiefile:
                        opts["cookiefile"] = cookiefile
                    elif cfb:
                        opts["cookiesfrombrowser"] = (cfb,)

                    try:
                        info = extract_info_with_format_rotation(opts)
                        if not isinstance(info, dict):
                            continue
                        if has_login_required(logger.lines):
                            force_browser_cookie = True
                            disabled_web = disabled_web or ("web" in client_set)
                            last_err = RuntimeError("需要登录或验证（可能触发人机校验）。已尝试自动读取浏览器 Cookie；如仍失败请手动开启或提供 cookies 文件。")
                            continue
                        if has_po_token_required(logger.lines) and any(c in {"android", "ios", "mweb"} for c in client_set):
                            last_err = RuntimeError("该客户端需要 PO Token，已降级到其他客户端。")
                            disabled_web = disabled_web or ("web" in client_set)
                            continue
                        if has_js_challenge_failure(logger.lines):
                            disabled_web = True
                            last_err = RuntimeError("JS challenge 失败，可能导致格式缺失。建议升级 yt-dlp 或更换网络。")
                            continue
                        last_video_id = (info or {}).get("id") or last_video_id
                        last_err = None
                        subtitles = info.get("subtitles") or {}
                        auto = info.get("automatic_captions") or {}
                        merged: dict = {}
                        if isinstance(subtitles, dict):
                            merged.update(subtitles)
                        if isinstance(auto, dict):
                            for k, v in auto.items():
                                if k not in merged:
                                    merged[k] = v
                        preferred = [l for l in sub_langs if l and l != "all"]
                        try_langs = preferred or [str(k) for k in merged.keys()]
                        sess = build_download_session()
                        ua = random.choice(user_agents)
                        for lang in try_langs:
                            tracks = merged.get(lang)
                            if not isinstance(tracks, list) or not tracks:
                                continue
                            track = None
                            for t in tracks:
                                if isinstance(t, dict) and (t.get("ext") == "vtt" or str(t.get("ext") or "") == "vtt"):
                                    track = t
                                    break
                            if track is None:
                                for t in tracks:
                                    if isinstance(t, dict):
                                        track = t
                                        break
                            if not isinstance(track, dict):
                                continue
                            url = str(track.get("url") or "")
                            if not url:
                                continue
                            url = maybe_force_vtt(url)
                            try:
                                r = sess.get(url, headers={"User-Agent": ua, "Accept-Encoding": "identity"})
                                if r.status_code == 429:
                                    retry_after = r.headers.get("Retry-After", "")
                                    try:
                                        wait_s = int(retry_after)
                                    except Exception:
                                        wait_s = 20
                                    last_err = RuntimeError(f"timedtext 触发 429（Retry-After={retry_after or 'n/a'}）")
                                    time.sleep(max(8, min(90, wait_s)))
                                    continue
                                r.raise_for_status()
                            except requests.exceptions.HTTPError as e:
                                last_err = e
                                continue
                            except requests.exceptions.RequestException:
                                last_err = None
                                continue
                            raw = r.text or ""
                            if raw.strip():
                                file_name = f"{last_video_id or 'video'}.{lang}.vtt"
                                vtt_path = Path(tmp) / file_name
                                vtt_path.write_text(raw, encoding="utf-8", errors="ignore")
                                return last_video_id, [vtt_path]
                        last_err = RuntimeError("yt-dlp 未下载到字幕文件（可能无字幕或被限制）。")
                        return last_video_id, []
                    except DownloadError as e:
                        msg = strip_ansi(str(e))
                        has_cookie_in_log = any(CookieManager.is_cookie_error(line) for line in logger.lines)
                        if CookieManager.is_cookie_error(msg) or has_cookie_in_log:
                            last_cookie_error = RuntimeError(CookieManager.get_fatal_msg(msg, cfb))
                            if cfb:
                                disabled_browsers.add(cfb)
                            continue
                        else:
                            last_err = RuntimeError(msg + "\n\n" + "\n".join(logger.lines[-80:]))
                            if not cookiefile and not cfb:
                                no_cookie_error = last_err
                            elif no_cookie_error:
                                last_err = no_cookie_error
                        if "HTTP Error 429" in msg or "429" in msg:
                            time.sleep(2.0 * (attempt + 1))
                            continue
                        if cfb and CookieManager.is_cookie_error(msg):
                            continue
                        continue
                    except Exception as e:
                        last_err = RuntimeError(repr(e) + "\n\n" + "\n".join(logger.lines[-80:]))
                        if not cookiefile and not cfb:
                            no_cookie_error = last_err
                        else:
                            if no_cookie_error:
                                last_err = no_cookie_error
                        if cfb and CookieManager.is_cookie_error(e):
                            continue
                        continue
            return last_video_id, []

        def try_timedtext_direct(video_id: str, langs_try: list[str]) -> tuple[str, str] | None:
            if not video_id:
                return None
            sess = build_download_session()
            ua = random.choice(user_agents)
            for lang in langs_try:
                for kind in ("", "asr"):
                    params = {"v": video_id, "lang": lang, "fmt": "vtt"}
                    if kind:
                        params["kind"] = kind
                    try:
                        r = sess.get(
                            "https://www.youtube.com/api/timedtext",
                            params=params,
                            headers={"User-Agent": ua, "Accept-Encoding": "identity"},
                        )
                        if r.status_code == 429:
                            retry_after = r.headers.get("Retry-After", "")
                            try:
                                wait_s = int(retry_after)
                            except Exception:
                                wait_s = 20
                            last_err = RuntimeError(f"timedtext 触发 429（Retry-After={retry_after or 'n/a'}）")
                            time.sleep(max(8, min(90, wait_s)))
                            continue
                        r.raise_for_status()
                    except requests.exceptions.RequestException as e:
                        last_err = e
                        continue
                    raw = r.text or ""
                    if is_html_like_text(raw):
                        last_err = RuntimeError("timedtext 返回 HTML 源码")
                        continue
                    text = vtt_to_text(raw) if raw else ""
                    if text.strip():
                        label = f"{video_id} {lang}{' asr' if kind else ''}".strip()
                        return label, text
            return None

        def choose_vtt(vtts: list[Path], lang: str | None) -> Path:
            if lang:
                for p in vtts:
                    if f".{lang}." in p.name or p.name.endswith(f".{lang}.vtt"):
                        return p
            return vtts[0]

        direct = try_extract_and_download_by_url()
        if direct is not None:
            return direct

        preferred = [l for l in langs if l and l != "all"]
        if preferred:
            for lang in preferred:
                last_video_id, vtts = download_with_lang([lang])
                if vtts:
                    chosen = choose_vtt(vtts, lang)
                    raw = chosen.read_text(encoding="utf-8", errors="ignore")
                    text = vtt_to_text(raw)
                    label = f"{last_video_id or ''}{' ' if last_video_id else ''}{chosen.name}".strip()
                    return label, text
        last_video_id, vtts = download_with_lang(["all"])
        if vtts:
            chosen = choose_vtt(vtts, None)
            raw = chosen.read_text(encoding="utf-8", errors="ignore")
            text = vtt_to_text(raw)
            label = f"{last_video_id or ''}{' ' if last_video_id else ''}{chosen.name}".strip()
            return label, text

        preferred = [l for l in langs if l and l != "all"]
        try_video_id = last_video_id
        if not try_video_id:
            try:
                try_video_id = extract_video_id(video_url)
            except Exception:
                try_video_id = ""
        timedtext = try_timedtext_direct(try_video_id, preferred or ["en", "zh-Hans", "zh", "zh-Hant"])
        if timedtext is not None:
            return timedtext

        if last_err is not None:
            raise last_err
        if last_cookie_error:
            raise last_cookie_error
        raise RuntimeError("yt-dlp 兜底抓取失败（未知原因）。")


def transcribe_audio_with_whisper(audio_path: str, model_name: str, language: str, proxy_url: str = None, status_callback=None, fast_mode: bool = False, force_cpu: bool = False) -> str:
    # 优先尝试使用 faster-whisper (速度快 4-10 倍)
    # 记录 faster-whisper 的错误信息，以便 fallback 时展示
    fw_error = None
    
    if status_callback: status_callback("Initializing Whisper engine...")

    try:
        from faster_whisper import WhisperModel
        
        # 设置临时环境变量以支持代理下载模型
        original_http = os.environ.get("HTTP_PROXY")
        original_https = os.environ.get("HTTPS_PROXY")
        
        if proxy_url:
            os.environ["HTTP_PROXY"] = proxy_url
            os.environ["HTTPS_PROXY"] = proxy_url
            print(f"Setting proxy for faster-whisper model download: {proxy_url}")
        
        def _safe_status(message: str) -> None:
            try:
                if status_callback:
                    status_callback(message)
            except Exception:
                pass

        def _get_wav_duration_seconds(path: str) -> float | None:
            try:
                with wave.open(path, "rb") as wf:
                    frames = wf.getnframes()
                    rate = wf.getframerate()
                    if rate <= 0:
                        return None
                    return float(frames) / float(rate)
            except Exception:
                return None

        def _transcribe_in_worker(q: "queue.Queue[object]", transcribe_kwargs: dict) -> None:
            try:
                q.put(("begin", time.time()))
                segments_local, info_local = model.transcribe(audio_path, **transcribe_kwargs)

                text_segments_local: list[str] = []
                last_emit = time.time()
                seg_count = 0
                last_end_s: float | None = None
                for seg in segments_local:
                    seg_count += 1
                    text_segments_local.append(seg.text)
                    try:
                        last_end_s = float(getattr(seg, "end", None))
                    except Exception:
                        last_end_s = None
                    now = time.time()
                    if now - last_emit >= 10:
                        q.put(("progress", {"segments": seg_count, "end_s": last_end_s, "t": now}))
                        last_emit = now
                q.put(("done", {"text": "".join(text_segments_local).strip(), "info": info_local}))
            except Exception as e:
                q.put(("error", e))

        try:
            # 映射模型名称 (faster-whisper 使用同样的名称)
            model_size = (model_name or "").strip() or "base"
            
            # 初始化模型 (使用 int8 量化，CPU 上极快)
            # 缓存模型以避免重复加载
            cache_key = f"faster_whisper_{model_size}"
            cache = getattr(transcribe_audio_with_whisper, "_model_cache", None)
            if not isinstance(cache, dict):
                cache = {}
                setattr(transcribe_audio_with_whisper, "_model_cache", cache)
                
            model = cache.get(cache_key)
            cache_info = cache.get(f"{cache_key}_info", {}) if isinstance(cache, dict) else {}
            cached_device = str(cache_info.get("device") or "").lower()
            cached_compute = str(cache_info.get("compute_type") or "").lower()
            gpu_failed_flag = bool(cache.get(f"{cache_key}_gpu_failed", False))
            gpu_retry_once = bool(cache.get(f"{cache_key}_gpu_retry_once", False))
            cuda_ready = False
            cuda_reason = ""
            ct2_version = ""
            ct2_cuda_count = None
            try:
                import ctranslate2
                ct2_version = str(getattr(ctranslate2, "__version__", "")) or "unknown"
                if hasattr(ctranslate2, "get_cuda_device_count"):
                    ct2_cuda_count = ctranslate2.get_cuda_device_count()
                    cuda_ready = ct2_cuda_count > 0
                    if not cuda_ready:
                        cuda_reason = "ctranslate2 reports no CUDA devices"
                        _safe_status("CUDA 不可用：ctranslate2 未检测到设备")
            except Exception as e:
                cuda_reason = f"ctranslate2 check failed: {e}"
            if model is not None:
                if status_callback:
                    status_callback(f"Using cached Whisper model ({cached_device or 'cpu'} {cached_compute or 'int8'})")
                if not force_cpu and cuda_ready and cached_device == "cpu":
                    if not gpu_failed_flag:
                        cache.pop(cache_key, None)
                        cache.pop(f"{cache_key}_info", None)
                        model = None
                    elif not gpu_retry_once:
                        cache[f"{cache_key}_gpu_retry_once"] = True
                        if status_callback:
                            status_callback("CUDA 可用，尝试重新加载 GPU 模型")
                        cache.pop(cache_key, None)
                        cache.pop(f"{cache_key}_info", None)
                        model = None

            if model is None:
                # 自动检测最佳设备
                # 1. 尝试 GPU (CUDA)
                # 2. 回退到 CPU
                
                cpu_threads = os.cpu_count() or 4
                
                # 尝试加载 GPU
                # 注意：ctranslate2 需要 CUDA 11.x 或 12.x 运行时库
                device = "cpu"
                compute_type = "int8"
                
                # 尝试自动配置 NVIDIA 库路径 (从 pip 包)
                try:
                    import nvidia.cublas
                    import nvidia.cudnn
                    
                    # 尝试不同的子模块结构，适配不同版本的 nvidia 包
                    cublas_path = None
                    cudnn_path = None
                    
                    # Case 1: nvidia.cublas.lib
                    try:
                        import nvidia.cublas.lib
                        cublas_path = os.path.dirname(nvidia.cublas.lib.__file__)
                    except ImportError:
                        pass
                        
                    # Case 2: nvidia.cublas.bin (some versions)
                    if not cublas_path:
                        try:
                            import nvidia.cublas.bin
                            cublas_path = os.path.dirname(nvidia.cublas.bin.__file__)
                        except ImportError:
                            pass
                            
                    # Case 3: 根目录 (如果 __init__.py 就在 lib 旁)
                    if not cublas_path and hasattr(nvidia.cublas, '__file__') and nvidia.cublas.__file__:
                        cublas_path = os.path.join(os.path.dirname(nvidia.cublas.__file__), "lib")
                        if not os.path.exists(cublas_path):
                            cublas_path = os.path.join(os.path.dirname(nvidia.cublas.__file__), "bin")

                    # 同理 cuDNN
                    try:
                        import nvidia.cudnn.lib
                        cudnn_path = os.path.dirname(nvidia.cudnn.lib.__file__)
                    except ImportError:
                        pass
                        
                    if not cudnn_path and hasattr(nvidia.cudnn, '__file__') and nvidia.cudnn.__file__:
                         cudnn_path = os.path.join(os.path.dirname(nvidia.cudnn.__file__), "lib")
                         if not os.path.exists(cudnn_path):
                             cudnn_path = os.path.join(os.path.dirname(nvidia.cudnn.__file__), "bin")

                    if cublas_path and os.path.exists(cublas_path) and cublas_path not in os.environ["PATH"]:
                        os.environ["PATH"] = cublas_path + os.pathsep + os.environ["PATH"]
                        if hasattr(os, "add_dll_directory"):
                            try:
                                os.add_dll_directory(cublas_path)
                            except Exception:
                                pass
                        print(f"Added NVIDIA cuBLAS to PATH: {cublas_path}")
                    
                    if cudnn_path and os.path.exists(cudnn_path) and cudnn_path not in os.environ["PATH"]:
                        os.environ["PATH"] = cudnn_path + os.pathsep + os.environ["PATH"]
                        if hasattr(os, "add_dll_directory"):
                            try:
                                os.add_dll_directory(cudnn_path)
                            except Exception:
                                pass
                        print(f"Added NVIDIA cuDNN to PATH: {cudnn_path}")
                        
                except Exception as e_nvidia:
                    print(f"Failed to auto-configure NVIDIA libraries from pip packages: {e_nvidia}")

                diag_lines = []
                try:
                    diag_lines.append(f"Python {platform.python_version()} {platform.architecture()[0]}")
                    diag_lines.append(f"OS {platform.system()} {platform.release()}")
                    if ct2_version:
                        diag_lines.append(f"ctranslate2 {ct2_version}")
                    if ct2_cuda_count is not None:
                        diag_lines.append(f"ctranslate2 CUDA devices {ct2_cuda_count}")
                    if status_callback and diag_lines:
                        status_callback(" | ".join(diag_lines))
                except Exception:
                    pass

                try:
                    dll_dirs = []
                    env_paths = []
                    env_cuda_vars = []
                    for k, v in os.environ.items():
                        if k.startswith("CUDA_PATH") and v:
                            env_paths.append(v)
                            env_cuda_vars.append(f"{k}={v}")
                    for base in env_paths:
                        cand = os.path.join(base, "bin")
                        if os.path.isdir(cand):
                            dll_dirs.append(cand)
                    cuda_root = r"C:\Program Files\NVIDIA GPU Computing Toolkit\CUDA"
                    if os.path.isdir(cuda_root):
                        for name in os.listdir(cuda_root):
                            cand = os.path.join(cuda_root, name, "bin")
                            if os.path.isdir(cand):
                                dll_dirs.append(cand)
                    site_paths = [p for p in sys.path if p and "site-packages" in p.lower()]
                    for sp in site_paths:
                        pattern = os.path.join(sp, "nvidia", "cublas", "**", "cublas64_*.dll")
                        for dll in glob.glob(pattern, recursive=True):
                            dll_dirs.append(os.path.dirname(dll))
                    seen = set()
                    added_dirs = []
                    for d in dll_dirs:
                        if d in seen:
                            continue
                        seen.add(d)
                        if d not in os.environ["PATH"]:
                            os.environ["PATH"] = d + os.pathsep + os.environ["PATH"]
                            added_dirs.append(d)
                        if hasattr(os, "add_dll_directory"):
                            try:
                                os.add_dll_directory(d)
                            except Exception:
                                pass
                    if added_dirs:
                        print(f"Added CUDA DLL dirs: {len(added_dirs)}")
                        if status_callback:
                            status_callback(f"CUDA DLL dirs added: {len(added_dirs)}")
                    if status_callback:
                        preview_dirs = dll_dirs[:3]
                        status_callback(f"CUDA 搜索目录: {len(dll_dirs)} | 预览: {', '.join(preview_dirs) if preview_dirs else 'none'}")
                        if env_cuda_vars:
                            status_callback(f"CUDA 环境变量: {', '.join(env_cuda_vars[:3])}{'...' if len(env_cuda_vars) > 3 else ''}")
                except Exception as e_cuda_dirs:
                    print(f"Failed to auto-configure CUDA DLL dirs: {e_cuda_dirs}")

                # 显式打印正在加载模型，这在日志中可见
                print(f"Initializing faster-whisper model '{model_size}'...")
                
                import ctypes
                dll_candidates = [
                    "cublas64_12.dll",
                    "cublas64_11.dll",
                    "cublas64_10.dll",
                    "cublas64_10_2.dll",
                    "cudart64_12.dll",
                    "cudart64_11.dll",
                    "cudart64_10.dll",
                    "cudart64_10_2.dll",
                    "cudnn64_9.dll",
                    "cudnn64_8.dll",
                ]
                found_dlls = []
                failed_dlls = []
                found_paths = {}
                search_dirs = []
                try:
                    search_dirs.extend([p for p in os.environ.get("PATH", "").split(os.pathsep) if p])
                except Exception:
                    pass
                extra_dirs = locals().get("dll_dirs", [])
                if isinstance(extra_dirs, list):
                    search_dirs.extend(extra_dirs)
                seen_dirs = set()
                search_dirs = [d for d in search_dirs if d and not (d in seen_dirs or seen_dirs.add(d))]
                for name in dll_candidates:
                    loaded = False
                    try:
                        ctypes.CDLL(name)
                        loaded = True
                        found_paths[name] = name
                    except Exception:
                        pass
                    if not loaded:
                        for d in search_dirs:
                            dll_path = os.path.join(d, name)
                            if os.path.isfile(dll_path):
                                try:
                                    ctypes.CDLL(dll_path)
                                    loaded = True
                                    found_paths[name] = dll_path
                                    break
                                except Exception:
                                    pass
                    if loaded:
                        found_dlls.append(name)
                    else:
                        failed_dlls.append(name)

                has_cuda_libs = len(found_dlls) > 0
                if status_callback:
                    status_callback(f"CUDA DLL 探测: 命中 {len(found_dlls)} / 候选 {len(dll_candidates)}")
                    if found_paths:
                        preview_paths = list(found_paths.values())[:3]
                        status_callback(f"CUDA DLL 路径预览: {', '.join(preview_paths)}")
                    has_cublas = any(n.startswith("cublas") for n in found_dlls)
                    has_cudart = any(n.startswith("cudart") for n in found_dlls)
                    has_cudnn = any(n.startswith("cudnn") for n in found_dlls)
                    status_callback(f"CUDA DLL 族: cublas={has_cublas} cudart={has_cudart} cudnn={has_cudnn}")
                    if search_dirs:
                        for i in range(0, len(search_dirs), 8):
                            status_callback("CUDA 搜索目录清单: " + "; ".join(search_dirs[i:i+8]))
                    if found_paths:
                        items = [f"{k} -> {found_paths[k]}" for k in sorted(found_paths.keys())]
                        for i in range(0, len(items), 10):
                            status_callback("CUDA DLL 命中清单: " + "; ".join(items[i:i+10]))
                    if failed_dlls:
                        for i in range(0, len(failed_dlls), 10):
                            status_callback("CUDA DLL 未命中清单: " + ", ".join(failed_dlls[i:i+10]))
                if not has_cuda_libs:
                    msg = "⚠️ CUDA DLL 未找到，跳过 GPU，强制 CPU。"
                    print(msg)
                    if status_callback:
                        status_callback(msg)
                if not cuda_ready:
                    cuda_ready = has_cuda_libs
                    if not cuda_ready and not cuda_reason:
                        cuda_reason = "CUDA DLL 未找到"
                if status_callback and found_dlls:
                    status_callback(f"CUDA DLL 就绪: {', '.join(found_dlls[:3])}{'...' if len(found_dlls) > 3 else ''}")

                try:
                    if force_cpu:
                        raise RuntimeError("Forced CPU mode")
                    if not cuda_ready:
                        raise RuntimeError(f"CUDA not ready (Skipping): {cuda_reason}")

                    # 静默尝试 GPU (float16)
                    # 增加 download_root 参数，确保模型下载到项目目录下，方便管理且避免权限问题
                    # 使用 os.getcwd() 可能会变，建议用固定相对路径 "models"
                    model_dir = str(get_models_dir())
                    os.makedirs(model_dir, exist_ok=True)
                    
                    print(f"Attempting to load on GPU (cuda)... Download root: {model_dir}")
                    if status_callback: status_callback(f"Attempting to load Whisper on GPU (CUDA)...")
                    
                    # 1. 尝试 GPU float16 (最佳性能，但显存要求高)
                    try:
                        model = WhisperModel(model_size, device="cuda", compute_type="float16", download_root=model_dir)
                        device = "cuda"
                        compute_type = "float16"
                        print(f"✅ faster-whisper loaded on GPU (cuda) with float16. Threads: {cpu_threads}")
                        if status_callback: status_callback(f"✅ Whisper loaded on GPU (CUDA float16)")
                    except Exception as e_f16:
                        print(f"⚠️ GPU float16 load failed: {e_f16}")
                        
                        # 2. 尝试 GPU int8 (节省显存，适合 4GB 显存如 1050 Ti 跑 large 模型)
                        try:
                            print(f"Attempting to load on GPU (cuda) int8...")
                            if status_callback: status_callback(f"GPU float16 failed. Trying GPU int8...")
                            model = WhisperModel(model_size, device="cuda", compute_type="int8", download_root=model_dir)
                            device = "cuda"
                            compute_type = "int8"
                            print(f"✅ faster-whisper loaded on GPU (cuda) with int8. Threads: {cpu_threads}")
                            if status_callback: status_callback(f"✅ Whisper loaded on GPU (CUDA int8)")
                        except Exception as e_int8:
                            print(f"❌ GPU int8 load failed: {e_int8}")
                            import traceback
                            traceback.print_exc()
                            raise e_int8 # 抛出异常以触发外层的 CPU fallback
    
                except Exception as e:
                    print(f"GPU load failed completely ({e}), fallback to CPU int8")
                    if status_callback: status_callback(f"❌ GPU failed ({str(e)[:50]}...). Falling back to CPU.")
                    cache[f"{cache_key}_gpu_failed_reason"] = str(e)
                    # 3. 回退到 CPU int8
                    model_dir = str(get_models_dir())
                    model = WhisperModel(model_size, device="cpu", compute_type="int8", cpu_threads=cpu_threads, download_root=model_dir)
                    device = "cpu"
                    compute_type = "int8"
                    cache[f"{cache_key}_gpu_failed"] = True
                    print(f"✅ faster-whisper loaded on CPU with int8. Threads: {cpu_threads}")
                    if status_callback: status_callback(f"✅ Whisper loaded on CPU (int8)")
                
                # 记录 device 信息到 model 对象，方便外部读取（虽然 model 是 C++ 包装对象可能不支持随意 setattr）
                # 我们可以存到 cache 的 metadata 里
                cache[cache_key] = model
                cache[f"{cache_key}_info"] = {"device": device, "compute_type": compute_type}
                if device == "cuda":
                    cache[f"{cache_key}_gpu_failed"] = False
                    cache.pop(f"{cache_key}_gpu_failed_reason", None)
                
            # 开始转写
            lang = None if (not language or language == "auto") else language
            
            device_for_status = cache.get(f"{cache_key}_info", {}).get("device", "cpu")
            _safe_status(f"Transcribing audio with Whisper ({device_for_status})...")

            try:
                wav_seconds = _get_wav_duration_seconds(audio_path)
                if wav_seconds is not None:
                    _safe_status(f"Audio duration: {wav_seconds:.1f}s")

                used_device_now = str(cache.get(f"{cache_key}_info", {}).get("device", "cpu")).lower()
                used_compute_now = str(cache.get(f"{cache_key}_info", {}).get("compute_type", "int8")).lower()
                cpu_count = os.cpu_count() or 4
                running_on_render = bool(str(os.environ.get("RENDER_SERVICE_ID", "") or "").strip())
                audio_minutes = (float(wav_seconds) / 60.0) if wav_seconds is not None else None
                batch_size = 8
                if used_device_now == "cuda":
                    # 针对 int8 和不同模型大小优化 batch_size
                    base_batch = 8
                    if model_size in {"tiny", "base"}:
                        base_batch = 24
                    elif model_size in {"small"}:
                        base_batch = 16
                    elif model_size in {"medium"}:
                        base_batch = 12
                    else: # large
                        base_batch = 6
                    
                    if used_compute_now == "int8":
                        # int8 显存占用更小，可以增加 batch_size 以提高吞吐量
                        base_batch = int(base_batch * 1.5)
                    
                    batch_size = base_batch
                    
                    if audio_minutes is not None and audio_minutes >= 90:
                        # 长音频适当保守一点，避免显存碎片化
                        batch_size = max(4, min(batch_size, 16))
                else:
                    base_cpu_batch = 20 if cpu_count >= 16 else (16 if cpu_count >= 8 else 8)
                    if model_size in {"tiny", "base", "small"}:
                        batch_size = base_cpu_batch
                    else:
                        batch_size = min(10, base_cpu_batch)
                    if audio_minutes is not None and audio_minutes >= 60:
                        batch_size = max(8, min(batch_size, 12))
                    if running_on_render:
                        batch_size = min(batch_size, 4)
                        _safe_status("Render 内存保护：已降低 CPU batch_size")

                auto_fast_mode = False
                if not fast_mode and wav_seconds is not None:
                    if used_device_now == "cpu" and wav_seconds >= 20 * 60:
                        auto_fast_mode = True
                        _safe_status("Auto fast mode: long audio on CPU")
                    elif used_device_now == "cuda" and wav_seconds >= 60 * 60:
                        auto_fast_mode = True
                        _safe_status("Auto fast mode: very long audio on CUDA")

                effective_fast_mode = bool(fast_mode or auto_fast_mode)

                vad_parameters = {"min_silence_duration_ms": 500}
                if wav_seconds is not None and wav_seconds >= 20 * 60:
                    vad_parameters = {"min_silence_duration_ms": 800, "speech_pad_ms": 150}
                fast_chunk_len = 60
                if effective_fast_mode:
                    fast_chunk_len = 30 if used_device_now == "cuda" else 60
                    vad_parameters = {"min_silence_duration_ms": 1000, "speech_pad_ms": 80}
                    _safe_status(f"Fast mode enabled: chunk_length={fast_chunk_len}, timestamps=off")

                transcribe_kwargs = {
                    "beam_size": 1,
                    "best_of": 1,
                    "temperature": 0.0,
                    "language": lang,
                    "condition_on_previous_text": False,
                    "vad_filter": True,
                    "vad_parameters": vad_parameters,
                    "without_timestamps": True,
                    "batch_size": int(batch_size),
                }
                if effective_fast_mode:
                    transcribe_kwargs["chunk_length"] = fast_chunk_len
                _safe_status(f"Whisper参数: device={used_device_now}, batch={batch_size}, chunk={transcribe_kwargs.get('chunk_length', 'auto')}, vad={vad_parameters.get('min_silence_duration_ms')}")
                try:
                    import inspect

                    accepted = set(inspect.signature(model.transcribe).parameters.keys())
                    dropped = sorted([k for k in transcribe_kwargs.keys() if k not in accepted])
                    transcribe_kwargs = {k: v for k, v in transcribe_kwargs.items() if k in accepted}
                    if dropped:
                        _safe_status(f"faster-whisper 参数兼容：已忽略 {', '.join(dropped)}")
                except Exception:
                    pass

                q: "queue.Queue[object]" = queue.Queue()
                worker = threading.Thread(target=_transcribe_in_worker, args=(q, transcribe_kwargs), daemon=True)
                started_at = time.time()
                worker.start()

                last_progress_at = started_at
                last_seg_count = 0
                last_end_s: float | None = None
                last_heartbeat_at = started_at
                max_total_seconds = max(30.0 * 60.0, float(wav_seconds) * 8.0) if wav_seconds is not None else 30.0 * 60.0
                stall_seconds = 10.0 * 60.0
                if str(device_for_status).lower() == "cuda":
                    max_total_seconds = max(60.0 * 60.0, float(wav_seconds) * 12.0) if wav_seconds is not None else 60.0 * 60.0
                    stall_seconds = max(2.0 * 60.0, float(wav_seconds) * 0.06) if wav_seconds is not None else 2.0 * 60.0
                    stall_cap = max_total_seconds * 0.4
                    if stall_seconds > stall_cap:
                        stall_seconds = stall_cap

                result_text: str | None = None
                info = None
                while True:
                    try:
                        msg = q.get(timeout=2.0)
                    except Exception:
                        msg = None

                    now = time.time()
                    if now - last_heartbeat_at >= 30.0:
                        hb = f"Whisper heartbeat: {int(now - started_at)}s elapsed"
                        if wav_seconds is not None:
                            hb += f", audio {wav_seconds:.1f}s"
                        if last_seg_count:
                            hb += f", segments {last_seg_count}"
                        if last_end_s is not None:
                            hb += f", last_end {last_end_s:.1f}s"
                        _safe_status(hb)
                        last_heartbeat_at = now

                    if msg is None:
                        if now - started_at > max_total_seconds:
                            raise TimeoutError(f"Whisper transcribe timeout after {int(now - started_at)}s")
                        if last_progress_at and now - last_progress_at > stall_seconds:
                            raise TimeoutError(f"Whisper appears stalled (no progress for {int(now - last_progress_at)}s)")
                        continue

                    kind = msg[0]
                    if kind == "progress":
                        payload = msg[1] if len(msg) > 1 else {}
                        last_seg_count = int(payload.get("segments") or last_seg_count)
                        try:
                            last_end_s = float(payload.get("end_s")) if payload.get("end_s") is not None else last_end_s
                        except Exception:
                            pass
                        last_progress_at = now
                        if wav_seconds is not None and last_end_s is not None:
                            pct = max(0.0, min(100.0, (last_end_s / wav_seconds) * 100.0))
                            _safe_status(f"Whisper progress: {pct:.1f}% ({last_end_s:.1f}/{wav_seconds:.1f}s), segments={last_seg_count}")
                        else:
                            _safe_status(f"Whisper progress: segments={last_seg_count}")
                    elif kind == "done":
                        payload = msg[1] if len(msg) > 1 else {}
                        result_text = str(payload.get("text") or "").strip()
                        info = payload.get("info")
                        break
                    elif kind == "error":
                        err = msg[1] if len(msg) > 1 else RuntimeError("Unknown transcribe error")
                        raise err

                text = (result_text or "").strip()
            except TimeoutError as te:
                if device_for_status == "cuda":
                    _safe_status(f"CUDA transcribe stalled, fallback to CPU. Reason: {te}")
                    cache[f"{cache_key}_gpu_failed_reason"] = f"CUDA stalled: {te}"
                    cache[f"{cache_key}_gpu_failed"] = True
                    cache[f"{cache_key}_gpu_retry_once"] = True
                    cpu_threads = os.cpu_count() or 4
                    model_dir = str(get_models_dir())
                    os.makedirs(model_dir, exist_ok=True)
                    model = WhisperModel(model_size, device="cpu", compute_type="int8", cpu_threads=cpu_threads, download_root=model_dir)
                    cache[cache_key] = model
                    cache[f"{cache_key}_info"] = {"device": "cpu", "compute_type": "int8"}
                    cache[f"{cache_key}_gpu_failed"] = True
                    device_for_status = "cpu"
                    _safe_status("Retry transcribe on CPU (int8)...")
                    used_device_now = "cpu"
                    batch_size = 16 if (os.cpu_count() or 4) >= 8 else 8
                    if model_size not in {"tiny", "base", "small"}:
                        batch_size = min(8, batch_size)
                    transcribe_kwargs = {
                        "beam_size": 1,
                        "best_of": 1,
                        "temperature": 0.0,
                        "language": lang,
                        "condition_on_previous_text": False,
                        "vad_filter": True,
                        "vad_parameters": vad_parameters,
                        "without_timestamps": True,
                        "batch_size": int(batch_size),
                    }
                    if effective_fast_mode:
                        transcribe_kwargs["chunk_length"] = 60
                    try:
                        import inspect

                        accepted = set(inspect.signature(model.transcribe).parameters.keys())
                        dropped = sorted([k for k in transcribe_kwargs.keys() if k not in accepted])
                        transcribe_kwargs = {k: v for k, v in transcribe_kwargs.items() if k in accepted}
                        if dropped:
                            _safe_status(f"faster-whisper 参数兼容：已忽略 {', '.join(dropped)}")
                    except Exception:
                        pass

                    q = queue.Queue()
                    worker = threading.Thread(target=_transcribe_in_worker, args=(q, transcribe_kwargs), daemon=True)
                    started_at = time.time()
                    worker.start()

                    last_progress_at = started_at
                    last_seg_count = 0
                    last_end_s = None
                    last_heartbeat_at = started_at
                    max_total_seconds = max(60.0 * 60.0, float(wav_seconds) * 12.0) if wav_seconds is not None else 60.0 * 60.0
                    stall_seconds = 5.0 * 60.0  # CPU fallback: 5 分钟无进度则认定卡死

                    result_text = None
                    info = None
                    while True:
                        try:
                            msg = q.get(timeout=2.0)
                        except Exception:
                            msg = None

                        now = time.time()
                        if now - last_heartbeat_at >= 30.0:
                            hb = f"Whisper heartbeat: {int(now - started_at)}s elapsed"
                            if wav_seconds is not None:
                                hb += f", audio {wav_seconds:.1f}s"
                            if last_seg_count:
                                hb += f", segments {last_seg_count}"
                            if last_end_s is not None:
                                hb += f", last_end {last_end_s:.1f}s"
                            _safe_status(hb)
                            last_heartbeat_at = now

                        if msg is None:
                            if now - started_at > max_total_seconds:
                                raise TimeoutError(f"Whisper transcribe timeout after {int(now - started_at)}s")
                            if last_progress_at and now - last_progress_at > stall_seconds:
                                raise TimeoutError(f"Whisper appears stalled (no progress for {int(now - last_progress_at)}s)")
                            continue

                        kind = msg[0]
                        if kind == "progress":
                            payload = msg[1] if len(msg) > 1 else {}
                            last_seg_count = int(payload.get("segments") or last_seg_count)
                            try:
                                last_end_s = float(payload.get("end_s")) if payload.get("end_s") is not None else last_end_s
                            except Exception:
                                pass
                            last_progress_at = now
                            if wav_seconds is not None and last_end_s is not None:
                                pct = max(0.0, min(100.0, (last_end_s / wav_seconds) * 100.0))
                                _safe_status(f"Whisper progress: {pct:.1f}% ({last_end_s:.1f}/{wav_seconds:.1f}s), segments={last_seg_count}")
                            else:
                                _safe_status(f"Whisper progress: segments={last_seg_count}")
                        elif kind == "done":
                            payload = msg[1] if len(msg) > 1 else {}
                            result_text = str(payload.get("text") or "").strip()
                            info = payload.get("info")
                            break
                        elif kind == "error":
                            err = msg[1] if len(msg) > 1 else RuntimeError("Unknown transcribe error")
                            raise err

                    text = (result_text or "").strip()
                else:
                    raise
            
            # 获取实际使用的 device info
            dev_info = cache.get(f"{cache_key}_info", {})
            used_device = dev_info.get("device", "cpu")
            used_compute = dev_info.get("compute_type", "int8")
            
            # 附加 device 信息到文本末尾 (隐式传递给 UI)
            # 格式：<!-- FW_DEVICE: device_name (compute_type) -->
            debug_extra = ""
            if used_device != "cuda":
                gpu_failed_reason = cache.get(f"{cache_key}_gpu_failed_reason")
                if gpu_failed_reason:
                    reason_clean = str(gpu_failed_reason).replace("\n", " ").strip()
                    if len(reason_clean) > 140:
                        reason_clean = reason_clean[:140] + "..."
                    debug_extra = f" | GPU_FAIL: {reason_clean}"
            debug_tag = f"\n\n<!-- FW_DEVICE: {used_device.upper()} ({used_compute}){debug_extra} -->"
            
            if not text:
                 raise RuntimeError(f"faster-whisper 转写结果为空。Device: {used_device}")
            return text + debug_tag

        finally:
            # 恢复环境变量
            if proxy_url:
                if original_http:
                    os.environ["HTTP_PROXY"] = original_http
                else:
                    os.environ.pop("HTTP_PROXY", None)
                    
                if original_https:
                    os.environ["HTTPS_PROXY"] = original_https
                else:
                    os.environ.pop("HTTPS_PROXY", None)

    except ImportError:
        # 如果没有安装 faster-whisper，打印警告并回退到 openai-whisper
        print("Warning: 'faster-whisper' not installed. Using slower 'openai-whisper'. Install faster-whisper for 4x speedup.")
        pass
    except Exception as e:
        # faster-whisper 失败，回退到 openai-whisper
        fw_error = str(e)
        print(f"faster-whisper failed: {e}. Falling back to openai-whisper.")
        import traceback
        traceback.print_exc()
        pass

    # --- 以下是原有的 openai-whisper 逻辑 (作为兜底) ---
    try:
        import whisper  # type: ignore
    except Exception as e:
        msg = "未安装 openai-whisper，无法启用音频转写兜底。"
        if fw_error:
            msg += f"\n(此前 faster-whisper 也失败了: {fw_error})"
        raise RuntimeError(msg) from e
        
    print("Fallback to openai-whisper (Legacy)...")
    
    # 尝试获取 device (虽然 openai-whisper 自动处理，但我们可以尝试检测)
    import torch
    used_device = "CUDA" if torch.cuda.is_available() else "CPU"
    used_compute = "float16" if used_device == "CUDA" else "float32" # 简单假设


    model_name = (model_name or "").strip() or "base"
    language = (language or "").strip().lower()
    lang = None if (not language or language == "auto") else language

    cache = getattr(transcribe_audio_with_whisper, "_model_cache", None)
    if not isinstance(cache, dict):
        cache = {}
        setattr(transcribe_audio_with_whisper, "_model_cache", cache)
    model = cache.get(model_name)
    if model is None:
        model = whisper.load_model(model_name)
        cache[model_name] = model
    
    # 定义重试参数列表
    # 策略优化：优先使用 Greedy (beam_size=1)，极大提升速度 (3-5x)
    # 如果 Greedy 失败，再尝试更稳健的 Beam Search 或 强力兜底
    
    retry_configs = [
        # 1. 极速模式：Greedy Decoding
        {"fp16": False, "language": lang, "beam_size": 1}, 
        # 2. 强力兜底：禁用所有过滤，强制输出
        {"fp16": False, "language": lang, "logprob_threshold": None, "compression_ratio_threshold": None, "condition_on_previous_text": False, "no_speech_threshold": 0.95},
    ]

    last_exc = None
    for i, cfg in enumerate(retry_configs):
        try:
            # print(f"Whisper attempt {i+1} with config: {cfg}")
            result = model.transcribe(audio_path, **cfg)
            text = str((result or {}).get("text") or "").strip()
            if text:
                 # 附加 openai-whisper 的 device info
                extra_msg = f" | ⚠️ faster-whisper 启动失败: {fw_error}" if fw_error else ""
                debug_tag = f"\n\n<!-- FW_DEVICE: {used_device} ({used_compute}) [OpenAI-Whisper]{extra_msg} -->"
                return text + debug_tag
        except Exception as e:
            last_exc = e
            # 如果是 KeyError (Linear)，直接抛出，无法重试修复
            if isinstance(e, KeyError) and "Linear" in str(e):
                raise KeyError(f"Whisper 模型加载/推理时发生 KeyError: {e}。这通常是 pytorch/whisper 版本不兼容导致。") from e
            # 继续尝试下一个配置

    # 所有重试都失败
    if last_exc:
        msg = str(last_exc)
        # 获取文件大小信息，辅助 debug
        file_info = ""
        try:
            sz = os.path.getsize(audio_path)
            file_info = f" (音频文件大小: {sz} bytes)"
        except:
            pass
            
        # 附加 device 信息到错误消息 (隐式传递给 UI)
        # 既然 faster-whisper 的 device info 很有用，这里也尝试获取一下
        # 虽然这里是 openai-whisper，但我们可以简单返回 CPU
        # 或者我们统一异常格式
        
        if "cannot reshape tensor of 0 elements" in msg:
            raise RuntimeError(f"音频转写失败：未能检测到有效语音片段 (No speech detected)。请检查视频是否静音，或尝试更小的模型{file_info}。") from last_exc
        raise last_exc
    
    raise RuntimeError("音频转写结果为空 (所有重试方案均未产生文本)。")



def transcribe_video_audio_with_ytdlp(
    video_url: str,
    proxy_url: str,
    timeout_seconds: float,
    retries: int,
    cookies_file: str,
    cookies_from_browser: str,
    model_name: str,
    language: str,
    status_callback=None,
    fast_mode: bool = False,
    cookie_debug_summary: str = "",
) -> tuple[str, str]:
    class YdlLogger:
        def __init__(self) -> None:
            self.lines: list[str] = []

        def _add(self, level: str, msg: str) -> None:
            s = strip_ansi(str(msg or "")).strip()
            if not s:
                return
            self.lines.append(f"[{level}] {s}")
            if len(self.lines) > 400:
                self.lines = self.lines[-250:]

        def debug(self, msg: str) -> None:
            self._add("debug", msg)

        def info(self, msg: str) -> None:
            self._add("info", msg)

        def warning(self, msg: str) -> None:
            self._add("warning", msg)

        def error(self, msg: str) -> None:
            self._add("error", msg)

    try:
        from yt_dlp import YoutubeDL  # type: ignore
        from yt_dlp.utils import DownloadError  # type: ignore
    except Exception as e:
        raise RuntimeError(
            "未安装 yt-dlp，无法启用音频转写兜底。可执行：python -m pip install yt-dlp -i https://pypi.tuna.tsinghua.edu.cn/simple"
        ) from e

        # 移除 detect_js_runtime 的定义，改用 yt-dlp 原生解释器或 yt_dlp_ejs
        # 移除此处的本地 is_cookie_error 定义，改用顶层的 _is_cookie_error

    normalized_video_url = normalize_video_url(video_url)
    parsed_video_host = (urlparse(normalized_video_url).netloc or "").lower()
    is_youtube_url = ("youtube.com" in parsed_video_host) or ("youtu.be" in parsed_video_host)

    cache_path = _audio_cache_path(normalized_video_url)
    if cache_path and cache_path.exists():
        try:
            if cache_path.stat().st_size > 16 * 1024:
                if status_callback:
                    status_callback("检测到音频缓存，跳过下载")
                force_cpu_flag = bool(getattr(transcribe_video_audio_with_ytdlp, "_force_cpu", False))
                text = transcribe_audio_with_whisper(str(cache_path), model_name=model_name, language=language, proxy_url=proxy_url, status_callback=status_callback, fast_mode=fast_mode, force_cpu=force_cpu_flag)
                label = f"{_audio_cache_key(normalized_video_url)} | whisper:{(model_name or '').strip() or 'base'}"
                return label, text
        except Exception:
            pass

    with tempfile.TemporaryDirectory() as tmp:
        outtmpl = os.path.join(tmp, "%(id)s.%(ext)s")
        last_err: Exception | None = None
        last_cookie_error: RuntimeError | None = None
        last_video_id = ""
        last_attempt_note = ""
        last_debug_lines: list[str] = []
        
        # 记录无 Cookie 时的错误
        no_cookie_error: Exception | None = None
        force_browser_cookie = False

        js_runtime_available, js_runtime_name = detect_js_runtime()
        running_on_render = bool(str(os.environ.get("RENDER_SERVICE_ID", "") or "").strip())
        web_clients_enabled = bool(js_runtime_available and not running_on_render)
        has_cookie_hint = bool((cookies_file or "").strip()) or bool((cookies_from_browser or "").strip())
        if is_youtube_url:
            client_strategies = [["tv"], ["web_safari"]]
            if not has_cookie_hint:
                client_strategies.extend([["ios"], ["android"]])
            if not running_on_render:
                client_strategies.append(["mweb"])
            if web_clients_enabled:
                client_strategies.extend([["web_creator"], []])
        else:
            client_strategies = [[]]
        client_plan = ", ".join("+".join(cs) if cs else "default" for cs in client_strategies)

        if fast_mode:
            format_candidates = ["worstaudio/worst", "bestaudio/best", None]
        else:
            format_candidates = ["bestaudio/best", "worstaudio/worst", None]
        
        disabled_browsers: set[str] = set()
        disabled_clients_reason: dict[str, str] = {}
        cookie_file_exists = bool((cookies_file or "").strip()) and os.path.exists((cookies_file or "").strip())
        cookie_runtime_hint = (
            f"input_cookie_file={'yes' if (cookies_file or '').strip() else 'no'}; "
            f"input_cookie_exists={'yes' if cookie_file_exists else 'no'}; "
            f"input_browser={cookies_from_browser or 'none'}; "
            f"js_runtime_available={'yes' if js_runtime_available else 'no'}; "
            f"js_runtime={js_runtime_name}; "
            f"running_on_render={'yes' if running_on_render else 'no'}; "
            f"web_clients_enabled={'yes' if web_clients_enabled else 'no'}; "
            f"client_plan={client_plan}"
        )
        selected_audio_summary = ""
        runtime_version_summary = build_runtime_version_diagnostics()
        
        start_time = time.time()
        requested_paths: list[Path] = []
        download_ready = False
        last_err_type = ""
        last_traceback_text = ""
        
        for client_set in client_strategies:
            for attempt in range(max(1, int(retries)) + 1):
                for cookiefile, cfb in CookieManager.get_sources(cookies_file, cookies_from_browser, force_browser_cookie):
                    if cfb and cfb in disabled_browsers:
                        continue
                    for fmt in format_candidates:
                        if download_ready:
                            break
                        client_label = ",".join(client_set) if client_set else "default"
                        attempt_note = f"client={client_label} fmt={(fmt or 'default')} cookie={'file' if cookiefile else ('browser' if cfb else 'none')}"

                        def progress_hook(d):
                            if status_callback:
                                if d["status"] == "downloading":
                                    p = strip_ansi(str(d.get("_percent_str") or "")).replace("%", "")
                                    e = strip_ansi(str(d.get("_eta_str") or ""))
                                    status_callback(f"Downloading audio: {p}% (ETA: {e})")
                                elif d["status"] == "finished":
                                    status_callback("Download complete, preparing transcription...")

                        logger = YdlLogger()
                        opts: dict = {
                            "progress_hooks": [progress_hook],
                            "outtmpl": outtmpl,
                            "noplaylist": True,
                            "quiet": True,
                            "no_warnings": True,
                            "verbose": True,
                            "nocheckcertificate": True,
                            "socket_timeout": float(timeout_seconds),
                            "retries": 1,
                            "ignoreerrors": False,
                            "ignore_config": True,
                            "http_headers": {"Accept-Language": "en-US,en;q=0.9"},
                            "logger": logger,
                        }
                        if is_youtube_url and client_set:
                            opts["extractor_args"] = {"youtube": {"player_client": client_set}}
                        if ffmpeg_binary_path:
                            opts["ffmpeg_location"] = ffmpeg_binary_path
                        if proxy_url:
                            opts["proxy"] = proxy_url
                        if cookiefile:
                            opts["cookiefile"] = cookiefile
                        elif cfb:
                            opts["cookiesfrombrowser"] = (cfb,)

                        try:
                            with YoutubeDL(opts) as ydl:
                                def pick_audio_format_entry(info: dict, prefer_worst: bool, preferred_ext: str | None = None) -> dict | None:
                                    formats = info.get("formats") if isinstance(info, dict) else None
                                    if not formats:
                                        return None
                                    audio_formats = []
                                    for item in formats:
                                        if not isinstance(item, dict):
                                            continue
                                        acodec = item.get("acodec")
                                        if not acodec or acodec == "none":
                                            continue
                                        if preferred_ext:
                                            ext = str(item.get("ext") or "").lower()
                                            if ext != preferred_ext.lower():
                                                continue
                                        audio_formats.append(item)
                                    if not audio_formats:
                                        return None

                                    def score(item: dict) -> float:
                                        for key in ("abr", "tbr"):
                                            val = item.get(key)
                                            if isinstance(val, (int, float)):
                                                return float(val)
                                        return 0.0

                                    return min(audio_formats, key=score) if prefer_worst else max(audio_formats, key=score)

                                def pick_audio_format(info: dict, prefer_worst: bool, preferred_ext: str | None = None) -> str | None:
                                    picked = pick_audio_format_entry(info, prefer_worst, preferred_ext)
                                    fid = picked.get("format_id") if isinstance(picked, dict) else None
                                    return str(fid) if fid else None

                                def has_audio_format(info: dict) -> bool:
                                    return pick_audio_format(info, False, None) is not None or pick_audio_format(info, True, None) is not None

                                def summarize_audio_entry(info: dict, chosen_entry: dict | None) -> str:
                                    formats = info.get("formats") if isinstance(info, dict) else None
                                    audio_count = 0
                                    if isinstance(formats, list):
                                        for item in formats:
                                            if isinstance(item, dict) and item.get("acodec") not in (None, "", "none"):
                                                audio_count += 1
                                    if not isinstance(chosen_entry, dict):
                                        return f"audio_format_count={audio_count}; selected=none"
                                    return (
                                        f"audio_format_count={audio_count}; "
                                        f"selected_format_id={chosen_entry.get('format_id') or 'unknown'}; "
                                        f"selected_ext={chosen_entry.get('ext') or 'unknown'}; "
                                        f"selected_protocol={chosen_entry.get('protocol') or 'unknown'}"
                                    )

                                def direct_download_audio(format_info: dict | None, output_dir: str) -> Path | None:
                                    """在 yt-dlp 没有落地产物时，直接下载可访问的音频流。"""
                                    if not isinstance(format_info, dict):
                                        return None
                                    audio_url = str(format_info.get("url") or "").strip()
                                    if not audio_url:
                                        return None
                                    protocol = str(format_info.get("protocol") or "").lower()

                                    ext = str(format_info.get("ext") or "bin").lower() or "bin"
                                    format_id = str(format_info.get("format_id") or "direct")
                                    base_name = last_video_id or "audio"
                                    target_path = Path(output_dir) / f"{base_name}.{format_id}.{ext}"

                                    session = TimeoutSession(timeout_seconds=max(5.0, float(timeout_seconds)))
                                    session.trust_env = True
                                    if proxy_url:
                                        session.trust_env = False
                                        session.proxies = {"http": proxy_url, "https": proxy_url}

                                    headers = dict(opts.get("http_headers") or {})
                                    headers["Referer"] = normalized_video_url

                                    if "m3u8" in protocol or "dash" in protocol:
                                        if not ffmpeg_binary_path or not os.path.exists(ffmpeg_binary_path):
                                            logger.warning("Direct media download fallback skipped: ffmpeg binary is unavailable for streaming protocol")
                                            return None
                                        ffmpeg_output = Path(output_dir) / f"{base_name}.{format_id}.m4a"
                                        ffmpeg_cmd = [
                                            ffmpeg_binary_path,
                                            "-y",
                                            "-nostdin",
                                            "-loglevel",
                                            "error",
                                        ]
                                        user_agent = str(headers.get("User-Agent") or "").strip()
                                        referer = str(headers.get("Referer") or "").strip()
                                        if user_agent:
                                            ffmpeg_cmd.extend(["-user_agent", user_agent])
                                        if referer:
                                            ffmpeg_cmd.extend(["-headers", f"Referer: {referer}\r\n"])
                                        ffmpeg_cmd.extend([
                                            "-i",
                                            audio_url,
                                            "-vn",
                                            "-acodec",
                                            "copy",
                                            str(ffmpeg_output),
                                        ])
                                        try:
                                            completed = subprocess.run(
                                                ffmpeg_cmd,
                                                capture_output=True,
                                                text=True,
                                                timeout=max(30, int(float(timeout_seconds) * 2)),
                                                check=False,
                                            )
                                            if completed.returncode == 0 and ffmpeg_output.exists() and ffmpeg_output.stat().st_size > 16 * 1024:
                                                logger.info("Direct media download fallback succeeded via ffmpeg streaming capture")
                                                return ffmpeg_output
                                            stderr_tail = (completed.stderr or "").strip()[-500:]
                                            logger.warning(f"Direct media download fallback via ffmpeg failed: {stderr_tail}")
                                        except Exception as e:
                                            logger.warning(f"Direct media download fallback via ffmpeg failed: {e}")
                                        return None

                                    try:
                                        with session.get(audio_url, headers=headers, stream=True, timeout=max(5.0, float(timeout_seconds))) as resp:
                                            resp.raise_for_status()
                                            with open(target_path, "wb") as fp:
                                                for chunk in resp.iter_content(chunk_size=1024 * 512):
                                                    if chunk:
                                                        fp.write(chunk)
                                        if target_path.exists() and target_path.stat().st_size > 16 * 1024:
                                            return target_path
                                    except Exception as e:
                                        logger.warning(f"Direct media download fallback failed: {e}")
                                    return None

                                try:
                                    if client_set and any(c in disabled_clients_reason for c in client_set):
                                        reason = next((disabled_clients_reason.get(c) for c in client_set if c in disabled_clients_reason), "Client disabled")
                                        if last_err is None:
                                            last_err = RuntimeError(reason)
                                        continue

                                    info = ydl.extract_info(normalized_video_url, download=False, process=False)
                                except DownloadError as e:
                                    msg = strip_ansi(str(e))
                                    has_cookie_in_log = any(CookieManager.is_cookie_error(line) for line in logger.lines)
                                    if CookieManager.is_cookie_error(msg) or has_cookie_in_log:
                                        last_cookie_error = RuntimeError(CookieManager.get_fatal_msg(msg, cfb))
                                        if cfb:
                                            disabled_browsers.add(cfb)
                                        continue
                                    if "missing a URL" in msg or "SABR streaming" in msg:
                                        if "web_safari" in client_set:
                                            disabled_clients_reason["web_safari"] = "Web Safari 客户端不兼容 (SABR/Missing URL)。"
                                        last_err = RuntimeError("Web Safari 客户端不兼容。")
                                        last_attempt_note = attempt_note + " (web_safari issue)"
                                        last_debug_lines = logger.lines[-80:]
                                        continue
                                    if "po token" in msg.lower() or "challenge solving failed" in msg.lower() or "only images are available" in msg.lower():
                                        force_browser_cookie = True
                                        last_err = RuntimeError("YouTube 触发挑战校验（PO Token/JS Challenge）。已自动切换为浏览器 Cookie 方案重试。")
                                        last_attempt_note = attempt_note + " (challenge/po token)"
                                        last_debug_lines = logger.lines[-80:]
                                        continue
                                    if has_login_required([], msg):
                                        force_browser_cookie = True
                                        for c in client_set:
                                            if c in {"tv", "tv_embedded"}:
                                                disabled_clients_reason[c] = "需要登录或验证（可能触发人机校验）。已尝试自动读取浏览器 Cookie；如仍失败请手动开启或提供 cookies 文件。"
                                        last_err = RuntimeError("需要登录或验证（可能触发人机校验）。已尝试自动读取浏览器 Cookie；如仍失败请手动开启或提供 cookies 文件。")
                                        last_attempt_note = attempt_note + " (login required)"
                                        last_debug_lines = logger.lines[-80:]
                                        continue
                                    if "requested format not available" in msg.lower():
                                        last_err = e
                                        last_attempt_note = attempt_note
                                        last_debug_lines = logger.lines[-80:]
                                        continue
                                    if "drm protected" in msg.lower():
                                        if "tv" in client_set:
                                            disabled_clients_reason["tv"] = "TV 客户端遭遇 DRM 保护限制。"
                                        last_err = RuntimeError("当前客户端遭遇 DRM 保护限制。")
                                        last_attempt_note = attempt_note + " (DRM protected)"
                                        last_debug_lines = logger.lines[-80:]
                                        continue
                                    raise e

                                try:
                                    if has_login_required(logger.lines):
                                        force_browser_cookie = True
                                        for c in client_set:
                                            if c in {"tv", "tv_embedded"}:
                                                disabled_clients_reason[c] = "需要登录或验证（可能触发人机校验）。已尝试自动读取浏览器 Cookie；如仍失败请手动开启或提供 cookies 文件。"
                                        last_err = RuntimeError("需要登录或验证（可能触发人机校验）。已尝试自动读取浏览器 Cookie；如仍失败请手动开启或提供 cookies 文件。")
                                        last_attempt_note = attempt_note + " (login required)"
                                        last_debug_lines = logger.lines[-80:]
                                        continue
                                    if has_po_token_required(logger.lines) and any(c in {"android", "ios", "mweb"} for c in client_set):
                                        force_browser_cookie = True
                                        for c in client_set:
                                            if c in {"android", "ios", "mweb"}:
                                                disabled_clients_reason[c] = "该客户端需要 PO Token，已降级到其他客户端。"
                                        last_err = RuntimeError("该客户端需要 PO Token，已降级到其他客户端。")
                                        last_attempt_note = attempt_note + " (po token required)"
                                        last_debug_lines = logger.lines[-80:]
                                        continue
                                    if has_js_challenge_failure(logger.lines):
                                        web_like_client = (not client_set) or any(c in {"web", "web_creator", "web_safari"} for c in client_set)
                                        if web_like_client:
                                            force_browser_cookie = True
                                            if "web" in client_set or "web_creator" in client_set:
                                                disabled_clients_reason["web"] = "JS challenge 失败，可能导致格式缺失。建议升级 yt-dlp 或更换网络。"
                                            last_err = RuntimeError("JS challenge 失败，可能导致格式缺失。建议升级 yt-dlp 或更换网络。")
                                            last_attempt_note = attempt_note + " (js challenge failed)"
                                            last_debug_lines = logger.lines[-80:]
                                            continue
                                    if not has_audio_format(info):
                                        if "web" in client_set:
                                            disabled_clients_reason["web"] = "未检测到可用音频格式。建议升级 yt-dlp 或更换网络。"
                                        last_err = RuntimeError("未检测到可用音频格式。建议升级 yt-dlp 或更换网络。")
                                        last_attempt_note = attempt_note + " (no audio formats)"
                                        last_debug_lines = logger.lines[-80:]
                                        continue

                                    selected_audio_entry = pick_audio_format_entry(info, False, None)
                                    selected_audio_summary = summarize_audio_entry(info, selected_audio_entry)
                                    if fmt:
                                        prefer_worst = fmt in {"worstaudio/worst", "worstaudio", "worst"}
                                        chosen = pick_audio_format(info, prefer_worst, None)
                                        selected_audio_entry = pick_audio_format_entry(info, prefer_worst, None)
                                        selected_audio_summary = summarize_audio_entry(info, selected_audio_entry)
                                        if chosen:
                                            ydl.params["format"] = chosen
                                        else:
                                            last_err = RuntimeError("Requested format is not available")
                                            last_attempt_note = attempt_note
                                            last_debug_lines = logger.lines[-80:]
                                            continue

                                    requested_paths = []
                                    direct_file = direct_download_audio(selected_audio_entry, tmp)
                                    if direct_file:
                                        requested_paths.append(direct_file)
                                        if isinstance(info, dict):
                                            last_video_id = str(info.get("id") or last_video_id)
                                    else:
                                        info = ydl.extract_info(normalized_video_url, download=True)
                                        if isinstance(info, dict):
                                            last_video_id = str(info.get("id") or last_video_id)
                                            requested_downloads = info.get("requested_downloads") or []
                                            if isinstance(requested_downloads, list):
                                                for item in requested_downloads:
                                                    if isinstance(item, dict):
                                                        fp = item.get("filepath")
                                                        if fp:
                                                            requested_paths.append(Path(str(fp)))
                                            requested_formats = info.get("requested_formats") or []
                                            if isinstance(requested_formats, list):
                                                for item in requested_formats:
                                                    if isinstance(item, dict):
                                                        fp = item.get("filepath")
                                                        if fp:
                                                            requested_paths.append(Path(str(fp)))
                                            direct_fp = info.get("filepath")
                                            if direct_fp:
                                                requested_paths.append(Path(str(direct_fp)))

                                        if not requested_paths:
                                            if isinstance(selected_audio_entry, dict):
                                                selected_protocol = str(selected_audio_entry.get("protocol") or "").lower() or "unknown"
                                                selected_ext = str(selected_audio_entry.get("ext") or "").lower() or "unknown"
                                                last_err = RuntimeError(f"直链下载与 yt-dlp 下载均未产出文件（protocol={selected_protocol}, ext={selected_ext}）")
                                                last_err_type = type(last_err).__name__
                                                last_attempt_note = attempt_note + " (no output files)"
                                                last_debug_lines = logger.lines[-80:]
                                                continue
                                except Exception as e:
                                    last_err = e
                                    last_err_type = type(e).__name__
                                    last_traceback_text = traceback.format_exc()
                                    last_attempt_note = attempt_note
                                    last_debug_lines = logger.lines[-80:]
                                    continue

                                last_err = None
                                last_err_type = ""
                                last_traceback_text = ""
                                download_ready = True
                                break
                        except DownloadError as e:
                            msg = strip_ansi(str(e))
                            has_cookie_in_log = any(CookieManager.is_cookie_error(line) for line in logger.lines)
                            if CookieManager.is_cookie_error(msg) or has_cookie_in_log:
                                last_cookie_error = RuntimeError(CookieManager.get_fatal_msg(msg, cfb))
                                if cfb:
                                    disabled_browsers.add(cfb)
                                continue
                            if ("ejs" in msg.lower()) or ("challenge solving failed" in msg.lower()) or ("js runtimes: none" in msg.lower()) or ("only images are available" in msg.lower()):
                                disabled_clients_reason["web"] = "Web 客户端挑战失败，已禁用。"
                                last_attempt_note = attempt_note + " (EJS/JS runtime issue, web disabled)"
                                last_debug_lines = logger.lines[-80:]
                                continue
                            if not cookiefile and not cfb:
                                no_cookie_error = e
                                last_err = e
                            else:
                                last_err = no_cookie_error if no_cookie_error else e
                            last_err_type = type(last_err).__name__ if last_err else ""
                            last_traceback_text = traceback.format_exc()
                            last_attempt_note = attempt_note
                            last_debug_lines = logger.lines[-80:]
                            if "HTTP Error 429" in msg or "429" in msg:
                                time.sleep(2.0 * (attempt + 1))
                            continue
                        except Exception as e:
                            if not cookiefile and not cfb:
                                no_cookie_error = e
                                last_err = e
                            else:
                                if no_cookie_error:
                                    last_err = no_cookie_error
                                elif not CookieManager.is_cookie_error(str(e)):
                                    last_err = e
                            last_err_type = type(last_err).__name__ if last_err else ""
                            last_traceback_text = traceback.format_exc()
                            last_attempt_note = attempt_note
                            last_debug_lines = logger.lines[-80:]
                            continue

                    if download_ready:
                        break

                if last_err is None:
                    # 成功获取到信息（且下载成功），不需要再重试其他 client
                    pass

                media_suffixes = {".m4a", ".webm", ".mp3", ".wav", ".opus", ".aac", ".flac", ".ogg", ".mp4", ".mkv", ".mov", ".m4v", ".ts", ".m4s"}
                candidates = []
                temp_outputs = []
                for p in Path(tmp).rglob("*"):
                    if not p.is_file():
                        continue
                    temp_outputs.append(f"{p.name} ({p.stat().st_size} bytes)")
                    ext = p.suffix.lower()
                    if ext in media_suffixes:
                        candidates.append(p)
                if not candidates:
                    for p in requested_paths:
                        try:
                            if p.is_file() and p.suffix.lower() in media_suffixes:
                                candidates.append(p)
                        except Exception:
                            continue
                if not candidates:
                    detail_lines = []
                    if last_attempt_note:
                        detail_lines.append(f"最近尝试: {last_attempt_note}")
                    if last_err:
                        err_text = strip_ansi(str(last_err))
                        if err_text:
                            detail_lines.append(f"上一次错误: {err_text}")
                        else:
                            detail_lines.append(f"上一次错误: <empty message>; type={last_err_type or type(last_err).__name__}; repr={repr(last_err)}")
                    if last_traceback_text:
                        trace_tail = "\n".join(last_traceback_text.strip().splitlines()[-12:])
                        detail_lines.append("异常堆栈(尾部):\n" + trace_tail)
                    if last_debug_lines:
                        debug_seen: set[str] = set()
                        debug_tail: list[str] = []
                        for item in reversed(last_debug_lines[-80:]):
                            k = item.strip().lower()
                            if not k or k in debug_seen:
                                continue
                            debug_seen.add(k)
                            debug_tail.append(item.strip())
                            if len(debug_tail) >= 12:
                                break
                        debug_tail.reverse()
                        tail = "\n".join(debug_tail)
                        detail_lines.append("调试日志(尾部):\n" + tail)
                    if temp_outputs:
                        detail_lines.append("临时目录文件:\n" + "\n".join(temp_outputs[:12]))
                    if requested_paths:
                        detail_lines.append("yt-dlp 报告的输出路径:\n" + "\n".join(str(p) for p in requested_paths[:12]))
                    if selected_audio_summary:
                        detail_lines.append("音频格式诊断:\n" + selected_audio_summary)
                    if cookie_debug_summary:
                        detail_lines.append("Cookies 运行时诊断:\n" + cookie_debug_summary)
                    detail_lines.append("Cookies 传参诊断:\n" + cookie_runtime_hint)
                    detail_lines.append("运行版本诊断:\n" + runtime_version_summary)
                    detail_text = "\n".join(detail_lines)
                    if last_err and str(last_err).strip().startswith("未下载到音频文件"):
                        pass
                    elif detail_text:
                        last_err = RuntimeError("未下载到音频文件（可能被限制或链接无效）。\n" + detail_text)
                    else:
                        last_err = RuntimeError("未下载到音频文件（可能被限制或链接无效）。")
                    continue
                audio = max(candidates, key=lambda x: x.stat().st_size)
                cached_audio = None
                
                # 记录下载耗时
                download_duration = time.time() - start_time
                if status_callback: 
                    status_callback(f"Download finished in {download_duration:.1f}s. Audio extracted. Starting Whisper transcription...")
                
                if cache_path:
                    try:
                        if (not cache_path.exists()) or (cache_path.stat().st_size < audio.stat().st_size):
                            shutil.copy2(audio, cache_path)
                        if cache_path.exists() and cache_path.stat().st_size > 16 * 1024:
                            cached_audio = cache_path
                    except Exception:
                        cached_audio = None
                
                force_cpu_flag = bool(getattr(transcribe_video_audio_with_ytdlp, "_force_cpu", False))
                audio_path = str(cached_audio or audio)
                
                t_transcribe_start = time.time()
                text = transcribe_audio_with_whisper(audio_path, model_name=model_name, language=language, proxy_url=proxy_url, status_callback=status_callback, fast_mode=fast_mode, force_cpu=force_cpu_flag)
                transcribe_duration = time.time() - t_transcribe_start
                
                # 将耗时信息注入到 text 末尾的注释中，供 UI 解析
                timing_tag = f"<!-- TIMING: download={download_duration:.1f}, transcribe={transcribe_duration:.1f} -->"
                text += timing_tag
                
                label = f"{last_video_id or (cached_audio.stem if cached_audio else audio.stem)} | whisper:{(model_name or '').strip() or 'base'}"
                return label, text

        if last_err is not None:
            raise last_err
        if last_cookie_error:
            raise last_cookie_error
        raise RuntimeError("音频转写兜底失败（未知原因）。")


def fetch_available_models(api_key: str, base_url: str, proxy_url: str = None) -> list[str]:
    error_msg = ""
    # 1. 尝试使用 OpenAI SDK
    try:
        from openai import OpenAI
        import httpx
        
        client_kwargs = {"api_key": api_key}
        if base_url and base_url.strip():
            client_kwargs["base_url"] = base_url.strip()
        
        httpx_kwargs = {"verify": False}
        if proxy_url and proxy_url.strip():
            httpx_kwargs["proxy"] = proxy_url.strip()
        
        client_kwargs["http_client"] = httpx.Client(**httpx_kwargs)
        
        client = OpenAI(**client_kwargs)
        models_page = client.models.list()
        
        model_ids = []
        for m in models_page:
            if hasattr(m, "id"):
                model_ids.append(m.id)
            elif isinstance(m, dict):
                model_ids.append(m.get("id"))
                
        return sorted(model_ids)
    except Exception as e:
        error_msg = str(e)
    
    # 2. 如果 SDK 失败，尝试直接 Requests 请求 (兜底)
    try:
        import requests
        headers = {"Authorization": f"Bearer {api_key}"}
        
        target_url = base_url.strip()
        if not target_url.endswith("/"):
            target_url += "/"
        
        # 尝试猜测 models 端点
        if "v1" not in target_url:
            target_url += "v1/"
        target_url += "models"
        
        proxies = None
        if proxy_url and proxy_url.strip():
            proxies = {"http": proxy_url, "https": proxy_url}
            
        r = requests.get(target_url, headers=headers, proxies=proxies, verify=False, timeout=10)
        r.raise_for_status()
        data = r.json()
        
        model_ids = []
        data_list = data.get("data", [])
        if isinstance(data_list, list):
            for m in data_list:
                if isinstance(m, dict):
                    model_ids.append(m.get("id"))
        
        if model_ids:
            # 简单过滤，排除 clearly non-chat models
            filtered = [
                m for m in model_ids 
                if not any(x in m for x in ["dall-e", "tts", "whisper", "embedding", "babbage", "davinci", "curie", "ada"])
            ]
            return sorted(filtered) if filtered else sorted(model_ids)
            
    except Exception as e2:
        error_msg += f" | Direct HTTP failed: {e2}"

    raise RuntimeError(f"获取模型列表失败: {error_msg}")

AUTHORITATIVE_SOURCE_RULES = [
    {
        "label": "美国经济分析局",
        "url": "https://www.bea.gov/",
        "aliases": [
            "美国经济分析局",
            "bureau of economic analysis",
            "u.s. bureau of economic analysis",
            "bea",
            "pce",
            "个人消费支出物价指数",
            "个人消费支出价格指数",
            "美国pce",
            "美国gdp",
            "耐用品订单",
            "durable goods orders",
        ],
        "query_terms": [
            "BEA PCE inflation",
            "BEA GDP release",
            "BEA durable goods orders",
        ],
    },
    {
        "label": "加拿大统计局",
        "url": "https://www.statcan.gc.ca/en/start",
        "aliases": [
            "加拿大统计局",
            "statistics canada",
            "statcan",
            "canada gdp",
            "加拿大gdp",
        ],
        "query_terms": [
            "Statistics Canada GDP",
            "StatCan gross domestic product",
        ],
    },
    {
        "label": "欧盟统计局",
        "url": "https://ec.europa.eu/eurostat",
        "aliases": [
            "欧盟统计局",
            "欧元区通胀",
            "欧元区cpi",
            "eurostat",
            "euro area inflation",
            "eurozone inflation",
        ],
        "query_terms": [
            "Eurostat euro area inflation",
            "Eurostat euro area CPI",
        ],
    },
    {
        "label": "欧洲央行",
        "url": "https://www.ecb.europa.eu/",
        "aliases": [
            "欧洲央行",
            "ecb",
            "european central bank",
            "euro area inflation",
        ],
        "query_terms": [
            "ECB inflation outlook",
            "ECB euro area inflation",
        ],
    },
    {
        "label": "德国联邦统计局",
        "url": "https://www.destatis.de/EN/Home/_node.html",
        "aliases": [
            "德国gdp",
            "德国国内生产总值",
            "destatis",
            "federal statistical office of germany",
            "germany gdp",
        ],
        "query_terms": [
            "Destatis Germany GDP",
        ],
    },
    {
        "label": "法国国家统计与经济研究所",
        "url": "https://www.insee.fr/en/accueil",
        "aliases": [
            "法国gdp",
            "法国国内生产总值",
            "insee",
            "france gdp",
        ],
        "query_terms": [
            "INSEE France GDP",
        ],
    },
    {
        "label": "意大利国家统计局",
        "url": "https://www.istat.it/en/",
        "aliases": [
            "意大利gdp",
            "意大利国内生产总值",
            "istat",
            "italy gdp",
        ],
        "query_terms": [
            "ISTAT Italy GDP",
        ],
    },
    {
        "label": "日本统计局",
        "url": "https://www.stat.go.jp/english/",
        "aliases": [
            "日本cpi",
            "日本通胀",
            "日本统计局",
            "statistics bureau of japan",
            "stat.go.jp",
            "japan cpi",
        ],
        "query_terms": [
            "Japan CPI Statistics Bureau",
        ],
    },
    {
        "label": "日本经济产业省",
        "url": "https://www.meti.go.jp/english/",
        "aliases": [
            "日本工业生产",
            "meti",
            "japan industrial production",
            "日本经济产业省",
            "ministry of economy trade and industry japan",
        ],
        "query_terms": [
            "METI Japan industrial production",
        ],
    },
    {
        "label": "日本银行",
        "url": "https://www.boj.or.jp/en/",
        "aliases": [
            "日本银行",
            "日本央行",
            "boj",
            "bank of japan",
        ],
        "query_terms": [
            "Bank of Japan outlook",
        ],
    },
    {
        "label": "新西兰储备银行",
        "url": "https://www.rbnz.govt.nz/",
        "aliases": [
            "新西兰央行",
            "新西兰储备银行",
            "rbnz",
            "reserve bank of new zealand",
            "official cash rate",
        ],
        "query_terms": [
            "RBNZ official cash rate",
            "RBNZ monetary policy statement",
        ],
    },
    {
        "label": "澳大利亚统计局",
        "url": "https://www.abs.gov.au/",
        "aliases": [
            "澳大利亚通胀",
            "澳洲通胀",
            "australia inflation",
            "australian inflation",
            "abs",
            "australian bureau of statistics",
        ],
        "query_terms": [
            "ABS Australia CPI",
            "Australian Bureau of Statistics inflation",
        ],
    },
    {
        "label": "澳大利亚储备银行",
        "url": "https://www.rba.gov.au/",
        "aliases": [
            "澳大利亚央行",
            "澳洲央行",
            "rba",
            "reserve bank of australia",
        ],
        "query_terms": [
            "RBA rate decision",
            "RBA inflation outlook",
        ],
    },
    {
        "label": "中国人民银行",
        "url": "https://www.pbc.gov.cn/en/3688006/index.html",
        "aliases": [
            "中国人民银行",
            "中国央行",
            "人民银行",
            "央行",
            "pboc",
            "people's bank of china",
            "mlf",
            "中期借贷便利",
            "一年期mlf",
        ],
        "query_terms": [
            "PBOC one-year MLF rate",
            "People's Bank of China MLF operation",
        ],
    },
    {
        "label": "国家外汇管理局",
        "url": "https://www.safe.gov.cn/en/",
        "aliases": [
            "国家外汇管理局",
            "外汇管理局",
            "safe",
            "state administration of foreign exchange",
            "资本外流",
            "资金外流",
        ],
        "query_terms": [
            "SAFE China cross-border capital flows",
            "SAFE China balance of payments",
        ],
    },
    {
        "label": "中国证监会",
        "url": "https://www.csrc.gov.cn/csrc_en/",
        "aliases": [
            "中国证监会",
            "证监会",
            "csrc",
            "china securities regulatory commission",
            "跨境证券",
            "老虎证券",
            "富途",
            "长桥",
            "tiger brokers",
            "futu",
            "longbridge",
        ],
        "query_terms": [
            "CSRC cross-border brokerage",
            "China broker crackdown CSRC",
        ],
    },
    {
        "label": "中国国家统计局",
        "url": "https://www.stats.gov.cn/",
        "aliases": [
            "中国国家统计局",
            "国家统计局",
            "中国工业利润",
            "中国pmi",
            "工业利润",
            "制造业pmi",
            "nbs china",
            "national bureau of statistics of china",
        ],
        "query_terms": [
            "NBS China industrial profits",
            "NBS China PMI",
        ],
    },
    {
        "label": "中国物流与采购联合会",
        "url": "https://www.chinawuliu.com.cn/",
        "aliases": [
            "中国物流与采购联合会",
            "cflp",
            "china federation of logistics and purchasing",
            "中国pmi",
            "制造业pmi",
        ],
        "query_terms": [
            "CFLP PMI China",
        ],
    },
    {
        "label": "香港交易所",
        "url": "https://www.hkex.com.hk/",
        "aliases": [
            "香港ipo",
            "香港上市",
            "港交所",
            "hkex",
            "hong kong ipo",
            "hong kong exchange",
        ],
        "query_terms": [
            "HKEX IPO filing",
        ],
    },
    {
        "label": "美国众议院中国问题特别委员会",
        "url": "https://selectcommitteeontheccp.house.gov/",
        "aliases": [
            "中国问题特别委员会",
            "美国国会",
            "house select committee on the chinese communist party",
            "select committee on the ccp",
            "众议院中国委员会",
        ],
        "query_terms": [
            "House CCP committee CATL Hong Kong IPO",
        ],
    },
    {
        "label": "以色列总理办公室",
        "url": "https://www.gov.il/en/departments/prime_ministers_office",
        "aliases": [
            "以色列总理办公室",
            "内塔尼亚胡",
            "netanyahu",
            "prime minister's office israel",
            "israeli prime minister office",
        ],
        "query_terms": [
            "Netanyahu meeting statement",
        ],
    },
    {
        "label": "乌克兰国防部官网",
        "url": "https://mod.gov.ua/en",
        "aliases": ["乌克兰国防部", "ukrainian ministry of defence", "ministry of defence of ukraine", "mod ukraine"],
    },
    {
        "label": "基辅市政府",
        "url": "https://kyivcity.gov.ua/",
        "aliases": ["基辅市政府", "基辅市政", "kyiv city government", "kyiv city state administration", "kyiv city"],
    },
    {
        "label": "Defense News",
        "url": "https://www.defensenews.com/",
        "aliases": ["defense news", "defensenews"],
    },
    {
        "label": "斯洛伐克政府官网",
        "url": "https://www.vlada.gov.sk/en/",
        "aliases": ["斯洛伐克政府", "slovak government", "government of slovakia", "government office of the slovak republic"],
    },
    {
        "label": "斯洛伐克国防部官网",
        "url": "https://www.mosr.sk/mo-sr-en/",
        "aliases": ["斯洛伐克国防部", "slovak ministry of defense", "ministry of defense of the slovak republic"],
    },
    {
        "label": "塔斯社",
        "url": "https://tass.com/",
        "aliases": ["塔斯社", "tass"],
    },
    {
        "label": "美国国务院",
        "url": "https://www.state.gov/",
        "aliases": ["美国国务院", "u.s. department of state", "us department of state", "state department"],
    },
    {
        "label": "古巴外交部",
        "url": "https://cubaminrex.cu/en",
        "aliases": ["古巴外交部", "cuban ministry of foreign affairs", "ministry of foreign affairs of cuba", "cubaminrex"],
    },
    {
        "label": "匈牙利国家选举办公室",
        "url": "https://www.valasztas.hu/",
        "aliases": ["匈牙利国家选举办公室", "hungarian national election office", "national election office hungary", "valasztas"],
    },
    {
        "label": "UKMTO",
        "url": "https://www.ukmto.org/",
        "aliases": ["ukmto", "united kingdom maritime trade operations", "英国海事贸易行动中心"],
    },
    {
        "label": "Yahoo Finance",
        "url": "https://finance.yahoo.com/",
        "aliases": ["雅虎财经", "yahoo finance"],
        "query_terms": [
            "Yahoo Finance blue-chip stocks",
            "Yahoo Finance Walmart Realty Income Philip Morris",
        ],
    },
    {
        "label": "彭博社",
        "url": "https://www.bloomberg.com/",
        "aliases": ["彭博社", "彭博", "bloomberg"],
        "query_terms": ["Bloomberg markets", "Bloomberg politics"],
    },
    {
        "label": "台湾证券交易所",
        "url": "https://www.twse.com.tw/en/",
        "aliases": ["台湾证券交易所", "台灣證券交易所", "台湾股市", "台灣股市", "twse", "taiwan stock exchange"],
    },
    {
        "label": "美丽岛电子报",
        "url": "https://www.my-formosa.com/",
        "aliases": ["美丽岛电子报", "美麗島電子報", "美丽岛", "美麗島", "my-formosa"],
    },
    {
        "label": "路透社",
        "url": "https://www.reuters.com/",
        "aliases": ["路透社", "reuters"],
        "query_terms": ["Reuters world", "Reuters markets"],
    },
    {
        "label": "美联社",
        "url": "https://apnews.com/",
        "aliases": ["美联社", "associated press", "ap news", "apnews"],
    },
    {
        "label": "纽约时报",
        "url": "https://www.nytimes.com/",
        "aliases": ["纽约时报", "new york times", "nytimes", "nyt"],
        "query_terms": ["New York Times world", "New York Times business"],
    },
    {
        "label": "华尔街日报",
        "url": "https://www.wsj.com/",
        "aliases": ["华尔街日报", "wall street journal", "wsj"],
        "query_terms": ["Wall Street Journal world", "Wall Street Journal markets"],
    },
    {
        "label": "金融时报",
        "url": "https://www.ft.com/",
        "aliases": ["金融时报", "financial times", "ft.com"],
        "query_terms": ["Financial Times world", "Financial Times markets"],
    },
    {
        "label": "华盛顿邮报",
        "url": "https://www.washingtonpost.com/",
        "aliases": ["华盛顿邮报", "washington post", "wapo"],
    },
    {
        "label": "基辅独立报",
        "url": "https://kyivindependent.com/",
        "aliases": ["基辅独立报", "kyiv independent", "the kyiv independent"],
    },
    {
        "label": "乌克兰国家通讯社",
        "url": "https://www.ukrinform.net/",
        "aliases": ["ukrinform", "乌克兰国家通讯社", "ukrainian national news"],
    },
    {
        "label": "日经亚洲",
        "url": "https://asia.nikkei.com/",
        "aliases": ["日经", "nikkei", "nikkei asia"],
    },
    {
        "label": "Axios",
        "url": "https://www.axios.com/",
        "aliases": ["axios"],
    },
    {
        "label": "CNBC",
        "url": "https://www.cnbc.com/",
        "aliases": ["cnbc"],
    },
]


MAJOR_MEDIA_DOMAINS = {
    "reuters.com",
    "www.reuters.com",
    "apnews.com",
    "www.apnews.com",
    "bloomberg.com",
    "www.bloomberg.com",
    "nytimes.com",
    "www.nytimes.com",
    "washingtonpost.com",
    "www.washingtonpost.com",
    "cnn.com",
    "www.cnn.com",
    "bbc.com",
    "www.bbc.com",
    "bbc.co.uk",
    "www.bbc.co.uk",
    "wsj.com",
    "www.wsj.com",
    "ft.com",
    "www.ft.com",
    "theguardian.com",
    "www.theguardian.com",
    "caixin.com",
    "www.caixin.com",
    "axios.com",
    "www.axios.com",
    "cnbc.com",
    "www.cnbc.com",
    "nikkei.com",
    "www.nikkei.com",
    "asia.nikkei.com",
    "kyivindependent.com",
    "www.kyivindependent.com",
    "ukrinform.net",
    "www.ukrinform.net",
    "aljazeera.com",
    "www.aljazeera.com",
    "thehill.com",
    "www.thehill.com",
    "politico.com",
    "www.politico.com",
    "foxbusiness.com",
    "www.foxbusiness.com",
}

NOISY_FACT_CHECK_DOMAIN_TOKENS = (
    "google.",
    "bing.",
    "youtube.com",
    "youtu.be",
    "bilibili.com",
    "douyin.com",
    "tiktok.com",
    "x.com",
    "twitter.com",
    "facebook.com",
    "instagram.com",
    "weibo.com",
)

NOISY_FACT_CHECK_PATH_TOKENS = (
    "/search",
    "/video",
    "/videos",
    "/shorts",
    "/playlist",
    "/account",
    "/login",
    "/signin",
    "/signup",
    "/register",
    "/tag/",
    "/tags/",
    "/topic/",
    "/topics/",
    "/category/",
    "/categories/",
    "/biography/",
    "/encyclopedia/",
    "/dictionary/",
    "/hans/",
)

SYNDICATED_FACT_CHECK_DOMAINS = {
    "msn.com",
    "www.msn.com",
    "yahoo.com",
    "www.yahoo.com",
}

PRIMARY_SOURCE_DOMAIN_HINTS = (
    ("wall street journal", "www.wsj.com"),
    ("wsj", "www.wsj.com"),
    ("financial times", "www.ft.com"),
)

WIRE_SERVICE_DOMAIN_HINTS = (
    (r"\breuters\b", "www.reuters.com"),
    (r"\bassociated press\b|\bap\b", "apnews.com"),
    (r"\bbloomberg\b", "www.bloomberg.com"),
)


_FACT_CHECK_ARTICLE_MATCH_CACHE: dict[tuple[str, str], dict] = {}


def _env_flag(name: str, default: bool = False) -> bool:
    value = str(os.environ.get(name, "") or "").strip().lower()
    if not value:
        return default
    return value in {"1", "true", "yes", "on"}


def _env_int(name: str, default: int, *, minimum: int = 0, maximum: int | None = None) -> int:
    try:
        value = int(str(os.environ.get(name, "") or "").strip() or default)
    except Exception:
        value = default
    value = max(minimum, value)
    if maximum is not None:
        value = min(maximum, value)
    return value


def _build_search_url(query: str, *, engine: str = "google", news: bool = False) -> str:
    """构造通用搜索链接，优先返回可人工复核的检索入口。"""
    encoded_q = quote(query)
    if engine == "bing":
        if news:
            return f"https://www.bing.com/news/search?q={encoded_q}"
        return f"https://www.bing.com/search?q={encoded_q}"
    if news:
        return f"https://www.google.com/search?tbm=nws&q={encoded_q}"
    return f"https://www.google.com/search?q={encoded_q}"


def _normalize_fact_check_source_url(url: str) -> str:
    raw = str(url or "").strip()
    if not raw:
        return ""
    try:
        parsed = urlparse(raw)
        filtered_query_pairs = []
        for key, value in parse_qsl(parsed.query or "", keep_blank_values=True):
            lowered_key = str(key or "").strip().lower()
            if not lowered_key:
                continue
            if (
                lowered_key.startswith("utm_")
                or lowered_key in {
                    "ved", "ei", "usg", "at", "ref", "ref_src", "feature", "fbclid", "gclid",
                    "igshid", "mc_cid", "mc_eid", "src", "source", "spm", "from", "mkt_tok",
                }
            ):
                continue
            filtered_query_pairs.append((key, value))
        normalized_path = re.sub(r"/+$", "", parsed.path or "")
        normalized_query = urlencode(filtered_query_pairs, doseq=True)
        return urlunparse((parsed.scheme.lower(), parsed.netloc.lower(), normalized_path, parsed.params, normalized_query, ""))
    except Exception:
        return raw


def _extract_fact_check_tokens(text: str) -> list[str]:
    raw_tokens = re.findall(
        r"(?:19|20)\d{2}|\d+(?:\.\d+)?(?:%|万亿|亿|万美元|美元|元|桶)?|[A-Za-z][A-Za-z0-9._-]{2,}|[\u4e00-\u9fff]{2,16}",
        str(text or ""),
        re.I,
    )
    stopwords = {
        "相关",
        "有关",
        "消息",
        "报道",
        "报道称",
        "消息称",
        "视频",
        "新闻",
        "网页",
        "官网",
        "official",
        "statement",
        "reported",
        "report",
        "news",
    }
    tokens: list[str] = []
    seen: set[str] = set()
    for token in raw_tokens:
        cleaned = re.sub(r"\s+", " ", str(token or "")).strip().lower()
        if not cleaned or cleaned in seen or cleaned in stopwords:
            continue
        seen.add(cleaned)
        tokens.append(cleaned)
    return tokens


def _collect_fact_check_match_stats(
    haystack: str,
    *,
    claim_text: str = "",
    query_text: str = "",
) -> dict:
    tokens = _extract_fact_check_tokens(" ".join([claim_text, query_text]))
    matched_tokens: list[str] = []
    numeric_matches: list[str] = []
    substantive_matches: list[str] = []
    haystack_lower = str(haystack or "").lower()
    for token in tokens:
        if not token or token not in haystack_lower:
            continue
        matched_tokens.append(token)
        if re.search(r"\d", token):
            if not (token.isdigit() and len(token) < 3):
                numeric_matches.append(token)
            continue
        if len(token) >= 4:
            substantive_matches.append(token)
    return {
        "tokens": tokens,
        "matched_tokens": matched_tokens,
        "matched_token_count": len(matched_tokens),
        "numeric_tokens": [token for token in tokens if re.search(r"\d", token) and not (token.isdigit() and len(token) < 3)],
        "matched_numeric_tokens": numeric_matches,
        "matched_substantive_tokens": substantive_matches,
    }


def _is_noise_fact_check_hit(item: dict) -> bool:
    url = str(item.get("url") or "").strip()
    title = str(item.get("title") or "").strip().lower()
    source_name = str(item.get("source") or "").strip().lower()
    snippet = str(item.get("snippet") or "").strip().lower()
    if not url:
        return True
    parsed = urlparse(url)
    domain = (parsed.netloc or "").lower()
    path = (parsed.path or "").lower()
    normalized_url = _normalize_fact_check_source_url(url)
    if not normalized_url or not domain:
        return True
    if any(token in domain for token in NOISY_FACT_CHECK_DOMAIN_TOKENS):
        return True
    if any(token in path for token in NOISY_FACT_CHECK_PATH_TOKENS):
        return True
    if any(token in normalized_url.lower() for token in ("?s=", "?search=", "/search?", "tbm=nws", "bing.com/news/search")):
        return True
    noisy_title_tokens = [
        "watch video",
        "watch live",
        "playlist",
        "login",
        "sign in",
        "register",
        "subscribe",
        "results for",
        "search result",
        "topic",
        "tag",
        "biography",
        "encyclopedia",
        "dictionary",
        "britannica",
        "汉典",
    ]
    if any(token in title for token in noisy_title_tokens):
        return True
    if any(token in source_name for token in ("youtube", "bilibili", "douyin", "tiktok")):
        return True
    if not title and not snippet:
        return True
    return False


def _is_relevant_fact_check_hit(
    item: dict,
    *,
    claim_text: str = "",
    query_text: str = "",
) -> bool:
    haystack = " ".join(
        [
            str(item.get("title") or ""),
            str(item.get("source") or ""),
            str(item.get("snippet") or ""),
        ]
    ).lower()
    if not haystack:
        return False
    match_stats = _collect_fact_check_match_stats(
        haystack,
        claim_text=claim_text,
        query_text=query_text,
    )
    if len(match_stats.get("matched_substantive_tokens") or []) >= 2:
        return True
    if (
        len(match_stats.get("matched_substantive_tokens") or []) >= 1
        and len(match_stats.get("matched_numeric_tokens") or []) >= 1
    ):
        return True
    if _find_authoritative_sources(haystack):
        return True
    return False


def _score_fact_check_hit(
    item: dict,
    *,
    claim_text: str = "",
    query_text: str = "",
    preferred_domains: set[str] | None = None,
) -> int:
    preferred_domains = {str(domain or "").lower() for domain in (preferred_domains or set()) if str(domain or "").strip()}
    url = str(item.get("url") or "").strip()
    title = str(item.get("title") or "").strip()
    source_name = str(item.get("source") or "").strip()
    snippet = str(item.get("snippet") or "").strip()
    parsed = urlparse(str(item.get("source_url") or "").strip()) if _get_fact_check_hit_domain(item) != (urlparse(url).netloc or "").lower() else urlparse(url)
    domain = (parsed.netloc or "").lower()
    path = (parsed.path or "").lower()
    haystack = " ".join([title, source_name, snippet]).lower()
    score = 0
    preferred_domain = _is_preferred_fact_check_domain(domain, preferred_domains)
    match_stats = _collect_fact_check_match_stats(
        haystack,
        claim_text=claim_text,
        query_text=query_text,
    )

    if preferred_domains and any(domain == preferred or domain.endswith(f".{preferred}") for preferred in preferred_domains):
        score += 16
    if domain in MAJOR_MEDIA_DOMAINS:
        score += 12
    if ".gov" in domain or domain.startswith("gov.") or domain.endswith(".gov"):
        score += 11
    if any(token in domain for token in ("reuters.com", "apnews.com", "bloomberg.com", "nytimes.com", "washingtonpost.com", "bbc.com", "cnn.com", "wsj.com", "ft.com")):
        score += 8
    if path.count("/") >= 2:
        score += 3
    if any(token in path for token in ("/news/", "/world/", "/business/", "/politics/", "/markets/", "/article/", "/articles/", "/story/", "/stories/")):
        score += 5
    if path in {"", "/"}:
        score -= 6

    for token in _extract_fact_check_tokens(" ".join([claim_text, query_text])):
        if token in haystack:
            if re.search(r"\d", token):
                if token.isdigit() and len(token) < 3:
                    continue
                score += 5
            elif len(token) >= 6:
                score += 4
            else:
                score += 2

    if match_stats["numeric_tokens"] and not match_stats["matched_numeric_tokens"]:
        score -= 8
    if not preferred_domain and match_stats["matched_token_count"] <= 1:
        score -= 10
    elif not preferred_domain and len(match_stats["matched_substantive_tokens"]) <= 1:
        score -= 6
    if (
        not preferred_domain
        and re.search(r"[A-Za-z]{3,}", query_text or "")
        and re.search(r"[\u4e00-\u9fff]", title + snippet)
    ):
        score -= 6

    if title:
        title_lower = title.lower()
        for token in _extract_fact_check_tokens(claim_text):
            if token in title_lower:
                if token.isdigit() and len(token) < 3:
                    continue
                score += 2

    return score


def _rerank_fact_check_hits(
    items: list[dict],
    *,
    claim_text: str = "",
    query_text: str = "",
    preferred_domains: set[str] | None = None,
    max_items: int = 3,
) -> list[dict]:
    ranked: list[tuple[int, int, dict]] = []
    for idx, item in enumerate(items or []):
        if _is_noise_fact_check_hit(item):
            continue
        score = _score_fact_check_hit(
            item,
            claim_text=claim_text,
            query_text=query_text,
            preferred_domains=preferred_domains,
        )
        enriched = dict(item)
        enriched["match_score"] = score
        ranked.append((score, idx, enriched))
    ranked.sort(key=lambda entry: (-entry[0], entry[1]))
    return [item for _, _, item in ranked[:max_items]]


def _is_preferred_fact_check_domain(domain: str, preferred_domains: set[str] | None = None) -> bool:
    normalized_domain = str(domain or "").strip().lower()
    preferred = {
        str(item or "").strip().lower()
        for item in (preferred_domains or set())
        if str(item or "").strip()
    }
    if not normalized_domain:
        return False
    return (
        normalized_domain in MAJOR_MEDIA_DOMAINS
        or any(
            normalized_domain == item or normalized_domain.endswith(f".{item}")
            for item in preferred
        )
        or ".gov" in normalized_domain
        or normalized_domain.startswith("gov.")
        or normalized_domain.endswith(".gov")
    )


def _extract_fact_check_article_match_signals(
    url: str,
    *,
    claim_text: str = "",
    query_text: str = "",
    proxy_url: str = None,
) -> dict:
    normalized_url = _normalize_fact_check_source_url(url)
    match_basis = clean_document_text(" ".join([claim_text, query_text]))[:300].lower()
    default_result = {"article_match_score": 0, "article_match_tokens": []}
    if not normalized_url or not match_basis:
        return default_result
    cache_key = (normalized_url, match_basis)
    if cache_key in _FACT_CHECK_ARTICLE_MATCH_CACHE:
        return dict(_FACT_CHECK_ARTICLE_MATCH_CACHE[cache_key])
    try:
        article = extract_web_article_text(url, proxy_url=proxy_url)
        article_text = clean_document_text(
            str(article.get("clean_text") or article.get("raw_text") or "")
        )[:8000].lower()
        if not article_text:
            _FACT_CHECK_ARTICLE_MATCH_CACHE[cache_key] = dict(default_result)
            return dict(default_result)
        matched_tokens: list[str] = []
        score = 0
        for token in _extract_fact_check_tokens(" ".join([claim_text, query_text]))[:12]:
            if token.isdigit() and len(token) < 3:
                continue
            if token not in article_text:
                continue
            matched_tokens.append(token)
            if re.search(r"\d", token):
                score += 6
            elif re.search(r"[\u4e00-\u9fff]", token) and len(token) >= 4:
                score += 4
            elif len(token) >= 6:
                score += 4
            else:
                score += 2
        if len(matched_tokens) >= 3:
            score += 4
        if len(matched_tokens) >= 5:
            score += 4
        result = {
            "article_match_score": score,
            "article_match_tokens": matched_tokens[:5],
        }
    except Exception:
        result = dict(default_result)
    _FACT_CHECK_ARTICLE_MATCH_CACHE[cache_key] = dict(result)
    return dict(result)


def _refine_fact_check_hits_with_article_text(
    items: list[dict],
    *,
    claim_text: str = "",
    query_text: str = "",
    proxy_url: str = None,
    article_fetch_limit: int = 2,
) -> list[dict]:
    if not items or not str(claim_text or query_text).strip():
        return items
    ranked_items = [dict(item) for item in (items or [])]
    fetch_budget = max(0, min(len(ranked_items), int(article_fetch_limit or 0)))
    for idx in range(fetch_budget):
        item = ranked_items[idx]
        signals = _extract_fact_check_article_match_signals(
            str(item.get("url") or ""),
            claim_text=claim_text,
            query_text=query_text,
            proxy_url=proxy_url,
        )
        if not signals:
            continue
        item.update(signals)
        item["match_score"] = int(item.get("match_score") or 0) + int(
            signals.get("article_match_score") or 0
        )
    ranked_items.sort(
        key=lambda item: (-int(item.get("match_score") or 0), -int(item.get("article_match_score") or 0))
    )
    return ranked_items


def _prune_fact_check_hits(
    items: list[dict],
    *,
    preferred_domains: set[str] | None = None,
    max_items: int = 3,
) -> list[dict]:
    if not items:
        return []
    ranked_items = [dict(item) for item in items]
    top_score = int(ranked_items[0].get("match_score") or 0)
    kept: list[dict] = []
    for idx, item in enumerate(ranked_items):
        if len(kept) >= max_items:
            break
        if idx == 0:
            kept.append(item)
            continue
        score = int(item.get("match_score") or 0)
        domain = _get_fact_check_hit_domain(item)
        if score >= max(10, top_score - 10):
            kept.append(item)
            continue
        if _is_preferred_fact_check_domain(domain, preferred_domains) and score >= max(8, top_score - 14):
            kept.append(item)
    return kept


def _dedupe_fact_check_source_links(source_links: list[tuple[str, str]]) -> list[tuple[str, str]]:
    deduped: list[tuple[str, str]] = []
    seen: set[str] = set()
    for label, url in source_links:
        raw_url = str(url or "").strip()
        normalized = _normalize_fact_check_source_url(raw_url)
        if not raw_url or not normalized or normalized in seen:
            continue
        seen.add(normalized)
        deduped.append((str(label or "").strip() or raw_url, raw_url))
    return deduped


def _strip_html_tags(text: str) -> str:
    value = re.sub(r"<[^>]+>", " ", str(text or ""))
    value = re.sub(r"\s+", " ", value).strip()
    return value


def _unwrap_bing_news_redirect(url: str) -> str:
    raw = str(url or "").strip()
    if not raw:
        return ""
    try:
        parsed = urlparse(raw)
        if "bing.com" not in (parsed.netloc or "").lower():
            return raw
        target = parse_qs(parsed.query or "").get("url", [""])[0].strip()
        return target or raw
    except Exception:
        return raw


GOOGLE_NEWS_AGGREGATOR_DOMAINS = {
    "news.google.com",
    "www.news.google.com",
}


def _get_fact_check_hit_domain(item: dict) -> str:
    url = str(item.get("url") or "").strip()
    source_url = str(item.get("source_url") or "").strip()
    url_domain = (urlparse(url).netloc or "").lower()
    source_domain = (urlparse(source_url).netloc or "").lower()
    if url_domain in GOOGLE_NEWS_AGGREGATOR_DOMAINS and source_domain:
        return source_domain
    return url_domain or source_domain


def _extract_fact_check_site_query_domains(query: str) -> list[str]:
    domains: list[str] = []
    seen: set[str] = set()
    for match in re.findall(r"(?<!\S)site:([^\s/]+)", str(query or ""), flags=re.I):
        domain = str(match or "").strip().strip("()\"'").lower()
        if not domain or domain in seen:
            continue
        seen.add(domain)
        domains.append(domain)
    return domains


def _fact_check_hit_matches_domain(item: dict, domain: str) -> bool:
    normalized_domain = str(domain or "").strip().lower()
    if not normalized_domain:
        return False
    candidates = {
        (urlparse(str(item.get("url") or "").strip()).netloc or "").lower(),
        (urlparse(str(item.get("source_url") or "").strip()).netloc or "").lower(),
        _get_fact_check_hit_domain(item),
    }
    return any(
        candidate and (candidate == normalized_domain or candidate.endswith(f".{normalized_domain}"))
        for candidate in candidates
    )


def _extract_google_news_article_id(url: str) -> str:
    raw = str(url or "").strip()
    if not raw:
        return ""
    try:
        parsed = urlparse(raw)
        if "news.google.com" not in (parsed.netloc or "").lower():
            return ""
        path_parts = [part for part in (parsed.path or "").split("/") if part]
        for idx, part in enumerate(path_parts[:-1]):
            if part in {"articles", "read"}:
                return path_parts[idx + 1]
    except Exception:
        return ""
    return ""


def _get_response_set_cookie_headers(response) -> list[str]:
    raw_headers = getattr(getattr(response, "raw", None), "headers", None)
    if raw_headers is not None and hasattr(raw_headers, "get_all"):
        try:
            values = raw_headers.get_all("Set-Cookie") or []
            if values:
                return [str(value or "").strip() for value in values if str(value or "").strip()]
        except Exception:
            pass
    header_value = str(getattr(response, "headers", {}).get("Set-Cookie") or "").strip()
    return [header_value] if header_value else []


def _merge_cookie_header(cookie_header: str, set_cookie_headers: list[str]) -> str:
    cookies: dict[str, str] = {}
    for chunk in str(cookie_header or "").split(";"):
        name, sep, value = chunk.strip().partition("=")
        if sep and name:
            cookies[name] = value
    for header in set_cookie_headers or []:
        name, sep, value = str(header or "").split(";", 1)[0].strip().partition("=")
        if sep and name:
            cookies[name] = value
    return "; ".join(f"{name}={value}" for name, value in cookies.items())


def _extract_google_news_decode_params(article_id: str, proxy_url: str = None) -> tuple[str, str, str]:
    if not article_id:
        return "", "", ""
    cookie_header = "CONSENT=PENDING+987"
    candidate_urls = [
        f"https://news.google.com/articles/{article_id}",
        f"https://news.google.com/rss/articles/{article_id}?oc=5",
    ]
    page_headers = {
        "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
        "sec-ch-ua": '"Google Chrome";v="131", "Chromium";v="131", "Not_A Brand";v="24"',
        "sec-ch-ua-mobile": "?0",
        "sec-ch-ua-platform": '"macOS"',
        "sec-fetch-dest": "document",
        "sec-fetch-mode": "navigate",
        "sec-fetch-site": "none",
        "sec-fetch-user": "?1",
        "upgrade-insecure-requests": "1",
    }
    for candidate_url in candidate_urls:
        try:
            response = requests.get(
                candidate_url,
                headers={**page_headers, "Cookie": cookie_header},
                allow_redirects=False,
                **_build_requests_kwargs(proxy_url, timeout_seconds=8.0),
            )
        except Exception:
            continue
        cookie_header = _merge_cookie_header(cookie_header, _get_response_set_cookie_headers(response))
        response_text = str(response.text or "")
        if "consent.google.com" in str(response.headers.get("Location") or "").lower():
            continue
        signature_match = re.search(r'data-n-a-sg="([^"]+)"', response_text)
        timestamp_match = re.search(r'data-n-a-ts="([^"]+)"', response_text)
        if signature_match and timestamp_match:
            return (
                signature_match.group(1).strip(),
                timestamp_match.group(1).strip(),
                cookie_header,
            )
    return "", "", cookie_header


def _parse_google_batchexecute_response(response_text: str) -> str:
    payload = str(response_text or "")
    if not payload:
        return ""
    try:
        if payload.startswith(")]}'"):
            payload = payload.split("\n\n", 1)[1]
        parsed = json.loads(payload)
        inner_payload = parsed[0][2]
        if isinstance(inner_payload, str):
            inner = json.loads(inner_payload)
            decoded_url = str(inner[1] or "").strip()
            if decoded_url.startswith(("http://", "https://")):
                return decoded_url
    except Exception:
        pass
    match = re.search(r'\[\\"garturlres\\",\\"(.*?)\\",', payload)
    if not match:
        match = re.search(r'\["garturlres","(.*?)",', payload)
    if not match:
        return ""
    try:
        decoded_url = bytes(match.group(1), "utf-8").decode("unicode_escape")
    except Exception:
        decoded_url = match.group(1)
    return decoded_url if decoded_url.startswith(("http://", "https://")) else ""


def _decode_google_news_url_via_signature(article_id: str, proxy_url: str = None) -> str:
    signature, timestamp, cookie_header = _extract_google_news_decode_params(article_id, proxy_url=proxy_url)
    if not signature or not timestamp:
        return ""
    inner_payload = json.dumps(
        [
            "garturlreq",
            [
                ["X", "X", ["X", "X"], None, None, 1, 1, "US:en", None, 1, None, None, None, None, None, 0, 1],
                "X",
                "X",
                1,
                [1, 1, 1],
                1,
                1,
                None,
                0,
                0,
                None,
                0,
            ],
            article_id,
            int(timestamp),
            signature,
        ],
        separators=(",", ":"),
    )
    payload = json.dumps([[["Fbv4je", inner_payload]]], separators=(",", ":"))
    try:
        response = requests.post(
            "https://news.google.com/_/DotsSplashUi/data/batchexecute",
            params={"rpcids": "Fbv4je"},
            headers={
                "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36",
                "Accept": "*/*",
                "Accept-Language": "en-US,en;q=0.9",
                "sec-ch-ua": '"Google Chrome";v="131", "Chromium";v="131", "Not_A Brand";v="24"',
                "sec-ch-ua-mobile": "?0",
                "sec-ch-ua-platform": '"macOS"',
                "sec-fetch-dest": "empty",
                "sec-fetch-mode": "cors",
                "sec-fetch-site": "same-origin",
                "Origin": "https://news.google.com",
                "Referer": "https://news.google.com/",
                "x-same-domain": "1",
                "Content-Type": "application/x-www-form-urlencoded;charset=UTF-8",
                "Cookie": cookie_header,
            },
            data={"f.req": payload},
            **_build_requests_kwargs(proxy_url, timeout_seconds=10.0),
        )
        response.raise_for_status()
    except Exception:
        return ""
    return _parse_google_batchexecute_response(str(response.text or ""))


def _decode_google_news_url(url: str, proxy_url: str = None) -> str:
    raw = str(url or "").strip()
    if not raw:
        return ""
    try:
        parsed = urlparse(raw)
        if "news.google.com" not in (parsed.netloc or "").lower():
            return raw
        article_id = _extract_google_news_article_id(raw)
        if not article_id:
            return raw

        padded = article_id + ("=" * (-len(article_id) % 4))
        decoded_bytes = base64.urlsafe_b64decode(padded)
        if decoded_bytes.startswith(b"\x08\x13\x22"):
            decoded_bytes = decoded_bytes[3:]
        if decoded_bytes.endswith(b"\xd2\x01\x00"):
            decoded_bytes = decoded_bytes[:-3]
        if not decoded_bytes:
            return raw

        first_len = decoded_bytes[0]
        if first_len >= 0x80 and len(decoded_bytes) >= 2:
            candidate = decoded_bytes[2:first_len + 2].decode("latin1", errors="ignore")
        else:
            candidate = decoded_bytes[1:first_len + 1].decode("latin1", errors="ignore")
        if candidate.startswith(("http://", "https://")):
            return candidate
        if not candidate.startswith("AU_yqL"):
            return raw

        decoded_url = _decode_google_news_url_via_signature(article_id, proxy_url=proxy_url)
        if decoded_url:
            return decoded_url

        payload = (
            '[[["Fbv4je","[\\"garturlreq\\",[[\\"en-US\\",\\"US\\",'
            '[\\"FINANCE_TOP_INDICES\\",\\"WEB_TEST_1_0_0\\"],null,null,1,1,\\"US:en\\",null,180,'
            'null,null,null,null,null,0,null,null,[1608992183,723341000]],\\"en-US\\",\\"US\\",1,'
            '[2,3,4,8],1,0,\\"655000234\\",0,0,null,0],\\"'
            + article_id
            + '\\"]",null,"generic"]]]'
        )
        response = requests.post(
            "https://news.google.com/_/DotsSplashUi/data/batchexecute?rpcids=Fbv4je",
            headers={
                **_build_article_headers("https://news.google.com/"),
                "Content-Type": "application/x-www-form-urlencoded;charset=utf-8",
                "Referer": "https://news.google.com/",
            },
            data=f"f.req={quote(payload, safe='')}",
            **_build_requests_kwargs(proxy_url, timeout_seconds=10.0),
        )
        response.raise_for_status()
        decoded_url = _parse_google_batchexecute_response(str(response.text or ""))
        return decoded_url or raw
    except Exception:
        return raw


def _fetch_bing_news_results(query: str, proxy_url: str = None, max_items: int = 4) -> list[dict]:
    query_text = re.sub(r"\s+", " ", str(query or "")).strip()
    if not query_text or max_items <= 0:
        return []

    feed_url = f"https://www.bing.com/news/search?format=rss&q={quote(query_text)}"
    required_domains = _extract_fact_check_site_query_domains(query_text)
    try:
        import xml.etree.ElementTree as ET

        response = requests.get(
            feed_url,
            headers=_build_article_headers(feed_url),
            **_build_requests_kwargs(proxy_url, timeout_seconds=6.0),
        )
        response.raise_for_status()
        root = ET.fromstring(response.text)
    except Exception:
        return []

    items: list[dict] = []
    seen_urls: set[str] = set()
    for item in root.findall(".//item"):
        title = _strip_html_tags(item.findtext("title", default=""))
        link = _unwrap_bing_news_redirect(str(item.findtext("link", default="") or "").strip())
        description = _strip_html_tags(item.findtext("description", default=""))
        pub_date = _strip_html_tags(item.findtext("pubDate", default=""))
        source_name = ""

        source_node = item.find("source")
        if source_node is not None:
            source_name = _strip_html_tags(source_node.text or "")

        if not source_name:
            for child in list(item):
                tag_name = str(child.tag or "")
                if tag_name.lower().endswith("source"):
                    source_name = _strip_html_tags(child.text or "")
                    if source_name:
                        break

        normalized_link = _normalize_fact_check_source_url(link)
        if not title or not link or not normalized_link or normalized_link in seen_urls:
            continue
        candidate = {
            "title": title,
            "url": link,
            "source": source_name,
            "snippet": description,
            "published_at": pub_date,
        }
        if required_domains and not any(
            _fact_check_hit_matches_domain(candidate, domain) for domain in required_domains
        ):
            continue
        seen_urls.add(normalized_link)
        items.append(candidate)
        if len(items) >= max_items:
            break
    return items


def _fetch_google_news_results(query: str, proxy_url: str = None, max_items: int = 4) -> list[dict]:
    query_text = re.sub(r"\s+", " ", str(query or "")).strip()
    if not query_text or max_items <= 0:
        return []

    feed_url = f"https://news.google.com/rss/search?q={quote(query_text)}&hl=en-US&gl=US&ceid=US:en"
    try:
        import xml.etree.ElementTree as ET

        response = requests.get(
            feed_url,
            headers=_build_article_headers(feed_url),
            **_build_requests_kwargs(proxy_url, timeout_seconds=6.0),
        )
        response.raise_for_status()
        root = ET.fromstring(response.text)
    except Exception:
        return []

    items: list[dict] = []
    seen_urls: set[str] = set()
    for item in root.findall(".//item"):
        title = _strip_html_tags(item.findtext("title", default=""))
        link = _decode_google_news_url(str(item.findtext("link", default="") or "").strip(), proxy_url=proxy_url)
        description = _strip_html_tags(item.findtext("description", default=""))
        pub_date = _strip_html_tags(item.findtext("pubDate", default=""))
        source_name = ""

        source_node = item.find("source")
        if source_node is not None:
            source_name = _strip_html_tags(source_node.text or "")
            source_url = str(source_node.get("url") or "").strip()
        else:
            source_url = ""

        normalized_link = _normalize_fact_check_source_url(link)
        if not title or not link or not normalized_link or normalized_link in seen_urls:
            continue
        seen_urls.add(normalized_link)
        items.append(
            {
                "title": title,
                "url": link,
                "source": source_name,
                "source_url": source_url,
                "snippet": description,
                "published_at": pub_date,
            }
        )
        if len(items) >= max_items:
            break
    return items


def _fetch_bing_web_results(query: str, proxy_url: str = None, max_items: int = 4) -> list[dict]:
    query_text = re.sub(r"\s+", " ", str(query or "")).strip()
    if not query_text or max_items <= 0:
        return []

    search_url = f"https://www.bing.com/search?q={quote(query_text)}"
    required_domains = _extract_fact_check_site_query_domains(query_text)
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return []

    try:
        response = requests.get(
            search_url,
            headers=_build_article_headers(search_url),
            **_build_requests_kwargs(proxy_url, timeout_seconds=8.0),
        )
        response.raise_for_status()
    except Exception:
        return []

    html = str(response.text or "")
    if not html:
        return []

    soup = BeautifulSoup(html, "html.parser")
    items: list[dict] = []
    seen_urls: set[str] = set()
    for node in soup.select("li.b_algo"):
        anchor = node.select_one("h2 a") or node.select_one("a")
        if anchor is None:
            continue
        title = _strip_html_tags(anchor.get_text(" ", strip=True))
        url = str(anchor.get("href") or "").strip()
        normalized_url = _normalize_fact_check_source_url(url)
        if not title or not url or not normalized_url or normalized_url in seen_urls:
            continue

        snippet_node = node.select_one(".b_caption p") or node.select_one("p")
        snippet = _strip_html_tags(snippet_node.get_text(" ", strip=True) if snippet_node else "")
        cite_node = node.select_one("cite")
        source_name = _strip_html_tags(cite_node.get_text(" ", strip=True) if cite_node else "")
        candidate = {
            "title": title,
            "url": url,
            "source": source_name,
            "snippet": snippet,
            "published_at": "",
        }
        if required_domains and not any(
            _fact_check_hit_matches_domain(candidate, domain) for domain in required_domains
        ):
            continue

        seen_urls.add(normalized_url)
        items.append(candidate)
        if len(items) >= max_items:
            break
    return items


ANYSEARCH_AUTHORITY_PATTERNS = re.compile(
    r"\b(reuters|bloomberg|associated press|ap news|financial times|wall street journal|wsj|"
    r"cnbc|nikkei|caixin|politico|axios|the hill|al jazeera|bbc|cnn|guardian)\b",
    re.I,
)

ANYSEARCH_IMPORTANT_TOPIC_PATTERNS = re.compile(
    r"(财经|金融|股市|股票|市值|汇率|央行|利率|通胀|原油|油价|政策|政府|选举|战争|冲突|外交|制裁|"
    r"finance|financial|stock|market cap|inflation|central bank|interest rate|oil|policy|government|"
    r"election|war|conflict|diplomacy|sanction)",
    re.I,
)


def _fact_check_hit_domain(item: dict) -> str:
    try:
        return (urlparse(str(item.get("url") or "").strip()).netloc or "").lower()
    except Exception:
        return ""


def _has_syndicated_fact_check_hits(items: list[dict]) -> bool:
    for item in items or []:
        domain = _fact_check_hit_domain(item)
        source_name = str(item.get("source") or "").strip().lower()
        if domain in SYNDICATED_FACT_CHECK_DOMAINS:
            return True
        if any(token in source_name for token in ("reuters on", "bloomberg on", "associated press on", " ap on ")):
            return True
    return False


def _should_use_anysearch_for_fact_check(
    claim_text: str,
    query_text: str,
    *,
    news_hits: list[dict],
    web_hits: list[dict],
    preferred_domains: set[str],
) -> bool:
    combined = " ".join([str(claim_text or ""), str(query_text or "")])
    if preferred_domains:
        return True
    if ANYSEARCH_AUTHORITY_PATTERNS.search(combined):
        return True
    if _has_syndicated_fact_check_hits(news_hits) or _has_syndicated_fact_check_hits(web_hits):
        return True
    if not news_hits and not web_hits and ANYSEARCH_IMPORTANT_TOPIC_PATTERNS.search(combined):
        return True
    if len(news_hits or []) + len(web_hits or []) <= 1 and ANYSEARCH_IMPORTANT_TOPIC_PATTERNS.search(combined):
        return True
    return False


def _fetch_anysearch_results(query: str, proxy_url: str = None, max_items: int = 3) -> list[dict]:
    query_text = re.sub(r"\s+", " ", str(query or "")).strip()
    if not query_text or max_items <= 0:
        return []

    endpoint = str(os.environ.get("ANYSEARCH_API_URL") or "https://api.anysearch.com/v1/search").strip()
    headers = {"Content-Type": "application/json"}
    api_key = str(os.environ.get("ANYSEARCH_API_KEY") or "").strip()
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"
    payload = {
        "query": query_text,
        "max_results": max(1, min(5, max_items)),
        "content_types": ["news", "web"],
        "zone": "intl",
    }
    try:
        response = requests.post(
            endpoint,
            headers=headers,
            json=payload,
            **_build_requests_kwargs(proxy_url, timeout_seconds=8.0),
        )
        response.raise_for_status()
        data = response.json()
    except Exception:
        return []

    raw_results = []
    if isinstance(data, dict):
        result_data = data.get("data")
        if isinstance(result_data, dict):
            raw_results = result_data.get("results") or []
        elif isinstance(data.get("results"), list):
            raw_results = data.get("results") or []

    items: list[dict] = []
    seen_urls: set[str] = set()
    for raw in raw_results:
        if not isinstance(raw, dict):
            continue
        title = _strip_html_tags(str(raw.get("title") or "").strip())
        url = str(raw.get("url") or raw.get("link") or "").strip()
        normalized_url = _normalize_fact_check_source_url(url)
        if not title or not normalized_url or normalized_url in seen_urls:
            continue
        snippet = _strip_html_tags(str(raw.get("snippet") or raw.get("description") or "").strip())
        content = _strip_html_tags(str(raw.get("content") or "").strip())
        if content and len(snippet) < 120:
            snippet = content[:320]
        source = str(raw.get("source") or raw.get("site_name") or "").strip()
        if not source:
            source = urlparse(url).netloc
        seen_urls.add(normalized_url)
        items.append(
            {
                "title": title,
                "url": url,
                "source": source,
                "snippet": snippet,
                "published_at": str(raw.get("published_at") or raw.get("date") or "").strip(),
                "provider": "AnySearch",
            }
        )
        if len(items) >= max_items:
            break
    return items


def _guess_fact_check_primary_source_domain(source_name: str) -> str:
    source_text = re.sub(r"\s+", " ", str(source_name or "")).strip().lower()
    if not source_text:
        return ""
    matched_sources = _find_authoritative_sources(source_text)
    if matched_sources:
        return (urlparse(str(matched_sources[0].get("url") or "")).netloc or "").lower()
    for needle, domain in PRIMARY_SOURCE_DOMAIN_HINTS:
        if needle in source_text:
            return domain
    return ""


def _guess_fact_check_item_primary_source_domain(item: dict) -> str:
    source_url = str(item.get("source_url") or "").strip()
    source_domain = (urlparse(source_url).netloc or "").lower()
    if source_domain:
        return source_domain
    combined = " ".join(
        [
            str(item.get("source") or ""),
            str(item.get("title") or ""),
            str(item.get("snippet") or ""),
        ]
    ).strip()
    matched_sources = _find_authoritative_sources(combined)
    if matched_sources:
        return (urlparse(str(matched_sources[0].get("url") or "")).netloc or "").lower()
    source_domain = _guess_fact_check_primary_source_domain(str(item.get("source") or ""))
    if source_domain:
        return source_domain
    lowered = combined.lower()
    for pattern, domain in WIRE_SERVICE_DOMAIN_HINTS:
        if re.search(pattern, lowered, re.I):
            return domain
    return ""


def _slugify_fact_check_title(text: str) -> str:
    value = re.sub(r"\s+", " ", str(text or "")).strip()
    value = re.sub(
        r"\s*[-|]\s*(reuters|associated press|ap|bloomberg)\s*$",
        "",
        value,
        flags=re.I,
    )
    value = re.sub(
        r"\b(?:[A-Za-z]\.){2,}",
        lambda match: match.group(0).replace(".", ""),
        value,
    )
    value = value.lower().replace("&", " and ")
    value = re.sub(r"[’']", "", value)
    value = re.sub(r"[^a-z0-9]+", "-", value)
    value = re.sub(r"-{2,}", "-", value).strip("-")
    return value


def _build_reuters_title_variants(
    title: str,
    *,
    claim_text: str = "",
    query_text: str = "",
    snippet: str = "",
) -> list[str]:
    base_title = re.sub(r"\s+", " ", str(title or "")).strip(" -|")
    if not base_title:
        return []
    combined_text = clean_document_text(" ".join([claim_text, query_text, snippet])).lower()
    variants: list[str] = []
    seen: set[str] = set()

    def add(value: str) -> None:
        candidate = re.sub(r"\s+", " ", str(value or "")).strip(" -|,;:")
        lowered = candidate.lower()
        if not candidate or lowered in seen:
            return
        if len(_slugify_fact_check_title(candidate)) < 16:
            return
        seen.add(lowered)
        variants.append(candidate)

    add(base_title)
    without_parenthetical = re.sub(r"\s+\([^)]*\)\s*$", "", base_title).strip(" -|")
    add(without_parenthetical)
    for separator in (":", ";", ","):
        prefix = base_title.split(separator, 1)[0].strip(" -|")
        if len(_extract_fact_check_tokens(prefix)) >= 4:
            add(prefix)

    if (
        "no progress" in combined_text
        and "if no progress" not in base_title.lower()
        and re.search(r"\bstop mediating\b", base_title, re.I)
    ):
        add(f"{base_title} if no progress")
    if (
        "as early as next week" in combined_text
        and "as early as next week" not in base_title.lower()
        and re.search(r"\bvisit\b", base_title, re.I)
    ):
        add(f"{base_title} as early as next week")
    if "yonhap reports" in combined_text and "yonhap reports" not in base_title.lower():
        add(f"{base_title}, Yonhap reports")
    return variants[:6]


def _get_fact_check_candidate_dates(published_at: str) -> list[str]:
    published_text = re.sub(r"\s+", " ", str(published_at or "")).strip()
    if not published_text:
        return []
    try:
        parsed = parsedate_to_datetime(published_text)
    except Exception:
        return []
    if parsed.tzinfo is not None:
        parsed = parsed.astimezone().replace(tzinfo=None)
    dates: list[str] = []
    seen: set[str] = set()
    for delta_days in (0, -1, 1):
        candidate = (parsed + timedelta(days=delta_days)).strftime("%Y-%m-%d")
        if candidate in seen:
            continue
        seen.add(candidate)
        dates.append(candidate)
    return dates


def _guess_reuters_section_candidates(
    title: str,
    *,
    claim_text: str = "",
    query_text: str = "",
    snippet: str = "",
) -> list[str]:
    haystack = " ".join([title, claim_text, query_text, snippet]).lower()
    sections: list[str] = []
    seen: set[str] = set()

    def add(section: str) -> None:
        value = str(section or "").strip().strip("/")
        if not value or value in seen:
            return
        seen.add(value)
        sections.append(value)

    if re.search(r"\b(korea|north korea|south korea|china|xi|pyongyang|beijing|yonhap|asia)\b", haystack):
        add("world/asia-pacific")
        add("world/china")
    if re.search(r"\b(ukraine|russia|kyiv|moscow|europe|rubio|nato)\b", haystack):
        add("world/europe")
    if re.search(r"\b(iran|israel|gaza|hormuz|tehran|middle east|netanyahu)\b", haystack):
        add("world/middle-east")
        add("world")
        add("business/energy")
    if re.search(r"\b(catl|ipo|listing|hkex|hong kong|jpmorgan|bank of america|market|stocks)\b", haystack):
        add("business")
        add("markets")
        add("world/china")
    add("world")
    add("business")
    add("markets")
    return sections


def _is_probable_reuters_article_url(url: str) -> bool:
    raw = str(url or "").strip()
    if not raw:
        return False
    try:
        parsed = urlparse(raw)
    except Exception:
        return False
    domain = (parsed.netloc or "").lower()
    path = (parsed.path or "").lower()
    if domain and "reuters.com" not in domain:
        return False
    if not path or path in {"", "/"}:
        return False
    if any(
        token in path
        for token in (
            "/live-updates/",
            "/graphics/",
            "/video/",
            "/videos/",
            "/podcast/",
            "/podcasts/",
            "/pictures/",
            "/fact-check/",
            "/breakingviews/",
            "/plus/",
        )
    ):
        return False
    last_segment = [segment for segment in path.split("/") if segment]
    slug = last_segment[-1] if last_segment else ""
    if any(
        slug.startswith(prefix)
        for prefix in (
            "analysis-",
            "commentary-",
            "explainer-",
            "factbox-",
            "live-",
            "liveblog-",
            "picture-",
            "podcast-",
            "timeline-",
            "video-",
        )
    ):
        return False
    return True


def _fetch_fact_check_page_metadata(url: str, proxy_url: str = None) -> dict:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return {}

    try:
        response = requests.get(
            url,
            headers=_build_article_headers(url),
            **_build_requests_kwargs(proxy_url, timeout_seconds=20.0),
        )
        response.raise_for_status()
    except Exception:
        return {}

    html = str(response.text or "")
    if not html:
        return {}
    try:
        soup = BeautifulSoup(html, "html.parser")
    except Exception:
        return {}

    title = ""
    for selector in (
        'meta[property="og:title"]',
        'meta[name="twitter:title"]',
        'meta[name="title"]',
    ):
        node = soup.select_one(selector)
        content = str(node.get("content") or "").strip() if node else ""
        if content:
            title = content
            break
    if not title and soup.title and soup.title.string:
        title = soup.title.string.strip()

    canonical_url = ""
    canonical_node = soup.select_one('link[rel="canonical"]')
    if canonical_node:
        canonical_url = str(canonical_node.get("href") or "").strip()
    if not canonical_url:
        og_url = soup.select_one('meta[property="og:url"]')
        canonical_url = str(og_url.get("content") or "").strip() if og_url else ""

    published_meta = ""
    for selector in (
        'meta[property="article:published_time"]',
        'meta[name="article:published_time"]',
        'meta[name="pubdate"]',
        'meta[name="date"]',
    ):
        node = soup.select_one(selector)
        content = str(node.get("content") or "").strip() if node else ""
        if content:
            published_meta = content
            break

    return {
        "title": title,
        "canonical_url": canonical_url,
        "published_at": published_meta,
    }


def _recover_reuters_direct_article_hit(
    *,
    title: str,
    published_at: str = "",
    claim_text: str = "",
    query_text: str = "",
    snippet: str = "",
    proxy_url: str = None,
) -> dict | None:
    clean_title = re.sub(r"\s+", " ", str(title or "")).strip(" -|")
    title_variants = _build_reuters_title_variants(
        clean_title,
        claim_text=claim_text,
        query_text=query_text,
        snippet=snippet,
    )
    if not title_variants:
        return None

    dates = _get_fact_check_candidate_dates(published_at)
    sections = _guess_reuters_section_candidates(
        title_variants[0],
        claim_text=claim_text,
        query_text=query_text,
        snippet=snippet,
    )
    candidate_records: list[dict] = []
    seen_urls: set[str] = set()
    candidate_slugs: list[str] = []
    seen_slugs: set[str] = set()
    for variant in title_variants[:4]:
        slug = _slugify_fact_check_title(variant)
        if not slug or slug in seen_slugs:
            continue
        seen_slugs.add(slug)
        candidate_slugs.append(slug)
    for section in sections[:6]:
        for candidate_date in dates[:3]:
            for slug in candidate_slugs:
                url = f"https://www.reuters.com/{section}/{slug}-{candidate_date}/"
                if url in seen_urls:
                    continue
                seen_urls.add(url)
                candidate_records.append(
                    {
                        "url": url,
                        "title_variant": slug,
                    }
                )

    title_tokens: list[str] = []
    seen_title_tokens: set[str] = set()
    for variant in title_variants[:4]:
        for token in _extract_fact_check_tokens(variant):
            if len(token) < 4 and not re.search(r"\d", token):
                continue
            if token in seen_title_tokens:
                continue
            seen_title_tokens.add(token)
            title_tokens.append(token)
    for candidate_record in candidate_records[:10]:
        candidate_url = str(candidate_record.get("url") or "").strip()
        matched_tokens: list[str] = []
        try:
            article = extract_web_article_text(candidate_url, proxy_url=proxy_url)
            article_text = clean_document_text(
                str(article.get("clean_text") or article.get("raw_text") or "")
            ).lower()
        except Exception:
            article_text = ""

        if article_text:
            matched_tokens = [token for token in title_tokens[:12] if token in article_text]
        if len(matched_tokens) >= max(3, min(5, len(title_tokens[:12]))):
            return {
                "title": clean_title,
                "url": candidate_url,
                "source": "Reuters",
                "source_url": "https://www.reuters.com/",
                "snippet": snippet,
                "published_at": published_at,
                "article_match_tokens": matched_tokens[:5],
                "article_match_score": 12 + (2 * len(matched_tokens[:5])),
                "match_score": 38 + (2 * len(matched_tokens[:5])),
            }

        metadata = _fetch_fact_check_page_metadata(candidate_url, proxy_url=proxy_url)
        metadata_text = clean_document_text(
            " ".join(
                [
                    str(metadata.get("title") or ""),
                    str(metadata.get("canonical_url") or ""),
                    str(metadata.get("published_at") or ""),
                ]
            )
        ).lower()
        canonical_url = str(metadata.get("canonical_url") or "").strip() or candidate_url
        canonical_domain = (urlparse(canonical_url).netloc or "").lower()
        if not metadata_text or "reuters" not in metadata_text:
            continue
        if canonical_domain and "reuters.com" not in canonical_domain:
            continue
        if not _is_probable_reuters_article_url(canonical_url):
            continue
        matched_tokens = [token for token in title_tokens[:12] if token in metadata_text]
        if len(matched_tokens) < max(4, min(6, len(title_tokens[:12]))):
            continue
        return {
            "title": clean_title,
            "url": canonical_url,
            "source": "Reuters",
            "source_url": "https://www.reuters.com/",
            "snippet": snippet,
            "published_at": published_at,
            "article_match_tokens": matched_tokens[:5],
            "article_match_score": 8 + (2 * len(matched_tokens[:5])),
            "match_score": 32 + (2 * len(matched_tokens[:5])),
        }

    evidence_text = clean_document_text(" ".join([clean_title, claim_text, query_text, snippet])).lower()
    if re.search(r"\breuters\b", " ".join([clean_title, query_text, snippet]), re.I):
        best_guess: tuple[int, dict, list[str]] | None = None
        for candidate_record in candidate_records[:8]:
            candidate_url = str(candidate_record.get("url") or "").strip()
            if not _is_probable_reuters_article_url(candidate_url):
                continue
            candidate_slug = (urlparse(candidate_url).path or "").strip("/").split("/")[-1]
            candidate_slug = re.sub(r"-\d{4}-\d{2}-\d{2}$", "", candidate_slug)
            variant_tokens = [
                token
                for token in _extract_fact_check_tokens(candidate_slug.replace("-", " "))
                if len(token) >= 4 or re.search(r"\d", token)
            ]
            if len(variant_tokens) < 5:
                continue
            matched_tokens = [token for token in variant_tokens[:10] if token in evidence_text]
            if len(matched_tokens) < max(5, min(7, len(variant_tokens[:10]))):
                continue
            score = len(matched_tokens)
            if not best_guess or score > best_guess[0]:
                best_guess = (score, candidate_record, matched_tokens[:5])
        if best_guess:
            _, candidate_record, matched_tokens = best_guess
            return {
                "title": clean_title,
                "url": str(candidate_record.get("url") or "").strip(),
                "source": "Reuters",
                "source_url": "https://www.reuters.com/",
                "snippet": snippet,
                "published_at": published_at,
                "article_match_tokens": matched_tokens,
                "article_match_score": 0,
                "match_score": 24 + (2 * len(matched_tokens)),
                "recovery_method": "slug_guess",
            }
    return None


def _recover_syndicated_fact_check_hit(
    item: dict,
    *,
    claim_text: str = "",
    query_text: str = "",
    proxy_url: str = None,
) -> dict | None:
    url = str(item.get("url") or "").strip()
    source_name = str(item.get("source") or "").strip()
    title = re.sub(r"\s+", " ", str(item.get("title") or "")).strip(" -|")
    snippet = str(item.get("snippet") or "").strip()
    published_at = str(item.get("published_at") or "").strip()
    domain = (urlparse(url).netloc or "").lower()
    if not title:
        return None

    primary_domain = _guess_fact_check_item_primary_source_domain(item)
    if not primary_domain:
        return None

    actual_url_domain = (urlparse(url).netloc or "").lower()
    if actual_url_domain and (
        actual_url_domain == primary_domain or actual_url_domain.endswith(f".{primary_domain}")
    ):
        return None

    if (
        domain not in SYNDICATED_FACT_CHECK_DOMAINS
        and not re.search(r"\breuters\b|\bassociated press\b|\bbloomberg\b", " ".join([source_name, snippet]), re.I)
    ):
        return None

    query_candidates: list[str] = []
    seen_queries: set[str] = set()

    def add_query(value: str) -> None:
        query = re.sub(r"\s+", " ", str(value or "")).strip()
        lowered = query.lower()
        if not query or lowered in seen_queries:
            return
        seen_queries.add(lowered)
        query_candidates.append(query)

    add_query(f'site:{primary_domain} "{title}"')
    add_query(f"site:{primary_domain} {title}")
    title_tokens = [token for token in _extract_fact_check_tokens(title) if len(token) >= 4 or re.search(r"\d", token)]
    if title_tokens:
        add_query(f"site:{primary_domain} {' '.join(title_tokens[:10])}")

    def recover_from_google_hits(google_hits: list[dict]) -> dict | None:
        exact_queries: list[str] = []
        seen_exact_queries: set[str] = set()

        def add_exact_query(value: str) -> None:
            normalized = re.sub(r"\s+", " ", str(value or "")).strip()
            lowered = normalized.lower()
            if not normalized or lowered in seen_exact_queries:
                return
            seen_exact_queries.add(lowered)
            exact_queries.append(normalized)

        for google_hit in google_hits[:2]:
            google_title = re.sub(r"\s+", " ", str(google_hit.get("title") or "")).strip(" -|")
            google_source_name = str(google_hit.get("source") or "").strip()
            if google_source_name:
                google_title = re.sub(
                    rf"\s*[-|]\s*{re.escape(google_source_name)}\s*$",
                    "",
                    google_title,
                    flags=re.I,
                ).strip(" -|")
            if not google_title:
                continue
            if primary_domain in {"www.reuters.com", "reuters.com"}:
                recovered_direct = _recover_reuters_direct_article_hit(
                    title=google_title,
                    published_at=str(google_hit.get("published_at") or ""),
                    claim_text=claim_text,
                    query_text=f"{query_text} {title}",
                    snippet=str(google_hit.get("snippet") or ""),
                    proxy_url=proxy_url,
                )
                if recovered_direct:
                    return recovered_direct
            add_exact_query(f'site:{primary_domain} "{google_title}"')
            exact_title_tokens = [
                token
                for token in _extract_fact_check_tokens(google_title)
                if len(token) >= 4 or re.search(r"\d", token)
            ]
            if exact_title_tokens:
                add_exact_query(f"site:{primary_domain} {' '.join(exact_title_tokens[:12])}")

        for exact_query in exact_queries[:4]:
            direct_hits = _fetch_bing_web_results(exact_query, proxy_url=proxy_url, max_items=5)
            direct_hits = [hit for hit in direct_hits if _fact_check_hit_matches_domain(hit, primary_domain)]
            if not direct_hits:
                continue
            direct_hits = _rerank_fact_check_hits(
                direct_hits,
                claim_text=claim_text,
                query_text=f"{query_text} {title}",
                preferred_domains={primary_domain},
                max_items=3,
            )
            direct_hits = _refine_fact_check_hits_with_article_text(
                direct_hits,
                claim_text=claim_text,
                query_text=f"{query_text} {title}",
                proxy_url=proxy_url,
                article_fetch_limit=1,
            )
            direct_hits = _prune_fact_check_hits(
                direct_hits,
                preferred_domains={primary_domain},
                max_items=1,
            )
            if direct_hits:
                return dict(direct_hits[0])
        return None

    if primary_domain in {"www.reuters.com", "reuters.com"}:
        recovered_direct = _recover_reuters_direct_article_hit(
            title=title,
            published_at=published_at,
            claim_text=claim_text,
            query_text=query_text,
            snippet=snippet,
            proxy_url=proxy_url,
        )
        if recovered_direct:
            recovered_direct["recovered_from_source"] = source_name
            recovered_direct["recovered_from_url"] = url
            return recovered_direct

    for recovery_query in query_candidates[:2]:
        site_hits = _fetch_bing_web_results(recovery_query, proxy_url=proxy_url, max_items=5)
        site_hits = [hit for hit in site_hits if _fact_check_hit_matches_domain(hit, primary_domain)]
        if not site_hits:
            google_hits = _fetch_google_news_results(recovery_query, proxy_url=proxy_url, max_items=5)
            google_hits = [hit for hit in google_hits if _fact_check_hit_matches_domain(hit, primary_domain)]
            if not google_hits:
                continue
            recovered_from_google = recover_from_google_hits(google_hits)
            if not recovered_from_google:
                continue
            site_hits = [recovered_from_google]
        site_hits = _rerank_fact_check_hits(
            site_hits,
            claim_text=claim_text,
            query_text=f"{query_text} {title}",
            preferred_domains={primary_domain},
            max_items=3,
        )
        site_hits = _refine_fact_check_hits_with_article_text(
            site_hits,
            claim_text=claim_text,
            query_text=f"{query_text} {title}",
            proxy_url=proxy_url,
            article_fetch_limit=1,
        )
        site_hits = _prune_fact_check_hits(
            site_hits,
            preferred_domains={primary_domain},
            max_items=1,
        )
        if not site_hits:
            continue
        recovered = dict(site_hits[0])
        recovered["recovered_from_source"] = source_name
        recovered["recovered_from_url"] = url
        return recovered
    return None


def _recover_syndicated_fact_check_hits(
    items: list[dict],
    *,
    claim_text: str = "",
    query_text: str = "",
    proxy_url: str = None,
    recovery_limit: int = 2,
) -> list[dict]:
    if not items:
        return []
    recovered_items: list[dict] = []
    seen_urls: set[str] = set()
    recovery_budget = max(0, int(recovery_limit or 0))
    for item in items:
        chosen = dict(item)
        if recovery_budget > 0:
            recovered = _recover_syndicated_fact_check_hit(
                item,
                claim_text=claim_text,
                query_text=query_text,
                proxy_url=proxy_url,
            )
            if recovered:
                chosen = recovered
                recovery_budget -= 1
        normalized_url = _normalize_fact_check_source_url(str(chosen.get("url") or ""))
        if not normalized_url or normalized_url in seen_urls:
            continue
        seen_urls.add(normalized_url)
        recovered_items.append(chosen)
    return recovered_items


def _find_authoritative_sources(text: str) -> list[dict]:
    """根据文本中的机构名/媒体名匹配官网与权威媒体站点。"""
    haystack = str(text or "").lower()
    if not haystack:
        return []
    matched: list[dict] = []
    seen_urls: set[str] = set()
    for rule in AUTHORITATIVE_SOURCE_RULES:
        aliases = [str(alias or "").lower() for alias in rule.get("aliases", [])]
        if not any(alias and alias in haystack for alias in aliases):
            continue
        url = str(rule.get("url") or "").strip()
        if not url or url in seen_urls:
            continue
        seen_urls.add(url)
        matched.append(rule)
    return matched


def _extract_year_tokens(text: str) -> list[str]:
    """提取文本中的年份，用于约束强时效搜索。"""
    years = re.findall(r"\b(?:19|20)\d{2}\b", str(text or ""))
    deduped: list[str] = []
    seen: set[str] = set()
    for year in years:
        if year in seen:
            continue
        seen.add(year)
        deduped.append(year)
    return deduped


FACT_CHECK_ENGLISH_QUERY_REPLACEMENTS = [
    (r"川普|特朗普", "Trump"),
    (r"卢比奥", "Rubio"),
    (r"习近平", "Xi Jinping"),
    (r"朝鲜", "North Korea"),
    (r"乌克兰", "Ukraine"),
    (r"基辅", "Kyiv"),
    (r"伊朗", "Iran"),
    (r"霍尔木兹海峡", "Strait of Hormuz"),
    (r"核协议", "nuclear deal"),
    (r"谈判", "talks"),
    (r"美国国务卿", "U.S. Secretary of State"),
    (r"退出.*?调解", "stop mediating"),
    (r"无成果", "no progress"),
    (r"可能访问", "may visit"),
    (r"访问", "visit"),
    (r"宁德时代", "CATL"),
    (r"香港IPO|香港上市", "Hong Kong IPO"),
    (r"摩根大通", "JPMorgan"),
    (r"美国银行", "Bank of America"),
    (r"格雷厄姆", "Graham"),
    (r"内塔尼亚胡", "Netanyahu"),
    (r"重新开放", "reopen"),
    (r"清除水雷", "clear mines"),
    (r"解除封锁", "lift blockade"),
    (r"恢复.*?石油出口", "restore oil exports"),
]


FACT_CHECK_MAJOR_MEDIA_SITE_RULES = [
    (
        re.compile(r"(雅虎财经|yahoo finance|沃尔玛|walmart|realty income|philip morris|blue chip|蓝筹股)", re.I),
        ["finance.yahoo.com", "www.reuters.com", "www.bloomberg.com", "www.wsj.com"],
    ),
    (
        re.compile(r"(朝鲜|north korea|访问|visit|pyongyang)", re.I),
        ["www.reuters.com", "www.bloomberg.com", "www.nytimes.com", "www.wsj.com"],
    ),
    (
        re.compile(r"(伊朗|iran|hormuz|霍尔木兹|核协议|nuclear deal|谈判|sanctions)", re.I),
        ["www.reuters.com", "www.bloomberg.com", "www.wsj.com", "www.ft.com"],
    ),
    (
        re.compile(r"(乌克兰|ukraine|基辅|kyiv|无人机|drone|导弹|missile|rubio)", re.I),
        ["www.reuters.com", "www.nytimes.com", "www.wsj.com", "www.ft.com"],
    ),
    (
        re.compile(r"(ipo|上市|香港ipo|香港上市|宁德时代|catl|jpmorgan|bank of america)", re.I),
        ["www.reuters.com", "www.bloomberg.com", "www.wsj.com", "www.ft.com"],
    ),
    (
        re.compile(r"(mlf|中国央行|人民银行|pboc|people's bank of china|净利差|银行资金成本)", re.I),
        ["www.reuters.com", "www.bloomberg.com", "www.wsj.com", "www.ft.com"],
    ),
    (
        re.compile(r"(跨境证券|老虎|富途|长桥|futu|tiger brokers|longbridge|资本外流|资金外流)", re.I),
        ["www.reuters.com", "www.bloomberg.com", "www.wsj.com", "www.ft.com"],
    ),
    (
        re.compile(r"(deepseek|阿里巴巴|alibaba|护照|passport|出国需.*批准|exit control|ai人才|ai talent|manus)", re.I),
        ["www.bloomberg.com", "www.reuters.com", "www.wsj.com", "www.ft.com"],
    ),
]


def _suggest_fact_check_major_media_domains(text: str) -> list[str]:
    value = re.sub(r"\s+", " ", str(text or "")).strip()
    if not value:
        return []
    domains: list[str] = []
    seen: set[str] = set()
    for pattern, candidates in FACT_CHECK_MAJOR_MEDIA_SITE_RULES:
        if not pattern.search(value):
            continue
        for domain in candidates:
            normalized = str(domain or "").strip().lower()
            if not normalized or normalized in seen:
                continue
            seen.add(normalized)
            domains.append(normalized)
    return domains


FACT_CHECK_CONTEXT_PREFIXES = [
    "主要内容",
    "财经头条",
    "视频观点",
    "中国经济",
    "AI人才边控",
    "伊朗战争",
    "乌克兰战争",
    "跨境证券整治",
]


def _strip_fact_check_context_prefixes(text: str) -> str:
    value = re.sub(r"\s+", " ", str(text or "")).strip()
    if not value:
        return ""
    for _ in range(3):
        updated = value
        for prefix in FACT_CHECK_CONTEXT_PREFIXES:
            updated = re.sub(rf"^{re.escape(prefix)}\s*[：:]\s*", "", updated, flags=re.I).strip()
        if updated == value:
            break
        value = updated
    return value


def _is_generic_media_query(query: str) -> bool:
    value = re.sub(r"\s+", " ", str(query or "")).strip()
    if not value:
        return False
    value = re.sub(r"^site:[^\s]+\s+", "", value, flags=re.I).strip()
    return bool(
        re.fullmatch(
            r"(?:Bloomberg|Reuters|New York Times|Wall Street Journal|Financial Times)(?:\s+(?:world|markets|business|politics))+",
            value,
            flags=re.I,
        )
    )


def _generate_fact_check_english_queries(claim: str) -> list[str]:
    text = _strip_fact_check_context_prefixes(claim)
    text = re.sub(r"\s+", " ", str(text or "")).strip()
    if not text:
        return []
    candidates: list[str] = []
    seen: set[str] = set()

    def add(value: str) -> None:
        query = re.sub(r"\s+", " ", str(value or "")).strip(" ,，。；;：:")
        lowered = query.lower()
        if not query or lowered in seen:
            return
        seen.add(lowered)
        candidates.append(query)

    normalized = text
    for pattern, replacement in FACT_CHECK_ENGLISH_QUERY_REPLACEMENTS:
        normalized = re.sub(pattern, f" {replacement} ", normalized, flags=re.I)
    normalized = re.sub(r"[\u4e00-\u9fff]+", " ", normalized)
    normalized = re.sub(r"[，。、“”‘’（）()：:；;、]", " ", normalized)
    normalized = re.sub(r"\s+", " ", normalized).strip()
    topic_haystack = f"{text} {normalized}".strip()

    if re.search(r"(伊朗|Iran|核协议|nuclear deal|霍尔木兹|Strait of Hormuz|亚伯拉罕协议|Abraham Accords)", topic_haystack, re.I):
        add("Trump Iran nuclear deal talks sanctions progress")
        add("Iran Strait of Hormuz reopen clear mines draft agreement")
        add("Iran oil exports 60-day nuclear talks Reuters")
        if re.search(r"(亚伯拉罕协议|Abraham Accords|停火|ceasefire|沙特|Saudi|海湾|Gulf)", topic_haystack, re.I):
            add("Trump Abraham Accords ceasefire Saudi Arabia Iran Reuters")
    if re.search(r"(卢比奥|\bRubio\b)", topic_haystack, re.I):
        add("Rubio says U.S. will stop mediating Ukraine peace talks")
        add("Rubio Ukraine talks no progress Reuters")
        add("Rubio Ukraine mediation Kyiv Independent")
    elif re.search(r"(乌克兰|Ukraine|基辅|Kyiv|俄罗斯|Russia)", topic_haystack, re.I):
        add("Russia threatens strikes on Kyiv defence sites Reuters")
        add("Russia urges foreigners to leave Kyiv Reuters")
        add("Kyiv air defence Patriot supply Reuters")
    if re.search(r"(习近平|Xi Jinping|朝鲜|North Korea)", topic_haystack, re.I):
        add("Xi Jinping may visit North Korea Reuters")
        add("Xi visit Pyongyang next week Reuters")
        add("North Korea visit capital markets Reuters")
    if re.search(r"(宁德时代|CATL|香港IPO|Hong Kong IPO|摩根大通|JPMorgan|美国银行|Bank of America)", topic_haystack, re.I):
        add("CATL Hong Kong IPO JPMorgan Bank of America Reuters")
        add("House CCP committee CATL Hong Kong IPO")
        add("HKEX CATL Hong Kong listing filing")
    if re.search(r"(MLF|中期借贷便利|中国央行|人民银行|PBOC|People's Bank of China)", topic_haystack, re.I):
        add("PBOC cuts one-year MLF rate to 1.45% Reuters")
        add("People's Bank of China one-year MLF 1.45% Bloomberg")
        add("PBOC MLF operation banking net interest margin")
    if re.search(r"(富途|老虎|长桥|跨境证券|资金外流|资本外流|Futu|Tiger Brokers|Longbridge|capital outflow|cross-border brokerage)", topic_haystack, re.I):
        add("China Futu Tiger Brokers Longbridge crackdown Reuters")
        add("China cross-border brokerage crackdown Bloomberg")
        add("China capital outflow 2025 Reuters")
    if re.search(r"(DeepSeek|Alibaba|阿里巴巴|护照|passport|AI人才|AI talent|边控|exit control|Manus)", topic_haystack, re.I):
        add("China expands exit controls to AI talent Bloomberg")
        add("China AI talent passport controls Alibaba DeepSeek Reuters")
        add("Bloomberg China AI talent exit controls")
    if re.search(r"(雅虎财经|Yahoo Finance|沃尔玛|Walmart|Realty Income|菲利普莫里斯|Philip Morris)", topic_haystack, re.I):
        add("Yahoo Finance three blue-chip stocks Walmart Realty Income Philip Morris")
        add("Walmart Realty Income Philip Morris defensive stocks Yahoo Finance")
    if re.search(r"[A-Za-z]{3,}", normalized):
        add(normalized)
    return candidates[:8]


def _score_fact_check_query(query: str) -> int:
    value = re.sub(r"\s+", " ", str(query or "")).strip()
    if not value:
        return -999
    lowered = value.lower()
    score = 0
    if lowered.startswith("site:"):
        score += 20
    if re.search(r"\b(?:pce|gdp|cpi|pmi|bea|rbnz|rba|ecb|eurostat|destatis|insee|istat|statcan|boj|meti|hkex|catl|rubio|netanyahu|hormuz|kyiv|north korea|strait of hormuz|trump|iran|axios|state department|xi|jpmorgan)\b", lowered, re.I):
        score += 10
    if re.search(r"[A-Za-z]{3,}", value):
        score += 4
    if re.search(r"\b(?:19|20)\d{2}\b", value):
        score += 3
    if len(value) <= 90:
        score += 2
    if _is_generic_media_query(value):
        score -= 12
    if any(
        domain in lowered
        for domain in [
            "finance.yahoo.com",
            "pbc.gov.cn",
            "safe.gov.cn",
            "csrc.gov.cn",
        ]
    ):
        score += 7
    if re.search(
        r"\b(?:reuters|associated press|bloomberg|axios|ap|new york times|nyt|wall street journal|wsj|financial times)\b",
        lowered,
        re.I,
    ):
        score += 5
    if "kyivindependent.com" in lowered or re.search(r"\bkyiv independent\b", lowered, re.I):
        score += 6
    if re.search(r"\brubio\b", lowered, re.I) and re.search(r"\bukraine\b", lowered, re.I):
        score += 4
    if re.search(r"\bxi\b", lowered, re.I) and re.search(r"\bnorth korea\b", lowered, re.I):
        score += 4
    if re.search(r"\bsays\b", lowered, re.I):
        score += 2
    if re.search(r"\brubio\b", lowered, re.I) and re.search(r"\b(?:mediating|mediation|talks|peace)\b", lowered, re.I):
        if not re.search(r"\bukraine\b", lowered, re.I):
            score -= 8
    if re.search(r"\bxi\b", lowered, re.I) and re.search(r"\b(?:visit|visits|visiting)\b", lowered, re.I):
        if not re.search(r"\bnorth korea\b|\bpyongyang\b", lowered, re.I):
            score -= 8
    if re.search(r"\btrump\b", lowered, re.I) and re.search(r"\b(?:nuclear|sanctions|talks|deal)\b", lowered, re.I):
        if not re.search(r"\biran\b|\bhormuz\b", lowered, re.I):
            score -= 8
    return score


def _prepare_fact_check_queries(claim: str, queries: list[str] | None) -> list[str]:
    """为事实核查生成更精确的检索词，补充时间与站点限定。"""
    candidates: list[str] = []
    seen: set[str] = set()

    def add_query(value: str) -> None:
        query = _strip_fact_check_context_prefixes(value)
        query = re.sub(r"\s+", " ", str(query or "")).strip()
        if not query:
            return
        lowered = query.lower()
        if lowered in seen:
            return
        seen.add(lowered)
        candidates.append(query)

    claim_text = _strip_fact_check_context_prefixes(claim)
    add_query(claim_text)
    for item in queries or []:
        add_query(item)

    relaxed_query_seeds = [claim_text, *(str(item).strip() for item in (queries or []))]
    for raw_query in relaxed_query_seeds:
        if not raw_query:
            continue
        relaxed = raw_query
        relaxed = re.sub(r"\b(?:19|20)\d{2}[-/.年]\d{1,2}[-/.月]\d{1,2}日?(?:前后)?\b", " ", relaxed)
        relaxed = re.sub(r"\b(?:19|20)\d{2}年\d{1,2}月(?:\d{1,2}日)?(?:前后)?\b", " ", relaxed)
        relaxed = re.sub(r"\b\d{1,2}月\d{1,2}日(?:前后)?\b", " ", relaxed)
        relaxed = re.sub(r"(前后|上述|相关|有关|消息称|报道称|据称|被指|这番|言论|表态)", " ", relaxed)
        relaxed = re.sub(r"\s+", " ", relaxed).strip(" ,，。；;：:")
        if relaxed and relaxed != raw_query:
            add_query(relaxed)

    years = _extract_year_tokens(" ".join(candidates))
    english_queries = _generate_fact_check_english_queries(claim_text)
    for english_query in english_queries:
        add_query(english_query)
    authoritative_sources = _find_authoritative_sources(" ".join([*candidates, *english_queries]))

    for source in authoritative_sources[:5]:
        for query_term in source.get("query_terms", []) or []:
            add_query(str(query_term or "").strip())
            for year in years[:2]:
                if year and year not in str(query_term):
                    add_query(f"{query_term} {year}")

    for base_query in list(candidates):
        if re.search(r"(大选|选举|election)", base_query, re.I):
            add_query(f"{base_query} 官方结果")
            add_query(f"{base_query} official result")
        if re.search(r"(声明|通报|公告|会谈|访问|statement|announcement|meeting|visit)", base_query, re.I):
            add_query(f"{base_query} 官方声明")
            add_query(f"{base_query} official statement")
        if re.search(r"(国防部|外交部|政府|minister|ministry|government)", base_query, re.I):
            add_query(f"{base_query} 官网")
        if re.search(r"(央行|利率|加息|降息|inflation|cpi|通胀)", base_query, re.I):
            add_query(f"{base_query} rate decision")
            add_query(f"{base_query} inflation")
        if re.search(r"(gdp|国内生产总值|耐用品订单|工业生产|工业利润|pmi)", base_query, re.I):
            add_query(f"{base_query} official data")
        if re.search(r"(朝鲜|north korea|访问|visit)", base_query, re.I):
            add_query(f"{base_query} Reuters")
            add_query(f"{base_query} AP")
        if re.search(r"(伊朗|hormuz|霍尔木兹|核协议|谈判)", base_query, re.I):
            add_query(f"{base_query} Reuters")
            add_query(f"{base_query} Axios")
            add_query(f"{base_query} State Department")
        if re.search(r"(乌克兰|基辅|无人机|导弹|ukraine|kyiv|missile|drone)", base_query, re.I):
            add_query(f"{base_query} Reuters")
            add_query(f"{base_query} Kyiv Independent")
        if re.search(r"(香港ipo|香港上市|ipo|宁德时代|catl)", base_query, re.I):
            add_query(f"{base_query} HKEX filing")
            add_query(f"{base_query} Reuters")

        if not str(base_query).lower().startswith("site:"):
            for domain in _suggest_fact_check_major_media_domains(base_query):
                add_query(f"site:{domain} {base_query}")

        for year in years[:2]:
            if year not in base_query:
                add_query(f"{base_query} {year}")

        for source in authoritative_sources[:4]:
            source_url = str(source.get("url") or "").strip()
            domain = urlparse(source_url).netloc
            if domain:
                if not str(base_query).lower().startswith("site:"):
                    add_query(f"site:{domain} {base_query}")
                for query_term in source.get("query_terms", []) or []:
                    query_term = str(query_term or "").strip()
                    if query_term:
                        add_query(f"site:{domain} {query_term}")

    if len(candidates) <= 12:
        return candidates

    general_queries = [item for item in candidates if not item.lower().startswith("site:")]
    site_queries = [item for item in candidates if item.lower().startswith("site:")]
    general_queries.sort(key=lambda item: (-_score_fact_check_query(item), candidates.index(item)))
    site_queries.sort(key=lambda item: (-_score_fact_check_query(item), candidates.index(item)))

    selected: list[str] = []
    reserved_site_slots = min(4, len(site_queries))
    selected.extend(general_queries[: max(0, 12 - reserved_site_slots)])
    selected.extend(site_queries[:reserved_site_slots])
    for item in candidates:
        if len(selected) >= 12:
            break
        if item in selected:
            continue
        selected.append(item)
    return selected[:12]


def _select_fact_check_queries(queries: list[str], max_queries: int = 4) -> list[str]:
    """优先保留少量高价值检索词，避免串行搜索过慢且来源重复。"""
    if max_queries <= 0:
        return []

    normalized_queries: list[str] = []
    seen: set[str] = set()
    for query in queries or []:
        value = re.sub(r"\s+", " ", str(query or "")).strip()
        lowered = value.lower()
        if not value or lowered in seen:
            continue
        seen.add(lowered)
        normalized_queries.append(value)

    if len(normalized_queries) <= max_queries:
        return normalized_queries

    general_queries = [item for item in normalized_queries if not item.lower().startswith("site:")]
    site_queries = [item for item in normalized_queries if item.lower().startswith("site:")]
    general_queries.sort(key=lambda item: (-_score_fact_check_query(item), normalized_queries.index(item)))
    site_queries.sort(key=lambda item: (-_score_fact_check_query(item), normalized_queries.index(item)))
    selected: list[str] = []

    preferred_site_limit = min(3 if max_queries >= 4 else 1, len(site_queries))
    preferred_general_limit = min(len(general_queries), max_queries - preferred_site_limit)
    selected.extend(general_queries[:preferred_general_limit])
    for site_query in site_queries[:preferred_site_limit]:
        if len(selected) >= max_queries:
            break
        if site_query not in selected:
            selected.append(site_query)

    for item in normalized_queries:
        if len(selected) >= max_queries:
            break
        if item in selected:
            continue
        selected.append(item)
    return selected[:max_queries]


def _search_single_claim_source(item: dict, proxy_url: str = None, search_mode: str = "standard") -> dict:
    claim = str(item.get("claim") or "").strip()
    queries = item.get("queries") or []
    if not queries:
        queries = [claim]
    prepared_queries = _prepare_fact_check_queries(claim, queries[:3])
    mode = str(search_mode or "standard").strip().lower()
    max_queries = 2 if mode == "fast" else 4
    selected_queries = _select_fact_check_queries(prepared_queries, max_queries=max_queries)
    search_markdown = perform_web_search(selected_queries, proxy=proxy_url, claim_text=claim, search_mode=mode)
    return {
        "claim": claim,
        "queries": selected_queries,
        "search_markdown": search_markdown,
    }


def perform_web_search(queries: list[str], proxy: str = None, claim_text: str = "", search_mode: str = "standard") -> str:
    if not queries:
        return ""

    mode = str(search_mode or "standard").strip().lower()
    is_fast_mode = mode == "fast"
    anysearch_enabled = _env_flag("ANYSEARCH_ENABLED", default=False)
    anysearch_limit = 0 if is_fast_mode else _env_int("ANYSEARCH_MAX_CALLS_PER_CLAIM", 2, minimum=0, maximum=4)
    anysearch_calls = 0
    results_text = []
    global_source_links: list[str] = []
    seen_source_urls: set[str] = set()

    def render_candidate_hits(items: list[dict], heading: str) -> None:
        if not items:
            return
        results_text.append(heading)
        for item in items:
            title = str(item.get("title") or "").strip()
            url = str(item.get("url") or "").strip()
            if not title or not url:
                continue
            results_text.append(f"  - [{title}]({url})")
            detail_parts = []
            source_name = str(item.get("source") or "").strip()
            published_at = str(item.get("published_at") or "").strip()
            snippet = str(item.get("snippet") or "").strip()
            match_tokens = [
                str(token or "").strip()
                for token in (item.get("article_match_tokens") or [])
                if str(token or "").strip()
            ]
            if source_name:
                detail_parts.append(f"来源：{source_name}")
            recovered_from_source = str(item.get("recovered_from_source") or "").strip()
            if recovered_from_source:
                detail_parts.append(f"追源：由 {recovered_from_source} 聚合页恢复")
            if str(item.get("recovery_method") or "").strip() == "slug_guess":
                detail_parts.append("方式：Reuters 直链推测")
            if published_at:
                detail_parts.append(f"时间：{published_at}")
            if snippet:
                detail_parts.append(f"摘要：{snippet[:180]}")
            if match_tokens:
                detail_parts.append(f"正文匹配：{'、'.join(match_tokens[:4])}")
            if detail_parts:
                results_text.append(f"    - {'；'.join(detail_parts)}")
            normalized_url = _normalize_fact_check_source_url(url)
            if normalized_url and normalized_url not in seen_source_urls:
                seen_source_urls.add(normalized_url)
                label = source_name or title
                global_source_links.append(f"- [{label}]({url})")

    # 事实核查统一收敛为可人工复核的 Google/Bing 新闻检索入口，
    # 同时附上已识别到的官网/权威媒体站点，避免模型只输出模糊机构名。
    def add_search_links(q_term: str) -> bool:
        try:
            matched_sources = _find_authoritative_sources(" ".join([claim_text, q_term]))
            preferred_domains = {
                (urlparse(str(source.get("url") or "").strip()).netloc or "").lower()
                for source in matched_sources
                if str(source.get("url") or "").strip()
            }
            results_text.append(f"### 搜索关键字: {q_term}")
            results_text.append(f"- [Google 新闻核查]({_build_search_url(q_term, engine='google', news=True)})")
            results_text.append(f"- [Google 网页核查]({_build_search_url(q_term, engine='google', news=False)})")
            results_text.append(f"- [Bing 新闻核查]({_build_search_url(q_term, engine='bing', news=True)})")
            results_text.append(f"- [Bing 网页核查]({_build_search_url(q_term, engine='bing', news=False)})")
            anysearch_had_hits = False

            news_hits = _fetch_bing_news_results(q_term, proxy_url=proxy, max_items=2)
            should_merge_authoritative_google_hits = bool(preferred_domains) or bool(
                re.search(r"\breuters\b|\bassociated press\b|\bap\b|\bbloomberg\b", q_term, re.I)
            )
            if len(news_hits) < 1 or should_merge_authoritative_google_hits:
                google_hits = _fetch_google_news_results(
                    q_term,
                    proxy_url=proxy,
                    max_items=3 if should_merge_authoritative_google_hits else 2,
                )
                if preferred_domains:
                    google_hits = [
                        hit
                        for hit in google_hits
                        if any(_fact_check_hit_matches_domain(hit, domain) for domain in preferred_domains)
                    ]
                for news in google_hits:
                    normalized_url = _normalize_fact_check_source_url(str(news.get("url") or ""))
                    if not normalized_url:
                        continue
                    if any(
                        _normalize_fact_check_source_url(str(existing.get("url") or "")) == normalized_url
                        for existing in news_hits
                    ):
                        continue
                    news_hits.append(news)
                    if len(news_hits) >= (4 if should_merge_authoritative_google_hits else 3):
                        break
            news_hits = _rerank_fact_check_hits(
                news_hits,
                claim_text=claim_text,
                query_text=q_term,
                preferred_domains=preferred_domains,
                max_items=3,
            )
            news_hits = _recover_syndicated_fact_check_hits(
                news_hits,
                claim_text=claim_text,
                query_text=q_term,
                proxy_url=proxy,
                recovery_limit=1 if is_fast_mode else 2,
            )
            news_hits = _prune_fact_check_hits(
                news_hits,
                preferred_domains=preferred_domains or set(MAJOR_MEDIA_DOMAINS),
                max_items=3,
            )
            if news_hits:
                render_candidate_hits(news_hits, "- 命中的候选报道：")

            web_hits = _fetch_bing_web_results(q_term, proxy_url=proxy, max_items=6)
            if news_hits:
                existing_urls = {
                    _normalize_fact_check_source_url(str(item.get("url") or ""))
                    for item in news_hits
                }
                web_hits = [
                    item for item in web_hits
                    if _normalize_fact_check_source_url(str(item.get("url") or "")) not in existing_urls
                ]
            web_hits = _rerank_fact_check_hits(
                web_hits,
                claim_text=claim_text,
                query_text=q_term,
                preferred_domains=preferred_domains or set(MAJOR_MEDIA_DOMAINS),
                max_items=4,
            )
            web_hits = [
                item
                for item in web_hits
                if _is_relevant_fact_check_hit(item, claim_text=claim_text, query_text=q_term)
            ]
            web_hits = _recover_syndicated_fact_check_hits(
                web_hits,
                claim_text=claim_text,
                query_text=q_term,
                proxy_url=proxy,
                recovery_limit=0 if is_fast_mode else 1,
            )
            if not is_fast_mode:
                web_hits = _refine_fact_check_hits_with_article_text(
                    web_hits,
                    claim_text=claim_text,
                    query_text=q_term,
                    proxy_url=proxy,
                    article_fetch_limit=2,
                )
            web_hits = _prune_fact_check_hits(
                web_hits,
                preferred_domains=preferred_domains or set(MAJOR_MEDIA_DOMAINS),
                max_items=3,
            )
            if web_hits:
                render_candidate_hits(web_hits[:3], "- 命中的候选网页：")

            if (
                anysearch_enabled
                and anysearch_calls < anysearch_limit
                and _should_use_anysearch_for_fact_check(
                    claim_text,
                    q_term,
                    news_hits=news_hits,
                    web_hits=web_hits,
                    preferred_domains=preferred_domains,
                )
            ):
                anysearch_calls += 1
                anysearch_hits = _fetch_anysearch_results(q_term, proxy_url=proxy, max_items=3)
                existing_urls = {
                    _normalize_fact_check_source_url(str(item.get("url") or ""))
                    for item in [*(news_hits or []), *(web_hits or [])]
                }
                anysearch_hits = [
                    item
                    for item in anysearch_hits
                    if _normalize_fact_check_source_url(str(item.get("url") or "")) not in existing_urls
                ]
                anysearch_hits = _rerank_fact_check_hits(
                    anysearch_hits,
                    claim_text=claim_text,
                    query_text=q_term,
                    preferred_domains=preferred_domains or set(MAJOR_MEDIA_DOMAINS),
                    max_items=3,
                )
                anysearch_hits = _recover_syndicated_fact_check_hits(
                    anysearch_hits,
                    claim_text=claim_text,
                    query_text=q_term,
                    proxy_url=proxy,
                    recovery_limit=1,
                )
                anysearch_hits = _prune_fact_check_hits(
                    anysearch_hits,
                    preferred_domains=preferred_domains or set(MAJOR_MEDIA_DOMAINS),
                    max_items=3,
                )
                if anysearch_hits:
                    anysearch_had_hits = True
                    render_candidate_hits(anysearch_hits[:3], "- AnySearch 补强命中的候选来源：")

            if not news_hits and not web_hits and not anysearch_had_hits:
                results_text.append("- 自动检索结果：当前未抓到可直接引用的候选网页或报道。")
                results_text.append("  - 注意：这不等于“全网不存在相关内容”，只表示本轮自动检索暂未命中，后续应继续放宽搜索词，或改用英文/原文机构名继续核对。")

            if matched_sources:
                results_text.append("- 权威站点参考：")
                for source_idx, source in enumerate(matched_sources[:4]):
                    label = str(source.get("label") or "").strip()
                    url = str(source.get("url") or "").strip()
                    if not label or not url:
                        continue
                    results_text.append(f"  - [{label}]({url})")
                    normalized_url = _normalize_fact_check_source_url(url)
                    if normalized_url and normalized_url not in seen_source_urls:
                        seen_source_urls.add(normalized_url)
                        global_source_links.append(f"- [{label}]({url})")
                    domain = urlparse(url).netloc
                    if domain:
                        site_query_base = re.sub(r"^\s*site:[^\s]+\s+", "", q_term, flags=re.I).strip()
                        site_query = f"site:{domain} {site_query_base}".strip()
                        results_text.append(f"  - [{label} 定向搜索]({_build_search_url(site_query, engine='google', news=True)})")
                        if source_idx < (1 if is_fast_mode else 2):
                            site_hits = _fetch_bing_web_results(site_query, proxy_url=proxy, max_items=4)
                            site_hits = [
                                item for item in site_hits
                                if _normalize_fact_check_source_url(str(item.get("url") or "")) not in seen_source_urls
                            ]
                            if not site_hits:
                                google_site_hits = _fetch_google_news_results(site_query, proxy_url=proxy, max_items=4)
                                google_site_hits = [
                                    item for item in google_site_hits
                                    if _fact_check_hit_matches_domain(item, domain)
                                    and _normalize_fact_check_source_url(str(item.get("url") or "")) not in seen_source_urls
                                ]
                                site_hits = _recover_syndicated_fact_check_hits(
                                    google_site_hits,
                                    claim_text=claim_text,
                                    query_text=site_query,
                                    proxy_url=proxy,
                                    recovery_limit=2,
                                )
                            site_hits = _rerank_fact_check_hits(
                                site_hits,
                                claim_text=claim_text,
                                query_text=site_query,
                                preferred_domains={domain},
                                max_items=3,
                            )
                            site_hits = [
                                item
                                for item in site_hits
                                if _is_relevant_fact_check_hit(item, claim_text=claim_text, query_text=site_query)
                            ]
                            if not is_fast_mode:
                                site_hits = _refine_fact_check_hits_with_article_text(
                                    site_hits,
                                    claim_text=claim_text,
                                    query_text=site_query,
                                    proxy_url=proxy,
                                    article_fetch_limit=1,
                                )
                            site_hits = _prune_fact_check_hits(
                                site_hits,
                                preferred_domains={domain},
                                max_items=2,
                            )
                            if site_hits:
                                render_candidate_hits(site_hits[:2], f"- {label} 定向命中的候选页面：")
            results_text.append("")
            if len(seen_source_urls) >= 4:
                return True
        except Exception:
            return False
        return False

    for q in queries:
        if add_search_links(q):
            break

    if global_source_links:
        results_text.append("### 已识别的权威站点")
        results_text.extend(global_source_links[:8])
        results_text.append("")

    return "\n".join(results_text)


def _extract_json_candidate(text: str) -> str:
    text = (text or "").strip()
    if not text:
        return ""
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
        text = text.strip()
    start = text.find("{")
    end = text.rfind("}")
    if start != -1 and end != -1 and end > start:
        return text[start:end + 1]
    return text


def _extract_json_string_field(raw_text: str, field_names: list[str]) -> str:
    raw_text = raw_text or ""
    for field_name in field_names:
        pattern = rf'"{re.escape(field_name)}"\s*:\s*("(?:\\.|[^"\\])*")'
        match = re.search(pattern, raw_text, re.S)
        if not match:
            continue
        try:
            value = json.loads(match.group(1))
            return str(value or "").strip()
        except Exception:
            continue
    return ""


def _parse_summary_payload(raw_text: str):
    visited = set()
    candidates = [raw_text, _extract_json_candidate(raw_text)]
    while candidates:
        candidate = (candidates.pop(0) or "").strip()
        if not candidate or candidate in visited:
            continue
        visited.add(candidate)
        try:
            data = json.loads(candidate)
        except Exception:
            continue
        if isinstance(data, str):
            candidates.append(data)
            candidates.append(_extract_json_candidate(data))
            continue
        if isinstance(data, dict):
            return data
    summary_md = _extract_json_string_field(raw_text, ["summary_markdown", "summary", "summary_md"])
    fact_md = _extract_json_string_field(raw_text, ["fact_check_markdown", "fact_check", "factcheck_markdown", "fact_check_md"])
    if summary_md:
        return {
            "summary_markdown": summary_md,
            "fact_check_markdown": fact_md,
        }
    return None


def _normalize_output_locale(raw_locale: str | None) -> str:
    normalized = str(raw_locale or "").strip().lower().replace("_", "-")
    if normalized.startswith("en"):
        return "en"
    return "zh"


def _is_english_output_locale(raw_locale: str | None) -> bool:
    return _normalize_output_locale(raw_locale) == "en"


def _build_fact_check_placeholder(ui_locale: str | None = None) -> str:
    if _is_english_output_locale(ui_locale):
        return "- Structured fact-check output was not generated successfully. Please try again."
    return "- 暂未成功生成结构化事实核查结果，请重新尝试生成。"


FACT_CHECK_SECTION_SPLIT_PATTERN = (
    r"(?=^(?:###\s*(?:条目|Item)\s*\d+|\d+\.\s*"
    r"(?:新闻/声明|关键声明|声明|新闻|Claim|Statement|News Claim|News)[:：]))"
)
FACT_CHECK_STRUCTURED_ITEM_PATTERN = (
    r"^(?:###\s*(?:条目|Item)\s*\d+|\d+\.\s*"
    r"(?:新闻/声明|关键声明|声明|新闻|Claim|Statement|News Claim|News)[:：])"
)


def _fact_check_text(ui_locale: str | None, key: str) -> str:
    locale = _normalize_output_locale(ui_locale)
    messages = {
        "zh": {
            "item": "条目",
            "candidate_claim": "候选声明",
            "claim_label": "新闻/声明",
            "conclusion_label": "来源定位",
            "rationale_label": "来源线索说明",
            "pending_label": "建议继续查看",
            "sources_label": "来源出处",
            "source_links_label": "来源链接",
            "search_terms_label": "搜索词",
            "search_results_label": "搜索结果",
            "fallback_claim_intro": "系统已先整理出本条可核查说法，并附上可直接复核的候选来源线索。",
            "fallback_conclusion": "已整理候选来源，建议继续人工查看",
            "fallback_rationale": "本条先展示系统已汇总的搜索入口和候选来源，方便你直接点开原文自行判断，不代替用户下真假结论。",
            "fallback_pending": "建议优先核对原始报道时间、数字口径、机构原文、社媒原帖和二次转载是否存在偏差。",
            "fallback_no_links": "当前未提取到可点击来源链接。",
            "fallback_no_queries": "未生成搜索词",
            "fallback_detail_rationale": "系统已先保留候选搜索词 `{search_text}` 与可人工复核来源，便于你继续点开原文核对；本条内容不直接作为最终定论。",
            "fallback_detail_pending": "建议继续核对原始报道、官方披露、社媒原帖、统计口径与发布时间是否一致。",
            "progress_extract_claims": "正在抽取关键声明...",
            "progress_search_sources": "正在检索外部来源...",
            "progress_generate_fact_check": "正在整理逐条来源结果...",
            "progress_fact_check_done": "关键声明来源分析完成。",
            "system_prompt": "你是一个严谨的新闻来源导航助手。你的任务是整理出处、报道页、官网页和社媒原帖，不替用户判断真假。",
            "retry_system_prompt": "你是一个严谨且避免模板化的新闻来源导航助手。请基于现有候选来源更明确地指出哪些页面值得用户自行查看。",
            "insufficient_evidence": "缺乏证据",
            "dubious": "存疑",
            "soften_rationale": "本轮已检索到与该说法相关的公开网页线索，至少可以确认存在可继续核对的对应页面；当前更适合先说明已命中的网页来源，再补充仍待核对的正文细节、发布时间与原始出处。",
        },
        "en": {
            "item": "Item",
            "candidate_claim": "Candidate Claim",
            "claim_label": "Claim",
            "conclusion_label": "Source Status",
            "rationale_label": "Source Notes",
            "pending_label": "Suggested Follow-ups",
            "sources_label": "Source Links",
            "source_links_label": "Source Links",
            "search_terms_label": "Search Queries",
            "search_results_label": "Search Results",
            "fallback_claim_intro": "The system has already extracted a checkable claim and preserved candidate sources for quick review.",
            "fallback_conclusion": "Candidate sources collected for manual review",
            "fallback_rationale": "This draft surfaces current search leads and candidate sources so the user can open the original pages directly instead of relying on an automatic truth verdict.",
            "fallback_pending": "Prioritize checking the original report date, exact figures, issuing institution, original social post, and whether secondary reposts introduced distortions.",
            "fallback_no_links": "No clickable source links were extracted this time.",
            "fallback_no_queries": "No search queries were generated",
            "fallback_detail_rationale": "The system preserved candidate search queries `{search_text}` and manually reviewable sources so you can open the original pages directly; this draft is not a final verdict.",
            "fallback_detail_pending": "Continue checking the original report, official disclosures, original social post, statistical methodology, and publication time for consistency.",
            "progress_extract_claims": "Extracting key claims...",
            "progress_search_sources": "Searching external sources...",
            "progress_generate_fact_check": "Organizing itemized source results...",
            "progress_fact_check_done": "Key-claim source analysis completed.",
            "system_prompt": "You are a rigorous source-discovery assistant. Focus on locating source pages, reports, official pages, and social posts instead of issuing truth verdicts.",
            "retry_system_prompt": "You are a rigorous source-discovery assistant who avoids generic wording. Point to the most useful pages for manual review based on the available candidate sources.",
            "insufficient_evidence": "Insufficient Evidence",
            "dubious": "Uncertain",
            "soften_rationale": "This round already captured public web pages related to the claim, so the output should first explain which matching pages were found and then note which publication details, article body, or original source still need follow-up verification.",
        },
    }
    return messages[locale][key]


def _build_video_summary_system_prompt(ui_locale: str | None = None) -> str:
    if _is_english_output_locale(ui_locale):
        return (
            "You are a professional video summarization assistant. "
            "Always return valid JSON. The summary must be clear, concise, and written in English."
        )
    return "你是一个专业的视频内容总结助手。请始终返回合法 JSON。总结必须分条清晰。"


def _build_video_summary_prompt(content: str, current_date: str, ui_locale: str | None = None) -> str:
    if _is_english_output_locale(ui_locale):
        return (
            f"You are a professional video summarization assistant. Today's real date is: {current_date}.\n"
            "Please summarize the following video transcript.\n"
            "[Output format requirements]\n"
            "Return JSON only and include exactly two fields: `summary_markdown` and `fact_check_markdown`.\n"
            "1. `summary_markdown` requirements:\n"
            "   - Must be Markdown.\n"
            "   - Must use concise bullet points rather than long paragraphs.\n"
            "   - Use this fixed structure:\n"
            "     ## Core Topic\n"
            "     - Summarize the video's main idea in 1 sentence.\n"
            "     ## Main Points\n"
            "     - List 6-12 key points.\n"
            "     - Each bullet should cover exactly one fact, view, or judgment.\n"
            "     - If a point is speculative or opinion-based, explicitly mark it as `Video Opinion` or `Speculation`.\n"
            "     ## Key Details\n"
            "     - List important numbers, dates, people, institutions, policy names, and similar details.\n"
            "     ## Conclusion (Optional)\n"
            "     - Only include this section if the video clearly presents a conclusion.\n"
            "     - Do not force a conclusion when the video does not provide one.\n"
            "2. `fact_check_markdown` requirements:\n"
            "   - Must return an empty string `\"\"`.\n"
            "   - Do not generate any fact-check items.\n\n"
            "**Transcript input:**\n"
            f"{content}"
        )
    return (
        f"你是一个专业的视频内容总结助手。当前真实日期是：{current_date}。\n"
        "请总结以下视频字幕。\n"
        "【输出格式要求】\n"
        "请严格输出为 JSON 格式，且只能包含以下两个字段：`summary_markdown`、`fact_check_markdown`。\n"
        "1. `summary_markdown` 的要求：\n"
        "   - 必须是 Markdown。\n"
        "   - 必须按“逐条列点”的形式输出，不要写成长篇大段落。\n"
        "   - 固定结构如下：\n"
        "     ## 核心主题\n"
        "     - 1 句话概括视频主旨。\n"
        "     ## 主要内容\n"
        "     - 逐条列出 6-12 条要点。\n"
        "     - 每条只讲一个事实、观点或判断，语言清晰直接。\n"
        "     - 若某条是推测、判断、观点，请明确标注“视频观点”或“推测”。\n"
        "     ## 关键信息\n"
        "     - 列出数字、时间、人物、机构、政策名称等关键信息。\n"
        "     ## 结论（可选）\n"
        "     - 只有当视频确实提出了明确结论时才输出本节。\n"
        "     - 如果视频没有清晰结论，就不要硬写结论。\n"
        "2. `fact_check_markdown` 的要求：\n"
        "   - 必须返回空字符串 `\"\"`。\n"
        "   - 不要生成任何事实核查条目。\n\n"
        "**字幕内容输入：**\n"
        f"{content}"
    )


def _build_video_summary_repair_prompt(content_str: str, ui_locale: str | None = None) -> str:
    if _is_english_output_locale(ui_locale):
        return (
            "Repair the content below into valid JSON with exactly two fields: "
            "`summary_markdown` and `fact_check_markdown`.\n"
            "- `summary_markdown` must preserve the existing summary information and format it as clean Markdown.\n"
            "- `fact_check_markdown` must be an empty string.\n"
            "- Return JSON only, with no explanation.\n\n"
            f"Original content:\n{content_str}"
        )
    return (
        "请将下面内容修复为合法 JSON，并且只能包含两个字段："
        "`summary_markdown` 和 `fact_check_markdown`。\n"
        "- `summary_markdown` 必须保留现有总结信息，并整理为清晰 Markdown。\n"
        "- `fact_check_markdown` 必须返回空字符串，不要补任何事实核查内容。\n"
        "- 只返回 JSON，不要解释。\n\n"
        f"原始内容：\n{content_str}"
    )


def _build_document_summary_system_prompt(ui_locale: str | None = None) -> str:
    if _is_english_output_locale(ui_locale):
        return (
            "You are a professional document summarization assistant. "
            "Always return valid JSON with only the field `summary_markdown`. "
            "Use Markdown and follow this fixed structure: `## Core Topic`, `## Main Points`, `## Key Details`, and `## Conclusion (Optional)`. "
            "Prefer concise bullet points instead of long paragraphs."
        )
    return (
        "你是一个专业的文档总结助手。请始终返回合法 JSON，且只能包含字段 `summary_markdown`。"
        "输出必须使用 Markdown，结构固定为：`## 核心主题`、`## 主要内容`、`## 关键信息`、`## 结论（可选）`。"
        "所有内容尽量分条列出，不要写成长篇大段落。"
    )


def _normalize_summary_payload(payload: dict | None, ui_locale: str | None = None) -> dict | None:
    if not isinstance(payload, dict):
        return None
    summary_text = (
        payload.get("summary_markdown")
        or payload.get("summary")
        or payload.get("summary_md")
        or ""
    )
    fact_check_text = (
        payload.get("fact_check_markdown")
        or payload.get("fact_check")
        or payload.get("factcheck_markdown")
        or payload.get("fact_check_md")
        or ""
    )
    summary_text = str(summary_text or "").strip()
    fact_check_text = str(fact_check_text or "").strip()
    if not summary_text:
        return None
    if not fact_check_text:
        fact_check_text = _build_fact_check_placeholder(ui_locale)
    return {
        "summary_markdown": summary_text,
        "fact_check_markdown": fact_check_text,
    }


def _is_placeholder_fact_check(text: str) -> bool:
    """判断事实核查内容是否仍是空洞占位文案。"""
    value = str(text or "").strip()
    if not value:
        return True
    placeholder_markers = [
        "暂未成功生成结构化事实核查结果",
        "未成功拆出结构化事实核查结果",
        "AI 未能稳定输出结构化事实核查结果",
        "模型未返回事实核查结果",
        "Structured fact-check output was not generated successfully",
        "Could not extract structured fact-check output",
        "AI failed to produce stable structured fact-check output",
        "The model did not return fact-check output",
    ]
    return any(marker in value for marker in placeholder_markers)


def _extract_markdown_links(markdown_text: str) -> list[tuple[str, str]]:
    text = str(markdown_text or "")
    if not text:
        return []
    pairs: list[tuple[str, str]] = []
    seen: set[str] = set()
    for label, url in re.findall(r"\[([^\]]+)\]\((https?://[^)\s]+)\)", text):
        clean_label = label.strip()
        clean_url = url.strip()
        if not clean_url or clean_url in seen:
            continue
        seen.add(clean_url)
        pairs.append((clean_label or clean_url, clean_url))
    return pairs


def _extract_fact_check_source_links(markdown_text: str) -> list[tuple[str, str]]:
    """仅保留实际可核查的外部来源，过滤搜索引擎结果页链接。"""
    filtered: list[tuple[str, str]] = []
    for label, url in _extract_markdown_links(markdown_text):
        domain = str(urlparse(url).netloc or "").lower()
        if not domain:
            continue
        if any(token in domain for token in ("google.", "bing.")):
            continue
        filtered.append((label, url))
    return _dedupe_fact_check_source_links(filtered)


def _extract_fact_check_hit_context(
    search_markdown: str,
    extra_links: list[tuple[str, str]] | None = None,
) -> dict[str, object]:
    """从搜索结果与已注入来源里提炼“找到了哪些类型的网页”。"""
    search_text = str(search_markdown or "")
    candidate_links = _extract_fact_check_source_links(search_text)
    if extra_links:
        candidate_links = _dedupe_fact_check_source_links(candidate_links + list(extra_links))

    has_candidate_news_hits = "命中的候选报道" in search_text
    has_candidate_web_hits = "命中的候选网页" in search_text
    has_site_hits = "定向命中的候选页面" in search_text
    has_authoritative_link = any(
        any(token in (url or "").lower() for token in [
            ".gov",
            "state.gov",
            "gov.",
            "mod.gov",
            "reuters.com",
            "apnews.com",
            "bloomberg.com",
            "nytimes.com",
            "washingtonpost.com",
        ])
        for _, url in candidate_links
    )
    preview_labels: list[str] = []
    seen_preview_tokens: set[str] = set()
    for label, url in candidate_links:
        parsed = urlparse(url)
        preview = str(label or "").strip() or str(parsed.netloc or "").lower()
        preview = re.sub(r"\s+", " ", preview).strip(" -")
        if not preview:
            continue
        preview_key = preview.lower()
        if preview_key in seen_preview_tokens:
            continue
        seen_preview_tokens.add(preview_key)
        preview_labels.append(preview)
        if len(preview_labels) >= 3:
            break
    return {
        "candidate_links": candidate_links,
        "candidate_link_count": len(candidate_links),
        "has_candidate_news_hits": has_candidate_news_hits,
        "has_candidate_web_hits": has_candidate_web_hits,
        "has_site_hits": has_site_hits,
        "has_authoritative_link": has_authoritative_link,
        "preview_labels": preview_labels,
    }


def _build_fact_check_hit_rationale(
    *,
    search_markdown: str,
    section_text: str = "",
    ui_locale: str = "zh",
) -> str:
    """根据命中的网页类型生成更贴题的判断依据，而不是复用固定模板。"""
    section_links = _extract_fact_check_source_links(section_text)
    hit_context = _extract_fact_check_hit_context(search_markdown, extra_links=section_links)
    candidate_link_count = int(hit_context.get("candidate_link_count") or 0)
    if candidate_link_count <= 0:
        return _fact_check_text(ui_locale, "soften_rationale")

    preview_labels = [str(item).strip() for item in hit_context.get("preview_labels") or [] if str(item).strip()]
    if ui_locale == "en":
        source_kinds: list[str] = []
        if hit_context.get("has_candidate_news_hits"):
            source_kinds.append("news reports")
        if hit_context.get("has_candidate_web_hits"):
            source_kinds.append("general web pages")
        if hit_context.get("has_site_hits") or hit_context.get("has_authoritative_link"):
            source_kinds.append("official or outlet pages")
        source_kind_text = ", ".join(source_kinds) if source_kinds else "public web pages"
        preview_text = ", ".join(preview_labels)
        preview_clause = f" including {preview_text}" if preview_text else ""
        return (
            f"This round already found {candidate_link_count} matching {source_kind_text}{preview_clause}. "
            "At minimum, there are public pages corresponding to the claim; the next step is to verify whether the body text, publication time, figures, and original sourcing fully match the video's wording."
        )

    source_kinds: list[str] = []
    if hit_context.get("has_candidate_news_hits"):
        source_kinds.append("新闻报道")
    if hit_context.get("has_candidate_web_hits"):
        source_kinds.append("普通网页")
    if hit_context.get("has_site_hits") or hit_context.get("has_authoritative_link"):
        source_kinds.append("官网/机构页面")
    source_kind_text = "、".join(source_kinds) if source_kinds else "公开网页"
    preview_text = "、".join(preview_labels)
    preview_clause = f"，例如 {preview_text}" if preview_text else ""
    return (
        f"本轮已检索到 {candidate_link_count} 个与该说法对应的{source_kind_text}{preview_clause}，"
        "至少可以确认网上存在可继续核对的相关页面；下一步应继续比对正文细节、发布时间、数字口径与原始出处是否完全一致。"
    )


def _soften_fact_check_wording(text: str) -> str:
    value = str(text or "")
    if not value:
        return value
    replacements = [
        (r"手动搜索未发现", "当前检索到的公开候选来源暂未形成直接佐证"),
        (r"未发现支持", "暂缺可直接支撑"),
        (r"未搜索到关于", "当前检索到的公开候选来源中，暂缺关于"),
        (r"未搜索到来自", "当前检索到的公开候选来源中，暂缺来自"),
        (r"未检索到可靠外部来源", "当前可交叉验证的公开候选来源仍不足"),
        (r"未找到关于", "当前检索到的公开候选来源中，暂缺关于"),
        (r"未找到来自", "当前检索到的公开候选来源中，暂缺来自"),
        (r"未发现关于", "当前检索到的公开候选来源中，暂缺关于"),
    ]
    for pattern, replacement in replacements:
        value = re.sub(pattern, replacement, value)
    return value


def _soften_absolute_negative_fact_check_with_sources(
    text: str,
    source_links: list[tuple[str, str]] | None,
) -> str:
    """当条目里已附带候选来源链接时，避免保留“未发现任何报道”这类绝对否定表述。"""
    value = str(text or "")
    if not value or not source_links:
        return value

    link_text = " ".join(f"{label} {url}" for label, url in source_links).lower()
    has_news_evidence = any(
        token in link_text
        for token in [
            "reuters",
            "路透",
            "bloomberg",
            "彭博",
            "twse",
            "taiwan stock exchange",
            "台灣證券交易所",
            "台湾证券交易所",
            "my-formosa",
            "美麗島",
            "美丽岛",
        ]
    )
    if not has_news_evidence:
        return value

    replacements = [
        (
            r"未发现任何主流财经新闻媒体[^。\n]*",
            "当前候选来源已包含主流财经媒体或权威站点链接，不宜直接下结论为“未发现报道”，应以这些链接的具体内容为准",
        ),
        (
            r"未发现来自[^。\n]*的报道",
            "当前候选来源已提供可进一步核对的相关报道或权威站点链接，不宜直接下结论为“未发现报道”",
        ),
        (
            r"未发现[^。\n]*证明",
            "当前候选来源已提供可进一步核对的相关链接，请直接依据链接内容判断是否构成证明",
        ),
    ]
    for pattern, replacement in replacements:
        value = re.sub(pattern, replacement, value)
    return value


def _enrich_fact_check_markdown_with_links(fact_md: str, search_results_md: str) -> str:
    fact_text = str(fact_md or "").strip()
    if not fact_text:
        return fact_text
    search_links = _extract_fact_check_source_links(search_results_md)
    if not search_links:
        return _soften_fact_check_wording(fact_text)
    return _enrich_fact_check_items_with_claim_sources(
        fact_text,
        [{"search_markdown": search_results_md}],
    )


def _enrich_fact_check_items_with_claim_sources(fact_md: str, claim_sources: list[dict] | None) -> str:
    fact_text = str(fact_md or "").strip()
    if not fact_text:
        return fact_text
    if not claim_sources:
        return fact_text

    # 兼容两种常见输出：
    # 1. `### 条目1` 这种结构化块
    # 2. `1. 新闻/声明：...` 这种编号列表
    sections = re.split(
        FACT_CHECK_SECTION_SPLIT_PATTERN,
        _soften_fact_check_wording(fact_text),
        flags=re.M,
    )
    updated_sections: list[str] = []
    claim_idx = 0

    for section in sections:
        stripped = section.strip()
        if not stripped:
            continue
        is_structured_item = bool(
            re.match(FACT_CHECK_STRUCTURED_ITEM_PATTERN, stripped)
        )
        if not is_structured_item:
            updated_sections.append(stripped)
            continue

        if len(claim_sources) == 1:
            current_sources = claim_sources[0]
        else:
            current_sources = claim_sources[claim_idx] if claim_idx < len(claim_sources) else {}
            claim_idx += 1
        source_links = _extract_fact_check_source_links(str(current_sources.get("search_markdown") or ""))
        if not source_links:
            updated_sections.append(stripped)
            continue

        if re.search(r"\[[^\]]+\]\(https?://", stripped):
            updated_sections.append(stripped)
            continue

        source_text = "；".join(f"[{label}]({url})" for label, url in source_links[:3])
        stripped = _soften_absolute_negative_fact_check_with_sources(stripped, source_links)
        source_label = _fact_check_text(
            "en" if re.search(r"(?:^###\s*Item\s*\d+|(?:^|\n)-\s*Sources?[:：])", stripped, re.I) else "zh",
            "sources_label",
        )
        source_links_label = _fact_check_text(
            "en" if source_label in {"Sources", "Source Links"} else "zh",
            "source_links_label",
        )

        # 如果模型已经输出了来源字段但没有链接，保留原有文字说明并额外补一行可点击链接。
        if re.search(r"(?:来源/出处|来源出处|Sources?|Source Links?)[:：]", stripped, re.I):
            stripped = re.sub(
                r"((?:来源/出处|来源出处|Sources?|Source Links?)[:：][^\n]*)",
                lambda m: f"{m.group(1)}\n- {source_links_label}: {source_text}",
                stripped,
                count=1,
                flags=re.I,
            )
        else:
            stripped = stripped.rstrip() + f"\n- {source_label}: {source_text}"
        updated_sections.append(stripped)

    return "\n\n".join(updated_sections)


def _normalize_fact_check_conclusions_with_sources(fact_md: str, claim_sources: list[dict] | None) -> str:
    """有候选网页时，避免把“已找到对应页面”的场景误写成统一的空泛结论。

    这里不把条目强行改为“属实”，只把更合适的场景收敛为“存疑”：
    - 搜索上下文里已出现候选报道、候选网页或定向站点页面；
    - 但模型仍输出“缺乏证据”；
    - 说明更接近“已找到对应网页，仍需继续核对细节”。
    """
    fact_text = str(fact_md or "").strip()
    if not fact_text or not claim_sources:
        return fact_text

    sections = re.split(
        FACT_CHECK_SECTION_SPLIT_PATTERN,
        fact_text,
        flags=re.M,
    )
    updated_sections: list[str] = []
    claim_idx = 0

    for section in sections:
        stripped = section.strip()
        if not stripped:
            continue
        is_structured_item = bool(
            re.match(FACT_CHECK_STRUCTURED_ITEM_PATTERN, stripped)
        )
        if not is_structured_item:
            updated_sections.append(stripped)
            continue

        if len(claim_sources) == 1:
            current_sources = claim_sources[0]
        else:
            current_sources = claim_sources[claim_idx] if claim_idx < len(claim_sources) else {}
            claim_idx += 1

        search_markdown = str(current_sources.get("search_markdown") or "")
        section_source_links = _extract_fact_check_source_links(stripped)
        hit_context = _extract_fact_check_hit_context(search_markdown, extra_links=section_source_links)
        candidate_link_count = int(hit_context.get("candidate_link_count") or 0)
        has_candidate_news_hits = bool(hit_context.get("has_candidate_news_hits")) and candidate_link_count >= 1
        has_candidate_web_hits = (
            bool(hit_context.get("has_candidate_web_hits"))
            or bool(hit_context.get("has_site_hits"))
        ) and candidate_link_count >= 1
        has_source_links = bool(re.search(r"\[[^\]]+\]\(https?://", stripped))
        has_authoritative_link = bool(hit_context.get("has_authoritative_link"))
        has_dynamic_fact_markers = bool(
            re.search(
                r"(市值|油价|股市|汇率|利率|GDP|CPI|失业率|销量|排名|第[一二三四五六七八九十\d]+|达到|升至|超越|\b\d+(?:\.\d+)?(?:万亿|亿|%|美元|元|桶)\b)",
                stripped,
                re.I,
            )
        )
        has_partial_support_rationale = any(
            marker in stripped
            for marker in [
                "需以",
                "需要来自",
                "官方数据",
                "实时数据",
                "权威金融数据机构",
                "动态变化",
                "未提供可验证的具体数据来源",
                "未能找到",
                "直接证实",
            ]
        )
        should_soften_to_dubious = (
            has_candidate_news_hits
            or has_candidate_web_hits
            or candidate_link_count >= 2
            or has_authoritative_link
            or (has_source_links and has_dynamic_fact_markers and has_partial_support_rationale)
        )
        if re.search(r"^\s*###\s*Item\s*\d+", stripped, re.I) or re.search(r"(?:Conclusion|Source Status):", stripped):
            conclusion_label = "Source Status" if "Source Status:" in stripped else "Conclusion"
            insufficient_value = "No Direct Source Yet" if conclusion_label == "Source Status" else "Insufficient Evidence"
            dubious_value = "Related Coverage Located" if conclusion_label == "Source Status" else "Uncertain"
        else:
            conclusion_label = "来源定位" if "来源定位" in stripped else "核查结论"
            insufficient_value = "暂未定位到直接来源" if conclusion_label == "来源定位" else "缺乏证据"
            dubious_value = "已找到相关来源" if conclusion_label == "来源定位" else "存疑"
        rationale_insert = _build_fact_check_hit_rationale(
            search_markdown=search_markdown,
            section_text=stripped,
            ui_locale="en" if conclusion_label == "Conclusion" else "zh",
        )
        rationale_line_pattern = r"(^-\s*(?:判断依据|依据|来源线索说明|Rationale|Source Notes)[:：]\s*)([^\n]*)"
        if not should_soften_to_dubious or f"{conclusion_label}：" + insufficient_value not in stripped and f"{conclusion_label}: {insufficient_value}" not in stripped:
            updated_sections.append(stripped)
            continue

        stripped = re.sub(
            rf"({re.escape(conclusion_label)}[:：]\s*){re.escape(insufficient_value)}",
            rf"\1{dubious_value}",
            stripped,
            count=1,
            flags=re.I,
        )
        generic_rationale_tokens = [
            "交易所公告",
            "原始统计口径",
            "一手权威来源",
            "更适合判为“存疑”",
            "indirect support",
            "exchange filing",
            "original statistical methodology",
            "primary authoritative source",
        ]
        if rationale_insert:
            if re.search(rationale_line_pattern, stripped, re.I | re.M):
                replace_existing_rationale = any(token in stripped for token in generic_rationale_tokens)
                if replace_existing_rationale:
                    stripped = re.sub(
                        rationale_line_pattern,
                        lambda m: f"{m.group(1)}{rationale_insert}",
                        stripped,
                        count=1,
                        flags=re.I | re.M,
                    )
                elif rationale_insert not in stripped:
                    stripped = re.sub(
                        rationale_line_pattern,
                        lambda m: (
                            f"{m.group(1)}{rationale_insert}；{m.group(2).strip()}"
                            if str(m.group(2) or "").strip()
                            else f"{m.group(1)}{rationale_insert}"
                        ),
                        stripped,
                        count=1,
                        flags=re.I | re.M,
                    )
            elif rationale_insert not in stripped:
                stripped = stripped.rstrip() + (
                    f"\n- {'Source Notes' if conclusion_label == 'Conclusion' else '来源线索说明'}: {rationale_insert}"
                )
        updated_sections.append(stripped)

    return "\n\n".join(updated_sections)


def _build_fact_check_fallback_markdown(
    *,
    claim_sources: list[dict] | None = None,
    search_results_md: str = "",
    ui_locale: str | None = None,
) -> str:
    """在模型未稳定输出时，基于候选来源生成最小可用的事实核查稿。"""
    item_label = _fact_check_text(ui_locale, "item")
    claim_label = _fact_check_text(ui_locale, "claim_label")
    conclusion_label = _fact_check_text(ui_locale, "conclusion_label")
    rationale_label = _fact_check_text(ui_locale, "rationale_label")
    pending_label = _fact_check_text(ui_locale, "pending_label")
    sources_label = _fact_check_text(ui_locale, "sources_label")
    candidate_claim_label = _fact_check_text(ui_locale, "candidate_claim")
    search_terms_label = _fact_check_text(ui_locale, "search_terms_label")
    sections = [
        f"### {item_label}1",
        f"- {claim_label}: {_fact_check_text(ui_locale, 'fallback_claim_intro')}",
        f"- {conclusion_label}: {_fact_check_text(ui_locale, 'fallback_conclusion')}",
        f"- {rationale_label}: {_fact_check_text(ui_locale, 'fallback_rationale')}",
        f"- {pending_label}: {_fact_check_text(ui_locale, 'fallback_pending')}",
    ]

    if claim_sources:
        rendered_sections = _render_fact_check_claim_sections(
            claim_sources,
            start_index=1,
            ui_locale=ui_locale,
        )
        if rendered_sections:
            return "\n\n".join(rendered_sections)

    search_links = _extract_fact_check_source_links(search_results_md)
    if search_links:
        source_text = "；".join(f"[{label}]({url})" for label, url in search_links[:8])
        sections.append(f"- {sources_label}: {source_text}")
    else:
        sections.append(f"- {sources_label}: {_fact_check_text(ui_locale, 'fallback_no_links')}")
    return "\n".join(sections)


def _render_fact_check_claim_sections(
    claim_sources: list[dict] | None,
    *,
    start_index: int = 1,
    ui_locale: str | None = None,
) -> list[str]:
    if not claim_sources:
        return []

    item_label = _fact_check_text(ui_locale, "item")
    claim_label = _fact_check_text(ui_locale, "claim_label")
    conclusion_label = _fact_check_text(ui_locale, "conclusion_label")
    rationale_label = _fact_check_text(ui_locale, "rationale_label")
    pending_label = _fact_check_text(ui_locale, "pending_label")
    sources_label = _fact_check_text(ui_locale, "sources_label")
    candidate_claim_label = _fact_check_text(ui_locale, "candidate_claim")
    rendered_sections: list[str] = []

    for offset, item in enumerate(claim_sources):
        item_index = start_index + offset
        claim = str(item.get("claim") or f"{candidate_claim_label}{item_index}").strip()
        query_list = item.get("queries") or []
        search_md = str(item.get("search_markdown") or "")
        links = _extract_fact_check_source_links(search_md)
        source_text = "；".join(f"[{label}]({url})" for label, url in links[:5]) if links else _fact_check_text(ui_locale, "fallback_no_links")
        search_text = " | ".join(str(query).strip() for query in query_list if str(query).strip()) or _fact_check_text(ui_locale, "fallback_no_queries")
        rendered_sections.append(
            "\n".join(
                [
                    f"### {item_label}{item_index}",
                    f"- {claim_label}: {claim}",
                    f"- {conclusion_label}: {_fact_check_text(ui_locale, 'fallback_conclusion')}",
                    f"- {rationale_label}: {_fact_check_text(ui_locale, 'fallback_detail_rationale').format(search_text=search_text)}",
                    f"- {pending_label}: {_fact_check_text(ui_locale, 'fallback_detail_pending')}",
                    f"- {sources_label}: {source_text}",
                ]
            )
        )
    return rendered_sections


def _count_structured_fact_check_items(fact_md: str) -> int:
    sections = re.split(FACT_CHECK_SECTION_SPLIT_PATTERN, str(fact_md or "").strip(), flags=re.M)
    return sum(
        1
        for section in sections
        if section.strip() and re.match(FACT_CHECK_STRUCTURED_ITEM_PATTERN, section.strip())
    )


def _ensure_fact_check_item_coverage(
    fact_md: str,
    claim_sources: list[dict] | None,
    *,
    ui_locale: str | None = None,
) -> str:
    fact_text = str(fact_md or "").strip()
    if not fact_text or not claim_sources:
        return fact_text

    expected_items = len([item for item in claim_sources if str(item.get("claim") or "").strip()])
    if expected_items <= 0:
        return fact_text

    actual_items = _count_structured_fact_check_items(fact_text)
    if actual_items <= 0:
        return _build_fact_check_fallback_markdown(claim_sources=claim_sources, ui_locale=ui_locale)
    if actual_items >= expected_items:
        return fact_text

    missing_sections = _render_fact_check_claim_sections(
        claim_sources[actual_items:expected_items],
        start_index=actual_items + 1,
        ui_locale=ui_locale,
    )
    if not missing_sections:
        return fact_text
    return fact_text.rstrip() + "\n\n" + "\n\n".join(missing_sections)


SUPPORTED_DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".markdown", ".pptx"}
DEFAULT_DOCUMENT_MAX_MB = int(os.environ.get("DOCUMENT_MAX_UPLOAD_MB", "20") or "20")
DEFAULT_DOCUMENT_DIRECT_SUMMARY_CHARS = int(os.environ.get("DOCUMENT_DIRECT_SUMMARY_CHARS", "15000") or "15000")
DEFAULT_DOCUMENT_CHUNK_CHARS = int(os.environ.get("DOCUMENT_CHUNK_CHARS", "12000") or "12000")


def validate_document_upload(file_name: str, file_size: int, max_size_mb: int = DEFAULT_DOCUMENT_MAX_MB) -> tuple[bool, str]:
    file_name = str(file_name or "").strip()
    suffix = os.path.splitext(file_name)[1].lower()
    if not file_name or suffix not in SUPPORTED_DOCUMENT_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_DOCUMENT_EXTENSIONS))
        return False, f"暂不支持该文档格式。当前支持：{supported}"
    if int(file_size or 0) <= 0:
        return False, "上传文件为空，请重新选择。"
    if int(file_size or 0) > max_size_mb * 1024 * 1024:
        return False, f"文档超过大小限制（{max_size_mb}MB），请拆分后再上传。"
    return True, ""


def clean_document_text(text: str) -> str:
    text = str(text or "")
    text = text.replace("\u00a0", " ").replace("\r\n", "\n").replace("\r", "\n")
    lines = []
    blank_pending = False
    for raw_line in text.split("\n"):
        line = re.sub(r"[ \t]+", " ", raw_line).strip()
        if not line:
            if not blank_pending:
                lines.append("")
                blank_pending = True
            continue
        blank_pending = False
        lines.append(line)
    cleaned = "\n".join(lines)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
    return cleaned


def _extract_pdf_text(file_bytes: bytes) -> dict:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise RuntimeError("未安装 pypdf，无法解析 PDF。") from e

    reader = PdfReader(BytesIO(file_bytes))
    page_texts = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            page_text = str(page.extract_text() or "").strip()
        except Exception:
            page_text = ""
        if page_text:
            page_texts.append(f"## 第{idx}页\n{page_text}")
    raw_text = "\n\n".join(page_texts).strip()
    ocr_used = False
    if not raw_text:
        ocr_text = _extract_pdf_text_via_ocr(file_bytes)
        if ocr_text:
            raw_text = ocr_text
            ocr_used = True

    return {
        "raw_text": raw_text,
        "page_count": len(reader.pages),
        "section_count": len(page_texts),
        "ocr_used": ocr_used,
    }


def _extract_docx_text(file_bytes: bytes) -> dict:
    try:
        from docx import Document
    except ImportError as e:
        raise RuntimeError("未安装 python-docx，无法解析 DOCX。") from e

    document = Document(BytesIO(file_bytes))
    parts = []
    paragraph_count = 0
    for para in document.paragraphs:
        text = str(para.text or "").strip()
        if not text:
            continue
        paragraph_count += 1
        style_name = str(getattr(getattr(para, "style", None), "name", "") or "").lower()
        heading_match = re.search(r"heading\s*(\d+)", style_name)
        if heading_match:
            level = max(1, min(6, int(heading_match.group(1))))
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)

    table_count = 0
    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [str(cell.text or "").strip() for cell in row.cells]
            cells = [cell for cell in cells if cell]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            table_count += 1
            parts.append(f"## 表格{table_count}\n" + "\n".join(rows))

    return {
        "raw_text": "\n\n".join(parts).strip(),
        "page_count": None,
        "section_count": paragraph_count + table_count,
    }


def _extract_pptx_text(file_bytes: bytes) -> dict:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("未安装 python-pptx，无法解析 PPTX。") from e

    prs = Presentation(BytesIO(file_bytes))
    slide_parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        shape_texts = []
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                shape_texts.append(text)
                continue
            table = getattr(shape, "table", None)
            if table:
                rows = []
                for row in table.rows:
                    cells = [str(cell.text or "").strip() for cell in row.cells]
                    cells = [cell for cell in cells if cell]
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    shape_texts.append("\n".join(rows))
        if shape_texts:
            slide_parts.append(f"## 第{idx}页幻灯片\n" + "\n\n".join(shape_texts))

    return {
        "raw_text": "\n\n".join(slide_parts).strip(),
        "page_count": len(prs.slides),
        "section_count": len(slide_parts),
    }


def _extract_plain_text(file_bytes: bytes) -> dict:
    encodings = ["utf-8", "utf-8-sig", "gb18030", "gbk", "big5", "latin-1"]
    last_error = None
    for encoding in encodings:
        try:
            return {
                "raw_text": file_bytes.decode(encoding).strip(),
                "page_count": None,
                "section_count": None,
            }
        except Exception as e:
            last_error = e
    raise RuntimeError(f"文本文件解码失败：{last_error}")


def _extract_pdf_text_via_ocr(file_bytes: bytes) -> str:
    if str(os.environ.get("DOCUMENT_PDF_OCR_ENABLED", "1") or "1").strip().lower() in {"0", "false", "no"}:
        return ""
    try:
        import fitz
        import numpy as np
        from rapidocr_onnxruntime import RapidOCR
    except ImportError:
        return ""

    try:
        ocr_engine = RapidOCR()
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        page_texts = []
        max_pages = int(os.environ.get("DOCUMENT_PDF_OCR_MAX_PAGES", "20") or "20")
        for idx, page in enumerate(doc, start=1):
            if idx > max_pages:
                break
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            channels = 3 if pix.n >= 3 else 1
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
            if channels == 1:
                img = img[:, :, 0]
            elif pix.n > 3:
                img = img[:, :, :3]
            result, _ = ocr_engine(img)
            if not result:
                continue
            page_lines = []
            for item in result:
                try:
                    text = str(item[1] or "").strip()
                except Exception:
                    text = ""
                if text:
                    page_lines.append(text)
            if page_lines:
                page_texts.append(f"## 第{idx}页\n" + "\n".join(page_lines))
        return "\n\n".join(page_texts).strip()
    except Exception as e:
        print(f"PDF OCR extraction failed: {e}")
        return ""


def _guess_remote_file_name(url: str, response=None) -> str:
    parsed = urlparse(str(url or "").strip())
    candidate = os.path.basename(parsed.path or "").strip()
    if candidate:
        return candidate
    content_type = str(getattr(response, "headers", {}).get("Content-Type", "") or "").lower()
    if "pdf" in content_type:
        return "remote.pdf"
    if "presentation" in content_type or "pptx" in content_type:
        return "remote.pptx"
    if "wordprocessingml" in content_type or "docx" in content_type:
        return "remote.docx"
    if "markdown" in content_type:
        return "remote.md"
    if "text/plain" in content_type:
        return "remote.txt"
    return "remote.html"


def _build_requests_kwargs(proxy_url: str = None, timeout_seconds: float = 25.0) -> dict:
    kwargs = {"timeout": timeout_seconds, "verify": False}
    if proxy_url and proxy_url.strip():
        kwargs["proxies"] = {
            "http": proxy_url.strip(),
            "https": proxy_url.strip(),
        }
    return kwargs


def _build_article_headers(url: str) -> dict:
    parsed = urlparse(str(url or "").strip())
    host = (parsed.netloc or "").lower()
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/136.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
    }
    if "mp.weixin.qq.com" in host:
        headers.update({
            "User-Agent": "Mozilla/5.0 (iPhone; CPU iPhone OS 16_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Mobile/15E148 MicroMessenger/8.0.40",
            "Referer": "https://mp.weixin.qq.com/",
        })
    return headers


def _looks_like_access_block_page(html: str, url: str = "") -> bool:
    sample = str(html or "")[:4000]
    host = (urlparse(str(url or "").strip()).netloc or "").lower()
    keywords = [
        "环境异常",
        "去验证",
        "请在微信客户端打开链接",
        "访问过于频繁",
        "当前访问环境异常",
        "完成验证",
    ]
    if any(keyword in sample for keyword in keywords):
        return True
    if "mp.weixin.qq.com" in host and "#js_content" not in sample and "js_content" not in sample:
        return True
    return False


def _extract_wechat_article_text(html: str, url: str) -> dict:
    try:
        from bs4 import BeautifulSoup
    except ImportError as e:
        raise RuntimeError("未安装 beautifulsoup4，无法解析微信公众号文章。") from e

    soup = BeautifulSoup(html, "html.parser")
    title_node = soup.select_one("#activity-name") or soup.select_one("h1")
    meta_title = soup.select_one('meta[property="og:title"]')
    title = title_node.get_text(" ", strip=True) if title_node else ""
    if not title and meta_title:
        title = str(meta_title.get("content") or "").strip()

    content_node = soup.select_one("#js_content")
    if not content_node:
        raise RuntimeError("未能从微信公众号页面中定位正文区域。")

    for tag in content_node.select("script, style, noscript"):
        tag.decompose()
    raw_text = content_node.get_text("\n", strip=True)
    clean_text = clean_document_text(raw_text)
    if not clean_text or len(clean_text) < 80:
        raise RuntimeError("微信公众号正文提取结果过短，疑似仍为异常页。")
    if title and title not in clean_text[:150]:
        raw_text = f"# {title}\n\n{raw_text}".strip()
        clean_text = clean_document_text(raw_text)

    return {
        "file_name": title or "wechat_article.html",
        "file_type": "html",
        "raw_text": raw_text,
        "clean_text": clean_text,
        "char_count": len(clean_text),
        "page_count": None,
        "section_count": None,
        "source_url": url,
        "ocr_used": False,
    }


def extract_web_article_text(url: str, proxy_url: str = None) -> dict:
    try:
        from bs4 import BeautifulSoup
    except ImportError as e:
        raise RuntimeError("未安装 beautifulsoup4，无法解析网页文章。") from e

    headers = _build_article_headers(url)
    resp = requests.get(url, headers=headers, **_build_requests_kwargs(proxy_url, timeout_seconds=30.0))
    resp.raise_for_status()
    html = resp.text
    if "mp.weixin.qq.com" in (urlparse(url).netloc or "").lower():
        return _extract_wechat_article_text(html, url)
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
        tag.decompose()

    title = ""
    if soup.title and soup.title.string:
        title = soup.title.string.strip()
    main = soup.find("article") or soup.find("main") or soup.body or soup
    blocks = []
    for node in main.find_all(["h1", "h2", "h3", "p", "li"], limit=500):
        text = node.get_text(" ", strip=True)
        if not text or len(text) < 2:
            continue
        if node.name in {"h1", "h2", "h3"}:
            level = {"h1": "#", "h2": "##", "h3": "###"}.get(node.name, "##")
            blocks.append(f"{level} {text}")
        else:
            blocks.append(text)

    raw_text = "\n\n".join(blocks).strip()
    if title and title not in raw_text[:200]:
        raw_text = f"# {title}\n\n{raw_text}".strip()
    clean_text = clean_document_text(raw_text)
    if not clean_text or _looks_like_access_block_page(html, url):
        raise RuntimeError("网页未提取到足够正文内容，请尝试更换文章链接。")
    return {
        "file_name": title or _guess_remote_file_name(url, resp) or "remote_article.html",
        "file_type": "html",
        "raw_text": raw_text,
        "clean_text": clean_text,
        "char_count": len(clean_text),
        "page_count": None,
        "section_count": None,
        "source_url": url,
        "ocr_used": False,
    }


def extract_document_from_url(url: str, proxy_url: str = None) -> dict:
    url = str(url or "").strip()
    if not url.lower().startswith(("http://", "https://")):
        raise RuntimeError("请输入完整的在线链接（以 http:// 或 https:// 开头）。")

    headers = _build_article_headers(url)
    resp = requests.get(url, stream=True, headers=headers, **_build_requests_kwargs(proxy_url, timeout_seconds=35.0))
    resp.raise_for_status()
    content = resp.content
    file_name = _guess_remote_file_name(url, resp)
    suffix = os.path.splitext(file_name)[1].lower()
    content_type = str(resp.headers.get("Content-Type", "") or "").lower()

    if not suffix:
        if "pdf" in content_type:
            file_name = "remote.pdf"
        elif "presentation" in content_type or "pptx" in content_type:
            file_name = "remote.pptx"
        elif "wordprocessingml" in content_type or "docx" in content_type:
            file_name = "remote.docx"
        elif "markdown" in content_type:
            file_name = "remote.md"
        elif "text/plain" in content_type:
            file_name = "remote.txt"
        else:
            file_name = "remote_article.html"
        suffix = os.path.splitext(file_name)[1].lower()

    if suffix in SUPPORTED_DOCUMENT_EXTENSIONS:
        ok, err = validate_document_upload(file_name, len(content))
        if not ok:
            raise RuntimeError(err)
        result = extract_document_text(content, file_name)
        result["source_url"] = url
        return result

    return extract_web_article_text(url, proxy_url=proxy_url)


def extract_document_text(file_bytes: bytes, file_name: str) -> dict:
    suffix = os.path.splitext(str(file_name or "").strip())[1].lower()
    if suffix == ".pdf":
        extracted = _extract_pdf_text(file_bytes)
    elif suffix == ".docx":
        extracted = _extract_docx_text(file_bytes)
    elif suffix == ".pptx":
        extracted = _extract_pptx_text(file_bytes)
    else:
        extracted = _extract_plain_text(file_bytes)

    raw_text = str(extracted.get("raw_text") or "").strip()
    clean_text = clean_document_text(raw_text)
    if not clean_text:
        raise RuntimeError("文档未提取到可用文本，请确认文档不是纯扫描图片，或先转换为可复制文本。")

    return {
        "file_name": file_name,
        "file_type": suffix.lstrip("."),
        "raw_text": raw_text,
        "clean_text": clean_text,
        "char_count": len(clean_text),
        "page_count": extracted.get("page_count"),
        "section_count": extracted.get("section_count"),
        "ocr_used": bool(extracted.get("ocr_used", False)),
    }


def split_document_chunks(text: str, max_chars: int = DEFAULT_DOCUMENT_CHUNK_CHARS) -> list[dict]:
    text = clean_document_text(text)
    if not text:
        return []

    paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]
    if not paragraphs:
        paragraphs = [text]

    chunks = []
    current_parts = []
    current_len = 0

    def flush_current():
        nonlocal current_parts, current_len
        if not current_parts:
            return
        chunk_text = "\n\n".join(current_parts).strip()
        chunks.append({
            "chunk_index": len(chunks) + 1,
            "text": chunk_text,
            "char_count": len(chunk_text),
        })
        current_parts = []
        current_len = 0

    for para in paragraphs:
        if len(para) > max_chars:
            flush_current()
            start = 0
            while start < len(para):
                end = min(start + max_chars, len(para))
                split_at = para.rfind("。", start, end)
                if split_at <= start + int(max_chars * 0.6):
                    split_at = para.rfind("，", start, end)
                if split_at <= start + int(max_chars * 0.5):
                    split_at = end
                else:
                    split_at += 1
                piece = para[start:split_at].strip()
                if piece:
                    chunks.append({
                        "chunk_index": len(chunks) + 1,
                        "text": piece,
                        "char_count": len(piece),
                    })
                start = split_at
            continue

        para_len = len(para) + 2
        if current_parts and current_len + para_len > max_chars:
            flush_current()
        current_parts.append(para)
        current_len += para_len

    flush_current()
    return chunks


def _build_openai_client(api_key: str, base_url: str, proxy_url: str = None):
    if not api_key:
        raise RuntimeError("请填写 API Key 以启用总结功能。")
    try:
        from openai import OpenAI
        import httpx
    except ImportError as e:
        raise RuntimeError("未安装 openai 依赖，无法进行总结。") from e

    client_kwargs = {"api_key": api_key}
    if base_url and base_url.strip():
        client_kwargs["base_url"] = base_url.strip()

    httpx_kwargs = {"verify": False}
    if proxy_url and proxy_url.strip():
        httpx_kwargs["proxy"] = proxy_url.strip()
    client_kwargs["http_client"] = httpx.Client(**httpx_kwargs)
    return OpenAI(**client_kwargs)


def _extract_completion_content(raw_resp) -> str:
    if isinstance(raw_resp, dict):
        choices = raw_resp.get("choices", [])
        if choices:
            return str(choices[0].get("message", {}).get("content", "") or "")
        return ""
    if hasattr(raw_resp, "choices") and raw_resp.choices:
        return str(raw_resp.choices[0].message.content or "")
    if isinstance(raw_resp, str):
        return raw_resp
    return ""


def _extract_summary_markdown(raw_text: str) -> str:
    payload = _parse_summary_payload(raw_text)
    if isinstance(payload, dict):
        summary_text = str(
            payload.get("summary_markdown")
            or payload.get("summary")
            or payload.get("summary_md")
            or ""
        ).strip()
        if summary_text:
            return summary_text
    return str(raw_text or "").strip()


def _summarize_document_passage(client, model: str, system_prompt: str, user_prompt: str, max_tokens: int) -> str:
    response = client.chat.completions.create(
        model=model.strip() or "gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.2,
        max_tokens=max_tokens,
        response_format={"type": "json_object"},
    )
    content = _extract_completion_content(response)
    if not content:
        raise RuntimeError("模型未返回有效总结内容。")
    return _extract_summary_markdown(content)


def summarize_document_text(
    text: str,
    api_key: str,
    base_url: str,
    model: str,
    proxy_url: str = None,
    progress_callback=None,
    ui_locale: str | None = None,
) -> dict:
    cleaned_text = clean_document_text(text)
    if not cleaned_text:
        raise RuntimeError("文档文本为空，无法总结。")

    client = _build_openai_client(api_key, base_url, proxy_url)
    content_len = len(cleaned_text)
    direct_system_prompt = _build_document_summary_system_prompt(ui_locale)

    if content_len <= DEFAULT_DOCUMENT_DIRECT_SUMMARY_CHARS:
        if callable(progress_callback):
            progress_callback(
                35,
                "Generating direct summary for the short document..."
                if _is_english_output_locale(ui_locale)
                else "文档较短，正在直接生成总结...",
            )
        prompt = (
            (
                "Please summarize the following document.\n"
                "Requirements:\n"
                "- Return JSON only\n"
                "- Include only the field `summary_markdown`\n"
                "- Output 6-10 bullets in `## Main Points`\n"
                "- Include `## Conclusion (Optional)` only when the document clearly presents a conclusion\n\n"
                f"Document body:\n{cleaned_text}"
            )
            if _is_english_output_locale(ui_locale)
            else (
                "请总结以下文档内容。\n"
                "要求：\n"
                "- 只返回 JSON\n"
                "- 仅包含字段 `summary_markdown`\n"
                "- `## 主要内容` 中输出 6-10 条要点\n"
                "- `## 结论（可选）` 只有在文档确实有明确结论时才输出\n\n"
                f"文档正文：\n{cleaned_text}"
            )
        )
        summary_markdown = _summarize_document_passage(
            client,
            model,
            direct_system_prompt,
            prompt,
            max_tokens=2200,
        )
        if callable(progress_callback):
            progress_callback(100, "Document summary completed." if _is_english_output_locale(ui_locale) else "文档总结完成。")
        return {
            "summary_markdown": summary_markdown,
            "strategy": "direct",
            "chunk_count": 1,
            "char_count": content_len,
        }

    chunks = split_document_chunks(cleaned_text, max_chars=DEFAULT_DOCUMENT_CHUNK_CHARS)
    if not chunks:
        raise RuntimeError("文档分块失败，无法继续总结。")

    chunk_summaries = []
    total_chunks = len(chunks)
    for idx, chunk in enumerate(chunks, start=1):
        if callable(progress_callback):
            progress_callback(
                20 + int(50 * idx / max(1, total_chunks)),
                (
                    f"Summarizing chunk {idx}/{total_chunks}..."
                    if _is_english_output_locale(ui_locale)
                    else f"正在总结第 {idx}/{total_chunks} 块..."
                ),
            )
        chunk_prompt = (
            (
                f"Please summarize the document excerpt below (chunk {idx}/{total_chunks}).\n"
                "Requirements:\n"
                "- Return JSON only\n"
                "- Include only the field `summary_markdown`\n"
                "- Use the structure `### Chunk Highlights` and `### Chunk Key Details`\n"
                "- Prefer bullet points instead of long paragraphs\n\n"
                f"Document excerpt:\n{chunk['text']}"
            )
            if _is_english_output_locale(ui_locale)
            else (
                f"请总结下面这段文档片段（第 {idx}/{total_chunks} 块）。\n"
                "要求：\n"
                "- 只返回 JSON\n"
                "- 仅包含字段 `summary_markdown`\n"
                "- 输出结构：`### 本块要点`、`### 本块关键信息`\n"
                "- 使用分条要点，不要长段落\n\n"
                f"文档片段：\n{chunk['text']}"
            )
        )
        chunk_summary = _summarize_document_passage(
            client,
            model,
            (
                "You are a document passage summarization assistant. Always return valid JSON containing only `summary_markdown`."
                if _is_english_output_locale(ui_locale)
                else "你是一个文档片段总结助手。请始终返回合法 JSON，且只能包含 `summary_markdown`。"
            ),
            chunk_prompt,
            max_tokens=1000,
        )
        chunk_summaries.append(
            (
                f"## Chunk {idx} Summary\n{chunk_summary}"
                if _is_english_output_locale(ui_locale)
                else f"## 第{idx}块摘要\n{chunk_summary}"
            )
        )

    if callable(progress_callback):
        progress_callback(
            80,
            "Merging chunk summaries into the final document summary..."
            if _is_english_output_locale(ui_locale)
            else "正在汇总整份文档摘要...",
        )
    merged_chunk_text = "\n\n".join(chunk_summaries)
    merge_prompt = (
        (
            "Below are chunk summaries from the same long document. Generate the final summary from them.\n"
            "Requirements:\n"
            "- Return JSON only\n"
            "- Include only the field `summary_markdown`\n"
            "- Use Markdown\n"
            "- Use the fixed structure: `## Core Topic`, `## Main Points`, `## Key Details`, `## Conclusion (Optional)`\n"
            "- Output 8-12 bullets in `## Main Points`\n"
            "- Avoid repeating the same meaning across chunks\n\n"
            f"Chunk summaries:\n{merged_chunk_text}"
        )
        if _is_english_output_locale(ui_locale)
        else (
            "下面是同一份长文档的分块摘要，请基于这些分块摘要生成最终总结。\n"
            "要求：\n"
            "- 只返回 JSON\n"
            "- 仅包含字段 `summary_markdown`\n"
            "- 使用 Markdown\n"
            "- 固定结构：`## 核心主题`、`## 主要内容`、`## 关键信息`、`## 结论（可选）`\n"
            "- `## 主要内容` 输出 8-12 条，按条列出\n"
            "- 不要重复每一块的相同意思\n\n"
            f"分块摘要：\n{merged_chunk_text}"
        )
    )
    summary_markdown = _summarize_document_passage(
        client,
        model,
        direct_system_prompt,
        merge_prompt,
        max_tokens=2500,
    )
    if callable(progress_callback):
        progress_callback(100, "Long document summary completed." if _is_english_output_locale(ui_locale) else "长文档总结完成。")
    return {
        "summary_markdown": summary_markdown,
        "strategy": "chunked",
        "chunk_count": total_chunks,
        "char_count": content_len,
    }


def _extract_claim_items(raw_text: str, max_claims: int) -> list[dict]:
    payload = _parse_summary_payload(raw_text)
    if isinstance(payload, dict):
        claim_list = payload.get("claims") or payload.get("items") or payload.get("key_claims") or []
    else:
        candidate = _extract_json_candidate(raw_text)
        try:
            parsed = json.loads(candidate)
        except Exception:
            parsed = None
        if isinstance(parsed, dict):
            claim_list = parsed.get("claims") or parsed.get("items") or parsed.get("key_claims") or []
        else:
            claim_list = []

    normalized = []
    for item in claim_list:
        if isinstance(item, str):
            claim = item.strip()
            queries = [claim] if claim else []
        elif isinstance(item, dict):
            claim = str(item.get("claim") or item.get("statement") or item.get("text") or "").strip()
            queries = item.get("queries") or item.get("search_queries") or []
            if isinstance(queries, str):
                queries = [queries]
            queries = [str(q).strip() for q in queries if str(q).strip()]
        else:
            continue
        claim = _normalize_fact_check_claim(claim)
        if claim:
            normalized.append({"claim": claim, "queries": queries[:2]})
        if len(normalized) >= max_claims:
            break
    return normalized


def _normalize_fact_check_claim(claim: str) -> str:
    text = re.sub(r"\s+", " ", str(claim or "").strip()).strip()
    if not text:
        return ""

    text = re.sub(r"^[\-*•\d\.\)\(（）：:、\s]+", "", text).strip()
    text = re.sub(r"[：:]\s*$", "", text).strip()
    text = _strip_fact_check_context_prefixes(text)

    trailing_phrases = [
        "其中，", "其中,", "但是，", "但是,", "不过，", "不过,", "并且，", "并且,",
        "而且，", "而且,", "此外，", "此外,", "另外，", "另外,", "例如，", "例如,",
        "比如，", "比如,", "随后，", "随后,", "然后，", "然后,", "同时，", "同时,",
    ]
    for phrase in trailing_phrases:
        if text.endswith(phrase):
            text = text[: -len(phrase)].rstrip(" ，,;；:：")
            break

    # 如果最后一个分句明显以承接词开头，说明模型大概率只截到半句，直接去掉尾巴。
    segments = re.split(r"([。！？；;])", text)
    merged = "".join(segments).strip()
    tail_match = re.search(
        r"(?:，|,)\s*(其中|但是|不过|并且|而且|此外|另外|例如|比如|随后|然后|同时)\s*$",
        merged,
    )
    if tail_match:
        merged = merged[:tail_match.start()].rstrip(" ，,;；:：")

    merged = re.sub(r"\s+", " ", merged).strip()
    if not merged:
        return ""
    if len(merged) < 8:
        return ""
    return merged


FACT_CHECK_COMMENTARY_PREFIX_PATTERN = re.compile(r"^(?:视频观点|观点|评论|分析|解读)\s*[：:]", re.I)
FACT_CHECK_OPINION_MARKERS = [
    "意在",
    "实质是",
    "终极盾牌",
    "被动局面",
    "救银行",
    "刺激经济",
    "逆向筛选",
    "难以接受",
    "谈判周期不匹配",
    "关键是",
]


def _is_likely_commentary_claim(raw_text: str) -> bool:
    value = re.sub(r"\s+", " ", str(raw_text or "")).strip()
    if not value:
        return False
    has_commentary_prefix = bool(FACT_CHECK_COMMENTARY_PREFIX_PATTERN.match(value))
    has_opinion_markers = any(marker in value for marker in FACT_CHECK_OPINION_MARKERS)
    has_hard_fact = bool(
        re.search(r"\b20\d{2}\b", value)
        or re.search(r"\d+(?:\.\d+)?\s*%", value)
        or re.search(r"\d+(?:,\d{3})+", value)
        or re.search(r"\d+(?:\.\d+)?\s*(?:亿|万亿|万人|亿美元|港元|导弹)\b", value)
        or re.search(r"(宣布|通报|处罚|要求|表示|报道称|报道|官方|央行|政府|reuters|bloomberg|wsj|ft)", value, re.I)
    )
    if has_commentary_prefix and not has_hard_fact:
        return True
    if has_commentary_prefix and has_opinion_markers and not re.search(r"(报道|官方|宣布|处罚|要求|央行|政府)", value, re.I):
        return True
    return has_opinion_markers and not has_hard_fact


def _extract_claims_from_summary_heuristic(summary_markdown: str, max_claims: int) -> list[dict]:
    summary_text = clean_document_text(summary_markdown or "")
    if not summary_text or max_claims <= 0:
        return []

    candidates: list[tuple[int, str]] = []
    lines = [line.strip() for line in str(summary_markdown or "").splitlines() if line.strip()]
    stop_markers = {"核心主题", "主要内容", "关键信息", "结论", "本块要点", "本块关键信息"}
    news_markers = [
        "表示", "称", "宣布", "通报", "回应", "指出", "发布", "发生", "遇袭", "逮捕",
        "选举", "关税", "制裁", "协议", "政策", "政府", "官方", "记者", "报道",
        "访问", "调解", "谈判", "整治", "承揽", "上市", "边控", "员工", "蓝筹股",
        "涉及", "针对", "进行", "开展", "部署", "要求", "呼吁", "强调", "指出",
    ]
    for raw_line in lines:
        raw_compact = re.sub(r"^#+\s*", "", raw_line).strip()
        raw_compact = re.sub(r"^[-*•]\s*", "", raw_compact).strip()
        if _is_likely_commentary_claim(raw_compact):
            continue
        cleaned_line = re.sub(r"^#+\s*", "", raw_line).strip()
        cleaned_line = re.sub(r"^[-*•]\s*", "", cleaned_line).strip()
        cleaned_line = _normalize_fact_check_claim(cleaned_line)
        if not cleaned_line or cleaned_line in stop_markers:
            continue
        if len(cleaned_line) > 140:
            continue
        score = 0
        if re.search(r"\d", cleaned_line):
            score += 2
        if any(marker in cleaned_line for marker in news_markers):
            score += 2
        if len(cleaned_line) >= 18:
            score += 1
        if score > 0:
            candidates.append((score, cleaned_line))

    seen: set[str] = set()
    results: list[dict] = []
    for _score, claim in sorted(candidates, key=lambda item: (-item[0], -len(item[1]))):
        lowered = claim.lower()
        if lowered in seen:
            continue
        seen.add(lowered)
        results.append({"claim": claim, "queries": [claim]})
        if len(results) >= max_claims:
            break
    return results


def _extract_claims_from_text_heuristic(text: str, max_claims: int) -> list[dict]:
    cleaned = clean_document_text(text or "")
    if not cleaned or max_claims <= 0:
        return []

    fragments = re.split(r"[。！？\n]+", cleaned)
    keywords = [
        "表示", "称", "宣布", "通报", "回应", "政府", "官方", "记者", "报道", "数据显示",
        "according to", "official", "访问", "调解", "谈判", "整治", "承揽", "上市", "边控", "涉及",
    ]
    scored: list[tuple[int, str]] = []
    for fragment in fragments:
        if _is_likely_commentary_claim(fragment):
            continue
        claim = _normalize_fact_check_claim(fragment)
        if not claim or len(claim) < 16 or len(claim) > 150:
            continue
        score = 0
        if re.search(r"\d", claim):
            score += 2
        if any(keyword.lower() in claim.lower() for keyword in keywords):
            score += 2
        if score > 0:
            scored.append((score, claim))

    seen: set[str] = set()
    results: list[dict] = []
    for _score, claim in sorted(scored, key=lambda item: (-item[0], -len(item[1]))):
        lowered = claim.lower()
        if lowered in seen:
            continue
        seen.add(lowered)
        results.append({"claim": claim, "queries": [claim]})
        if len(results) >= max_claims:
            break
    return results


def _build_heuristic_claim_items(text: str, summary_markdown: str, max_claims: int) -> list[dict]:
    claims = _extract_claims_from_summary_heuristic(summary_markdown, max_claims=max_claims)
    if len(claims) >= max_claims:
        return claims[:max_claims]

    supplemental = _extract_claims_from_text_heuristic(text, max_claims=max_claims)
    seen = {str(item.get("claim") or "").strip().lower() for item in claims}
    for item in supplemental:
        claim = str(item.get("claim") or "").strip()
        lowered = claim.lower()
        if not claim or lowered in seen:
            continue
        seen.add(lowered)
        claims.append(item)
        if len(claims) >= max_claims:
            break
    return claims[:max_claims]


def _supplement_claim_items(
    primary_claims: list[dict] | None,
    supplemental_claims: list[dict] | None,
    *,
    max_claims: int,
) -> list[dict]:
    merged: list[dict] = []
    seen: set[str] = set()
    for bucket in (primary_claims or [], supplemental_claims or []):
        for item in bucket:
            claim = _normalize_fact_check_claim(str(item.get("claim") or ""))
            if not claim:
                continue
            lowered = claim.lower()
            if lowered in seen:
                continue
            seen.add(lowered)
            queries = item.get("queries") or []
            if isinstance(queries, str):
                queries = [queries]
            normalized_queries = [str(query).strip() for query in queries if str(query).strip()]
            merged.append({"claim": claim, "queries": normalized_queries[:2]})
            if len(merged) >= max_claims:
                return merged
    return merged


def _looks_like_siliconflow_base_url(base_url_value: str) -> bool:
    lowered = str(base_url_value or "").strip().lower()
    return "siliconflow" in lowered


def _build_summary_model_candidates(model: str, base_url: str) -> list[str]:
    preferred = str(model or "").strip() or "gpt-3.5-turbo"
    candidates: list[str] = [preferred]
    if _looks_like_siliconflow_base_url(base_url):
        fallbacks = [
            "Pro/deepseek-ai/DeepSeek-V4-Flash",
            "deepseek-ai/DeepSeek-V4-Flash",
            "Pro/MiniMaxAI/MiniMax-M2.5",
            "MiniMaxAI/MiniMax-M2.5",
            "deepseek-ai/DeepSeek-V3.2",
        ]
        for fallback in fallbacks:
            if fallback and fallback not in candidates:
                candidates.append(fallback)
    return candidates


def extract_key_claims(
    text: str,
    summary_markdown: str,
    api_key: str,
    base_url: str,
    model: str,
    proxy_url: str = None,
    max_claims: int = 8,
) -> list[dict]:
    client = _build_openai_client(api_key, base_url, proxy_url)
    cleaned_text = clean_document_text(text)
    excerpt = _build_fact_check_excerpt(cleaned_text, max_chars=7000)
    summary_excerpt = clean_document_text(summary_markdown or "")[:2200]
    heuristic_claims = _build_heuristic_claim_items(cleaned_text, summary_excerpt, max_claims=max_claims)
    prompt = (
        f"请从下面文档中提取最值得做事实核查的 {max_claims} 条关键声明。\n"
        "要求：\n"
        "- 只提取可核查的硬信息，不要提取纯观点、情绪表达或修辞。\n"
        "- 优先提取：数字、时间、政策、官方表述、事件是否发生、人物公开言论、强因果判断。\n"
        "- 如果原文里存在多条不同主体、不同数字、不同时间点的声明，不要合并成一条笼统表述。\n"
        "- 如果文本里可核查点较多，请尽量提满，不要只返回 1-2 条过于宽泛的声明。\n"
        "- 每条 claim 必须是完整、自洽、可独立理解的一句话，不要只截半句。\n"
        "- 不要让 claim 以“其中、但是、并且、而且、例如、比如、随后、然后、同时”等承接词结尾。\n"
        "- 如果原文句子过长，请改写成完整但更短的独立陈述，保留主体、动作和关键数字/时间。\n"
        "- 每条 queries 最多给 2 个：1 个尽量精确，1 个尽量放宽。\n"
        "- 如果声明涉及外国人物、政府部门、城市、组织或国际事件，至少有 1 个 query 使用英文或原文名称，不要只给中文译名。\n"
        "- 如果声明里有日期，queries 不要都死扣完整日期；至少保留 1 个去掉具体日期的宽松版本，避免搜不到。\n"
        "- 只返回 JSON，对象格式如下：\n"
        "{\n"
        '  "claims": [\n'
        '    {"claim": "声明内容", "queries": ["搜索词1", "搜索词2"]}\n'
        "  ]\n"
        "}\n"
        "- 搜索词中尽量包含主体、时间、地点、数字、机构。\n\n"
        f"文档总结：\n{summary_excerpt}\n\n"
        f"文档正文节选：\n{excerpt}"
    )
    try:
        response = client.chat.completions.create(
            model=model.strip() or "gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "你是一个事实核查分析助手。请只返回合法 JSON。"},
                {"role": "user", "content": prompt},
            ],
            response_format={"type": "json_object"},
            max_tokens=1000,
            temperature=0.1,
        )
        content = _extract_completion_content(response)
        claims = _extract_claim_items(content, max_claims=max_claims)
    except Exception:
        claims = []
    claims = _supplement_claim_items(claims, heuristic_claims, max_claims=max_claims)
    if not claims:
        raise RuntimeError("未能从文档中提取到可核查的关键声明。")
    return claims


def classify_document_for_fact_check(
    text: str,
    summary_markdown: str,
    api_key: str,
    base_url: str,
    model: str,
    proxy_url: str = None,
) -> dict:
    client = _build_openai_client(api_key, base_url, proxy_url)
    cleaned_text = clean_document_text(text)
    excerpt = cleaned_text[:8000]
    prompt = (
        "请判断下面文档是否属于需要进行事实核查的类型。\n"
        "需要事实核查的典型类型包括：新闻、研究报告、时评/评论、政策解读、行业分析、带现实事件与数据结论的文章。\n"
        "通常不需要事实核查的类型包括：小说、散文、内部会议纪要、教程、产品文档、纯技术方案、个人笔记、合同草稿。\n"
        "请只返回 JSON，格式如下：\n"
        "{\n"
        '  "document_type": "news|research|commentary|policy_analysis|industry_analysis|meeting_notes|tutorial|technical_doc|fiction|personal_notes|other",\n'
        '  "should_fact_check": true,\n'
        '  "reason": "简短原因",\n'
        '  "recommended_claim_count": 5\n'
        "}\n"
        "- 如果文档明显属于新闻、研究、时评、政策解读、行业分析，should_fact_check 设为 true。\n"
        "- 否则设为 false。\n"
        "- recommended_claim_count 请根据文档信息密度给出 3-12 之间的整数。\n\n"
        f"文档总结：\n{summary_markdown[:3500]}\n\n"
        f"文档正文节选：\n{excerpt}"
    )
    response = client.chat.completions.create(
        model=model.strip() or "gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "你是一个文档类型判定助手。请只返回合法 JSON。"},
            {"role": "user", "content": prompt},
        ],
        response_format={"type": "json_object"},
        max_tokens=500,
        temperature=0.1,
    )
    content = _extract_completion_content(response)
    payload = _parse_summary_payload(content)
    if not isinstance(payload, dict):
        try:
            payload = json.loads(_extract_json_candidate(content))
        except Exception:
            payload = {}

    doc_type = str(payload.get("document_type") or "other").strip() or "other"
    should_fact_check = bool(payload.get("should_fact_check", False))
    reason = str(payload.get("reason") or "").strip()
    recommended_claim_count = int(payload.get("recommended_claim_count") or 5)
    if not (3 <= recommended_claim_count <= 12):
        recommended_claim_count = 5

    if not reason:
        lowered = (summary_markdown + "\n" + excerpt).lower()
        if any(k in lowered for k in ["研究", "报告", "统计", "数据", "survey", "research", "report"]):
            doc_type = "research" if doc_type == "other" else doc_type
            should_fact_check = True
            reason = "文档包含研究/数据结论，适合做关键声明核查。"
        elif any(k in lowered for k in ["评论", "时评", "社论", "news", "记者", "新华社", "路透", "政策"]):
            doc_type = "commentary" if doc_type == "other" else doc_type
            should_fact_check = True
            reason = "文档涉及现实事件、新闻或政策判断，适合做关键声明核查。"
        else:
            reason = "文档更偏说明、教程、纪要或普通资料，可直接总结，无需默认核查。"

    return {
        "document_type": doc_type,
        "should_fact_check": should_fact_check,
        "reason": reason,
        "recommended_claim_count": recommended_claim_count,
    }


def _build_fact_check_excerpt(text: str, max_chars: int = 7000) -> str:
    cleaned = clean_document_text(text or "")
    if len(cleaned) <= max_chars:
        return cleaned
    head_chars = max(2500, int(max_chars * 0.72))
    tail_chars = max(800, max_chars - head_chars - 40)
    return (
        cleaned[:head_chars].rstrip()
        + "\n...(中间内容省略)...\n"
        + cleaned[-tail_chars:].lstrip()
    )


def decide_video_fact_check_plan(text: str, summary_markdown: str) -> dict:
    content = str(text or "").strip()
    summary = clean_document_text(summary_markdown or "")
    skip_fact_check, skip_reason = _should_skip_fact_check_for_video_text(content)
    if skip_fact_check:
        return {
            "should_fact_check": False,
            "reason": f"当前视频更像测试/演示/技术链路内容，已跳过新闻核查（{skip_reason}）。",
            "recommended_claim_count": 0,
        }

    excerpt = _build_fact_check_excerpt(content, max_chars=5000)
    combined = f"{summary}\n{excerpt}".lower()
    negative_markers = [
        "教程", "教学", "代码", "编程", "开发", "安装", "评测", "开箱",
        "软件使用", "产品演示", "课程", "训练营", "实操", "案例讲解",
    ]
    news_markers = [
        "新闻", "报道", "记者", "新华社", "央视", "路透", "彭博", "法新社",
        "government", "official", "breaking", "reported", "according to",
    ]
    event_markers = [
        "发布", "宣布", "通报", "回应", "发生", "遇袭", "逮捕", "起火", "坠毁",
        "停火", "制裁", "投票", "选举", "协议", "关税", "经济数据", "失业率",
        "cpi", "gdp", "policy", "tariff", "sanction", "election",
    ]
    public_entity_markers = [
        "政府", "官方", "机构", "部门", "企业", "公司", "法院", "警方", "医院",
        "学校", "大学", "央行", "外交部", "商务部", "联合国", "白宫", "总统",
        "agency", "department", "ministry", "company", "corporation", "court",
        "police", "hospital", "university", "central bank", "white house",
        "president", "u.n.", "united nations",
    ]
    statement_markers = [
        "表示", "称", "指出", "披露", "证实", "否认", "承认", "公开表态", "时间线",
        "数据", "数字", "口径", "统计", "说法", "声明", "facts", "figures",
        "timeline", "statement", "claim", "claims", "said", "stated", "announced",
        "released", "confirmed", "denied",
    ]
    hard_fact_patterns = [
        r"\b20\d{2}\b",
        r"\b\d+(?:\.\d+)?\s*%",
        r"\b\d+(?:,\d{3})+\b",
        r"\b\d+(?:\.\d+)?\s*(?:亿|万亿|万人|万例|亿美元|亿元|人)\b",
    ]
    negative_hits = sum(1 for marker in negative_markers if marker in combined)
    news_hits = sum(1 for marker in news_markers if marker in combined)
    event_hits = sum(1 for marker in event_markers if marker in combined)
    public_entity_hits = sum(1 for marker in public_entity_markers if marker in combined)
    statement_hits = sum(1 for marker in statement_markers if marker in combined)
    hard_fact_hits = sum(1 for pattern in hard_fact_patterns if re.search(pattern, combined))
    content_len = len(content)
    summary_len = len(summary)

    if content_len < 180 and summary_len < 60:
        return {
            "should_fact_check": False,
            "reason": "当前视频可用于核查的文本过少，暂不启动新闻核查。",
            "recommended_claim_count": 0,
        }

    if negative_hits >= 3 and news_hits == 0 and event_hits == 0 and hard_fact_hits == 0 and content_len < 2500:
        return {
            "should_fact_check": False,
            "reason": "当前视频更像教程、产品演示或经验讲解，已默认跳过新闻核查。",
            "recommended_claim_count": 0,
        }

    if negative_hits >= 2 and news_hits == 0 and event_hits == 0:
        return {
            "should_fact_check": False,
            "reason": "当前视频更像观点表达、教程或日常内容，已跳过自动新闻核查。",
            "recommended_claim_count": 0,
        }

    clear_news_signal = news_hits >= 1 and event_hits >= 1
    clear_event_signal = event_hits >= 2 and hard_fact_hits >= 1
    clear_report_signal = news_hits >= 2 and hard_fact_hits >= 1
    clear_data_signal = news_hits >= 1 and hard_fact_hits >= 2 and content_len >= 1200
    clear_public_claim_signal = (
        content_len >= 900
        and hard_fact_hits >= 1
        and public_entity_hits >= 1
        and statement_hits >= 1
    )
    concise_public_claim_signal = (
        content_len >= 80
        and summary_len >= 100
        and hard_fact_hits >= 1
        and public_entity_hits >= 1
        and statement_hits >= 3
        and negative_hits == 0
    )
    clear_timeline_signal = (
        content_len >= 1400
        and hard_fact_hits >= 1
        and statement_hits >= 2
        and negative_hits == 0
    )
    should_fact_check = (
        clear_news_signal
        or clear_event_signal
        or clear_report_signal
        or clear_data_signal
        or clear_public_claim_signal
        or concise_public_claim_signal
        or clear_timeline_signal
    )

    if not should_fact_check:
        return {
            "should_fact_check": False,
            "reason": "当前视频未命中足够明确的新闻/事件信号，已跳过自动新闻核查。",
            "recommended_claim_count": 0,
        }

    signal_score = news_hits + event_hits + hard_fact_hits + public_entity_hits + statement_hits
    if content_len >= 12000 and signal_score >= 9:
        recommended_claim_count = 12
    elif (
        (content_len >= 8000 and signal_score >= 7)
        or (content_len >= 5000 and signal_score >= 8)
        or (content_len >= 3200 and clear_public_claim_signal and signal_score >= 7)
    ):
        recommended_claim_count = 8
    elif content_len >= 2500 and signal_score >= 3:
        recommended_claim_count = 5
    else:
        recommended_claim_count = 3

    reason = f"当前视频命中明显新闻/事件信号（得分 {signal_score}），将核查主要内容中的约 {recommended_claim_count} 条可核查声明。"
    return {
        "should_fact_check": True,
        "reason": reason,
        "recommended_claim_count": recommended_claim_count,
    }


def search_claim_sources(claim_items: list[dict], proxy_url: str = None, search_mode: str = "standard") -> list[dict]:
    items = [item for item in (claim_items or []) if str(item.get("claim") or "").strip()]
    if not items:
        return []
    if len(items) == 1:
        return [_search_single_claim_source(items[0], proxy_url=proxy_url, search_mode=search_mode)]

    from concurrent.futures import ThreadPoolExecutor

    max_workers = min(5, len(items))
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        return list(
            executor.map(
                _search_single_claim_source,
                items,
                [proxy_url] * len(items),
                [search_mode] * len(items),
            )
        )


def fact_check_document_claims(
    text: str,
    summary_markdown: str,
    api_key: str,
    base_url: str,
    model: str,
    proxy_url: str = None,
    max_claims: int = 8,
    progress_callback=None,
    ui_locale: str | None = None,
    search_mode: str = "standard",
) -> str:
    if max_claims <= 0:
        return ""

    summary_excerpt = clean_document_text(summary_markdown or "")[:2200]
    claims: list[dict] = []
    claim_sources: list[dict] = []
    if callable(progress_callback):
        progress_callback(5, _fact_check_text(ui_locale, "progress_extract_claims"))
    try:
        claims = extract_key_claims(
            text=text,
            summary_markdown=summary_excerpt,
            api_key=api_key,
            base_url=base_url,
            model=model,
            proxy_url=proxy_url,
            max_claims=max_claims,
        )
    except Exception as exc:
        print(f"Extract key claims failed: {exc}", flush=True)
        try:
            claims = _build_heuristic_claim_items(
                clean_document_text(text),
                clean_document_text(summary_markdown or "")[:2200],
                max_claims=max_claims,
            )
        except Exception:
            claims = []

    if not claims:
        return _build_fact_check_fallback_markdown(ui_locale=ui_locale)

    if callable(progress_callback):
        progress_callback(35, _fact_check_text(ui_locale, "progress_search_sources"))
    try:
        claim_sources = search_claim_sources(claims, proxy_url=proxy_url, search_mode=search_mode)
    except Exception as exc:
        print(f"Search claim sources failed: {exc}", flush=True)
        claim_sources = [
            {
                "claim": str(item.get("claim") or "").strip(),
                "queries": list(item.get("queries") or []),
                "search_markdown": "",
            }
            for item in claims
            if str(item.get("claim") or "").strip()
        ]

    compiled_sections = []
    candidate_claim_label = _fact_check_text(ui_locale, "candidate_claim")
    claim_label = _fact_check_text(ui_locale, "claim_label")
    search_terms_label = _fact_check_text(ui_locale, "search_terms_label")
    search_results_label = _fact_check_text(ui_locale, "search_results_label")
    for idx, item in enumerate(claim_sources, start=1):
        search_excerpt = str(item.get("search_markdown") or "").strip()[:2200]
        compiled_sections.append(
            f"### {candidate_claim_label}{idx}\n"
            f"- {claim_label}: {item['claim']}\n"
            f"- {search_terms_label}: {' | '.join(item.get('queries') or [])}\n"
            f"- {search_results_label}:\n{search_excerpt}"
        )

    if callable(progress_callback):
        progress_callback(70, _fact_check_text(ui_locale, "progress_generate_fact_check"))

    try:
        client = _build_openai_client(api_key, base_url, proxy_url)
        compiled_claim_text = "\n\n".join(compiled_sections)
        # 适当调大 max_tokens 以支持更多条目 (3-12条)
        if max_claims <= 1:
            final_max_tokens = 1000
        elif max_claims <= 3:
            final_max_tokens = 1600
        elif max_claims <= 6:
            final_max_tokens = 2200
        else:
            final_max_tokens = 3200
        prompt = (
            (
                "Based on the key claims and search results below, generate itemized source-discovery results in Markdown.\n"
                "Requirements:\n"
                f"- Return exactly {len(claim_sources)} items, one item for each candidate claim below, and keep the same order.\n"
                "- Use this structure for every item:\n"
                "### Item 1\n"
                "- Claim: ...\n"
                "- Source Status: Direct Source Located / Related Coverage Located / No Direct Source Yet\n"
                "- Source Notes:\n"
                "  - Located source pages: ...\n"
                "  - Why these pages are likely relevant to the claim: ...\n"
                "  - What still needs manual follow-up, if any: ...\n"
                "- Suggested Follow-ups: ...\n"
                "- Source Links: provide 2-4 external links in Markdown, such as [Reuters](https://...)\n"
                "- If the original text contains multiple distinct claims, split them into separate items instead of merging several numbers or events into one.\n"
                "- If a claim involves figures, dates, institutions, or rankings, explain as specifically as possible whether each element matches.\n"
                "- The task here is source discovery, not truth judgment. Focus on locating the most likely origin or report page for the claim.\n"
                "- Prefer article pages from major media, official statements, press releases, and outlet pages over generic topic pages, search pages, or videos.\n"
                "- Candidate web pages and site-specific official pages count as evidence candidates too; do not ignore them just because they are not from news RSS.\n"
                "- Your first job is to identify which matching web pages were found for each claim; even if the truth still needs follow-up, explain the matched reports, web pages, or official pages before discussing remaining gaps.\n"
                "- Start the rationale by summarizing the matched page types and source names instead of reusing a generic template across multiple items.\n"
                "- Do not reuse one-size-fits-all language such as exchange filings or original statistical methodology unless the collected candidate pages are actually about finance, exchanges, or statistics.\n"
                "- If the search context only says that no directly citable report was captured this round, phrase it as an automated retrieval miss rather than exaggerating it into `no mainstream or official coverage exists`.\n"
                "- When citing institutions, government departments, media outlets, or official sites in the rationale, prefer linking to the official or outlet page rather than naming it without a link.\n"
                "- Never treat the document itself as a source.\n"
                "- Even if search results are limited, list the candidate sources you do have instead of leaving the section empty.\n"
                "- Avoid truth-verdict language. If you found a likely source page, say so directly instead of judging whether the claim is true.\n"
                "- Return Markdown only, not JSON.\n\n"
                f"Document summary:\n{summary_excerpt}\n\n"
                f"Candidate claims and search results:\n{compiled_claim_text}"
            )
            if _is_english_output_locale(ui_locale)
            else (
                "请根据下面的关键声明与搜索结果，输出逐条来源导航 Markdown。\n"
                "要求：\n"
                f"- 必须严格返回 {len(claim_sources)} 条，下面每个候选声明都要对应 1 条，且顺序保持一致。\n"
                "- 每条都使用下面结构：\n"
                "### 条目1\n"
                "- 关键声明：...\n"
                "- 来源定位：已找到直接来源 / 已找到相关来源 / 暂未定位到直接来源\n"
                "- 来源线索说明：\n"
                "  - 已命中的来源网页：...\n"
                "  - 这些网页为什么像是这条新闻的出处或报道页：...\n"
                "  - 如果还有待人工确认的点，再补充说明：...\n"
                "- 建议继续查看：...\n"
                "- 来源出处：给出 2-4 个外部来源链接，格式如 [新华社](https://...)\n"
                "- 如果原文包含多条不同声明，请分别成条输出，不要把多个数字或多个事件揉成一条。\n"
                "- 如果一条声明涉及数字、时间、机构、排名，请尽量分别说明这些要素是否匹配。\n"
                "- 这一步的目标是定位来源网页，不是判断真伪。请优先回答“这条说法最像来自哪些网页”。\n"
                "- 优先选择大媒体报道页、官网声明页、新闻稿页、采访原文页，不要优先给专题页、搜索页、视频页或泛列表页。\n"
                "- 命中的候选网页、官网页面、定向站点页面同样属于证据候选，不要因为它们不是新闻 RSS 就忽略。\n"
                "- 你的首要任务是为每条声明指出当前已找到的对应网页来源；即使暂时无法判断真伪，也要先说明命中了哪些报道、普通网页或官网/机构页面。\n"
                "- `判断依据` 的第一句优先概括“已找到哪些网页类型与来源名”，不要把多条都写成同一套模板化理由。\n"
                "- 除非候选结果本身就是金融市场、交易所或统计口径主题，否则不要机械写“交易所公告 / 原始统计口径 / 一手权威来源”这类不贴题措辞。\n"
                "- 如果搜索上下文只写了“当前未抓到可直接引用的候选报道”，只能表述为“本轮自动检索暂未命中可直接引用来源”，不要夸大成“没有任何主流媒体或官方机构报道”。\n"
                "- 如果判断依据里提到具体机构、政府部门、媒体名或官网名，优先附上该机构/媒体的官网或栏目页链接，不要只写机构名称。\n"
                "- 禁止把文档本身当来源。\n"
                "- 如果搜索结果不足，也要写明目前查到的候选来源，不要空着。\n"
                "- 避免使用真假裁决式表述；只要找到了像出处的网页，就直接说明已找到哪些来源链接。\n"
                "- 只返回 Markdown，不要 JSON。\n\n"
                f"文档总结：\n{summary_excerpt}\n\n"
                f"候选声明与搜索结果：\n{compiled_claim_text}"
            )
        )
        response = client.chat.completions.create(
            model=model.strip() or "gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": _fact_check_text(ui_locale, "system_prompt")},
                {"role": "user", "content": prompt},
            ],
            max_tokens=final_max_tokens,
            temperature=0.15,
        )
        content = _extract_completion_content(response).strip()
    except Exception as exc:
        print(f"Generate fact check markdown failed: {exc}", flush=True)
        return _build_fact_check_fallback_markdown(claim_sources=claim_sources, ui_locale=ui_locale)

    has_any_candidate_pages = any(
        _extract_fact_check_source_links(str(item.get("search_markdown") or ""))
        for item in claim_sources
    )
    if (
        content
        and (
            (
                "Conclusion: Insufficient Evidence" in content
                and content.count("Conclusion: Insufficient Evidence") == max(1, content.count("### Item"))
            )
            or (
                "Source Status: No Direct Source Yet" in content
                and content.count("Source Status: No Direct Source Yet") == max(1, content.count("### Item"))
            )
            or (
                "核查结论：缺乏证据" in content
                and content.count("核查结论：缺乏证据") == max(1, content.count("### 条目"))
            )
            or (
                "来源定位：暂未定位到直接来源" in content
                and content.count("来源定位：暂未定位到直接来源") == max(1, content.count("### 条目"))
            )
            or (
                "Conclusion: Uncertain" in content
                and content.count("Conclusion: Uncertain") == max(1, content.count("### Item"))
                and any(token in content.lower() for token in [
                    "exchange filing",
                    "original statistical methodology",
                    "primary authoritative source",
                ])
            )
            or (
                "核查结论：存疑" in content
                and content.count("核查结论：存疑") == max(1, content.count("### 条目"))
                and any(token in content for token in [
                    "交易所公告",
                    "原始统计口径",
                    "一手权威来源",
                    "更适合判为“存疑”",
                ])
            )
            or (
                "Conclusion: Source Not Found Yet" in content
                and content.count("Conclusion: Source Not Found Yet") == max(1, content.count("### Item"))
            )
            or (
                "核查结论：暂未找到来源" in content
                and content.count("核查结论：暂未找到来源") == max(1, content.count("### 条目"))
            )
        )
        and has_any_candidate_pages
    ):
        try:
            retry_prompt = (
                (
                    "The previous draft was still too generic. Your first task is to identify which matching web pages were found for each claim before making any source-status summary.\n"
                    "For each item, first summarize the matched reports, general web pages, or official pages already found, and only then explain what still needs verification.\n"
                    "Do not reuse one-size-fits-all rationale such as exchange filings or original statistical methodology unless the collected pages are truly about that topic.\n"
                    "If candidate source pages are already present, do not label every item as `No Direct Source Yet`.\n"
                    "Keep the same Markdown structure and do not omit source links.\n\n"
                    f"Document summary:\n{summary_excerpt}\n\n"
                    f"Candidate claims and search results:\n{compiled_claim_text}"
                )
                if _is_english_output_locale(ui_locale)
                else (
                    "你上一版仍然过于模板化。你的首要任务不是先给终局判断，而是先为每条声明指出当前已找到的对应网页来源。\n"
                    "请逐条先概括已命中的报道、普通网页或官网/机构页面，再说明还待核对的细节；不要把多条都套成同一套“交易所公告 / 原始统计口径 / 一手权威来源”理由，除非搜索结果本身就是这个主题。\n"
                    "如果候选来源里已经有网页链接，不要把所有条目都写成“暂未定位到直接来源”。\n"
                    "请继续按原结构输出 Markdown，不要省略来源链接。\n\n"
                    f"文档总结：\n{summary_excerpt}\n\n"
                    f"候选声明与搜索结果：\n{compiled_claim_text}"
                )
            )
            retry_response = client.chat.completions.create(
                model=model.strip() or "gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": _fact_check_text(ui_locale, "retry_system_prompt")},
                    {"role": "user", "content": retry_prompt},
                ],
                max_tokens=final_max_tokens,
                temperature=0.1,
            )
            retry_content = _extract_completion_content(retry_response).strip()
            if retry_content:
                content = retry_content
        except Exception as retry_exc:
            print(f"Retry fact check markdown failed: {retry_exc}", flush=True)

    if callable(progress_callback):
        progress_callback(100, _fact_check_text(ui_locale, "progress_fact_check_done"))
    if not content:
        return _build_fact_check_fallback_markdown(claim_sources=claim_sources, ui_locale=ui_locale)
    final_markdown = _enrich_fact_check_items_with_claim_sources(content, claim_sources)
    final_markdown = _normalize_fact_check_conclusions_with_sources(final_markdown, claim_sources)
    final_markdown = _ensure_fact_check_item_coverage(final_markdown, claim_sources, ui_locale=ui_locale)
    if _is_placeholder_fact_check(final_markdown):
        return _build_fact_check_fallback_markdown(claim_sources=claim_sources, ui_locale=ui_locale)
    return final_markdown


def _should_skip_fact_check_for_video_text(text: str) -> tuple[bool, str]:
    """识别明显属于测试/产品说明/技术链路的字幕，避免误触发新闻事实核查。"""
    content = str(text or "").strip()
    if not content:
        return True, "empty_text"

    lowered = content.lower()
    strong_markers = [
        "模拟测试",
        "测试字幕",
        "测试文本",
        "回归验证",
        "链路验证",
        "产品主路径",
        "自动开始总结",
        "bridge payload",
        "payload_id",
        "ext_payload_id",
    ]
    if any(marker in content or marker in lowered for marker in strong_markers):
        return True, "strong_test_marker"

    paired_markers = [
        (("模拟", "测试"), "simulated_test_pair"),
        (("回归", "验证"), "regression_validation_pair"),
        (("技术", "链路"), "technical_pipeline_pair"),
    ]
    for keywords, reason in paired_markers:
        if all(keyword in content for keyword in keywords):
            return True, reason

    if "产品" in content and "验证" in content:
        product_validation_context_markers = [
            "主路径",
            "流程",
            "链路",
            "调试",
            "插件",
            "扩展",
            "bridge",
            "payload",
            "transcript",
            "api key",
            "render",
            "主站",
            "自动总结",
        ]
        context_hits = sum(1 for marker in product_validation_context_markers if marker in lowered or marker in content)
        if context_hits >= 2:
            return True, "product_validation_pair"

    # “内部”“测试”在新闻、采访、爆料类视频里也很常见，不能单独作为跳过核查的依据。
    internal_test_context_markers = [
        "bridge",
        "payload",
        "transcript",
        "api key",
        "插件",
        "扩展",
        "链路",
        "调试",
        "回归",
        "render",
        "主站",
    ]
    if "内部" in content and "测试" in content:
        context_hits = sum(1 for marker in internal_test_context_markers if marker in lowered or marker in content)
        if context_hits >= 2:
            return True, "internal_test_pair"

    if re.search(r"(这是一次|这是一段|本次|用于).{0,24}(模拟|测试|验证|调试|演示)", content):
        return True, "intro_test_pattern"

    technical_markers = [
        "chrome插件",
        "render主站",
        "render",
        "bridge",
        "payload",
        "transcript",
        "api key",
        "自动总结",
        "主链路",
        "产品技术流程",
        "技术流程",
    ]
    technical_hits = sum(1 for marker in technical_markers if marker in lowered or marker in content)
    if technical_hits >= 3 and any(word in content for word in ["测试", "验证", "调试", "回归"]):
        return True, "technical_test_context"

    return False, ""


def summarize_text(
    text: str,
    api_key: str,
    base_url: str,
    model: str,
    proxy_url: str = None,
    stream: bool = False,
    fact_check_model: str | None = None,
    enable_fact_check: bool = True,
    ui_locale: str | None = None,
):
    if not text or not text.strip():
        return "没有可总结的内容（文本为空）。"
    if not api_key:
        return "请填写 API Key 以启用总结功能。"

    if stream and str(fact_check_model or "").strip() and str(fact_check_model or "").strip() != str(model or "").strip():
        return "总结失败：当前双模型流水线不支持流式输出，请改用非流式模式。"

    try:
        client = _build_openai_client(api_key, base_url, proxy_url)
        summary_model = str(model or "").strip() or "gpt-3.5-turbo"
        fact_model = str(fact_check_model or "").strip() or summary_model

        content = text.strip()
        content_len = len(content)
        max_input_len = 16000 if content_len > 18000 else 20000
        if content_len > max_input_len:
            content = content[:max_input_len] + "\n...(内容过长已截断)..."

        from datetime import datetime
        current_date = datetime.now().strftime("%Y-%m-%d")

        fact_check_enabled = bool(enable_fact_check)
        if fact_check_enabled:
            skip_fact_check, skip_reason = _should_skip_fact_check_for_video_text(content)
            if skip_fact_check:
                fact_check_enabled = False
                print(f"SummarizeText: skip fact check, reason={skip_reason}", flush=True)
        else:
            print("SummarizeText: fact check disabled by caller", flush=True)

        prompt = _build_video_summary_prompt(content, current_date, ui_locale=ui_locale)
        max_tokens = 2600 if content_len < 12000 else 3000
        print(
            "SummarizeText: "
            f"content_len={content_len}, summary_model={summary_model}, fact_check_model={fact_model}, "
            f"fact_check_enabled={fact_check_enabled}"
        , flush=True)

        summary_model_candidates = _build_summary_model_candidates(summary_model, base_url)
        response = None
        content_str = ""
        normalized_payload = None
        summary_errors: list[str] = []
        active_summary_model = summary_model
        for candidate_model in summary_model_candidates:
            active_summary_model = candidate_model
            try:
                response = client.chat.completions.create(
                    model=candidate_model,
                    messages=[
                        {"role": "system", "content": _build_video_summary_system_prompt(ui_locale)},
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.2,
                    max_tokens=max_tokens,
                    response_format={"type": "json_object"},
                    stream=stream,
                )
                break
            except Exception as summary_exc:
                summary_errors.append(f"{candidate_model}: {summary_exc}")
                print(
                    f"SummarizeText: model retry candidate failed: {candidate_model} -> {summary_exc}",
                    flush=True,
                )
                response = None
        if response is None:
            raise RuntimeError(" | ".join(summary_errors) if summary_errors else "无法完成总结请求")
        if stream:
            return response

        if active_summary_model != summary_model:
            print(
                f"SummarizeText: fallback summary model applied -> {active_summary_model}",
                flush=True,
            )

        content_str = _extract_completion_content(response)
        if not content_str:
            return f"总结失败：无法解析响应内容。\n原始响应: {str(response)[:500]}"

        normalized_payload = _normalize_summary_payload(_parse_summary_payload(content_str), ui_locale=ui_locale)
        if not normalized_payload:
            try:
                repair_prompt = _build_video_summary_repair_prompt(content_str, ui_locale=ui_locale)
                repair_resp = client.chat.completions.create(
                    model=active_summary_model,
                    messages=[
                        {
                            "role": "system",
                            "content": (
                                "You are a JSON repair assistant. Return valid JSON only."
                                if _is_english_output_locale(ui_locale)
                                else "你是一个 JSON 修复助手，只返回合法 JSON。"
                            ),
                        },
                        {"role": "user", "content": repair_prompt},
                    ],
                    response_format={"type": "json_object"},
                    max_tokens=min(2200, max_tokens),
                    temperature=0.1,
                )
                repair_content = _extract_completion_content(repair_resp)
                normalized_payload = _normalize_summary_payload(_parse_summary_payload(repair_content), ui_locale=ui_locale)
            except Exception as repair_exc:
                print(f"Repair summary JSON failed: {repair_exc}", flush=True)

        if not normalized_payload:
            normalized_payload = {
                "summary_markdown": content_str,
                "fact_check_markdown": "",
            }

        summary_markdown = str(normalized_payload.get("summary_markdown") or "").strip() or content_str
        fact_check_markdown = ""
        if fact_check_enabled and len(content) > 100:
            try:
                max_claims = 8 if content_len >= 12000 else 6 if content_len >= 7000 else 5
                fact_check_markdown = fact_check_document_claims(
                    text=content,
                    summary_markdown=summary_markdown,
                    api_key=api_key,
                    base_url=base_url,
                    model=fact_model,
                    proxy_url=proxy_url,
                    max_claims=max_claims,
                    ui_locale=ui_locale,
                )
            except Exception as fact_exc:
                print(f"Video fact check pipeline failed: {fact_exc}", flush=True)
                fact_check_markdown = _build_fact_check_fallback_markdown(ui_locale=ui_locale)

        normalized_payload["summary_markdown"] = summary_markdown
        normalized_payload["fact_check_markdown"] = fact_check_markdown if fact_check_enabled else ""
        return json.dumps(normalized_payload, ensure_ascii=False)

    except Exception as e:
        msg = str(e)
        if "Model does not exist" in msg or "code': 20012" in msg:
            msg += "\n\n提示：你使用的 Base URL 可能不支持该模型，请在下拉框选择正确的模型（如 deepseek-chat）或检查 API 文档。"
        elif "Incorrect API key" in msg or "code': 401" in msg:
             msg += "\n\n提示：API Key 错误，请检查是否填写正确。"
        return f"总结失败：{msg}"


def get_video_transcript(
    api: YouTubeTranscriptApi,
    video_id: str,
    video_url: str,
    languages: list[str] | None = None,
) -> str:
    asr_force_cpu = bool(getattr(api, "_asr_force_cpu", False))
    
    # 1. YouTube 逻辑
    langs = expand_languages(languages or ["en"])
    disable_audio_transcribe_override = getattr(api, "_disable_audio_transcribe_override", None)
    if disable_audio_transcribe_override is None:
        disable_audio_transcribe = str(
            os.environ.get("DISABLE_AUDIO_TRANSCRIBE", "1") or "1"
        ).strip().lower() not in {"0", "false", "no"}
    else:
        disable_audio_transcribe = bool(disable_audio_transcribe_override)
    asr_enabled = bool(getattr(api, "_asr_enabled", False)) and not disable_audio_transcribe
    asr_model = str(getattr(api, "_asr_model", "") or "")
    asr_language = str(getattr(api, "_asr_language", "") or "")
    asr_fast_mode = bool(getattr(api, "_asr_fast_mode", False))
    status_cb = getattr(api, "_status_callback", None)
    local_fetch_node_mode = bool(str(os.environ.get("LOCAL_FETCH_NODE_MODE", "") or "").strip())
    running_on_render = bool(str(os.environ.get("RENDER_SERVICE_ID", "") or "").strip())
    disable_render_asr_fallback = running_on_render and str(
        os.environ.get("REMOTE_TRANSCRIBE_DISABLE_RENDER_ASR_FALLBACK", "0") or "0"
    ).strip().lower() not in {"0", "false", "no"}
    local_skip_transcript_api = (
        local_fetch_node_mode
        and str(os.environ.get("LOCAL_FETCH_SKIP_TRANSCRIPT_API", "1") or "1").strip().lower() not in {"0", "false", "no"}
    )
    remote_worker_mode = str(os.environ.get("REMOTE_TRANSCRIBE_MODE", "") or "").strip().lower()
    remote_worker_enabled = str(
        os.environ.get("REMOTE_TRANSCRIBE_ENABLED", "0") or "0"
    ).strip().lower() in {"1", "true", "yes"}
    prefer_remote_first = remote_worker_enabled and remote_worker_mode in {"prefer_remote", "remote_first", "force_remote"}
    remote_worker_summary = "disabled"
    try:
        if prefer_remote_first and not str(os.environ.get("LOCAL_FETCH_NODE_MODE", "") or "").strip():
            remote_worker_url = str(os.environ.get("REMOTE_TRANSCRIBE_URL", "") or "").strip()
            if remote_worker_url:
                try:
                    if callable(status_cb):
                        status_cb("优先调用本地抓取节点")
                    remote_text = try_fetch_transcript_via_remote_worker(
                        video_id=video_id,
                        video_url=video_url,
                        languages=langs,
                        api=api,
                    )
                    if remote_text and not is_html_like_text(remote_text):
                        return remote_text
                except Exception as remote_exc:
                    remote_worker_summary = f"{type(remote_exc).__name__}: {remote_exc}"
                    if callable(status_cb):
                        status_cb(f"本地抓取节点失败，回退 Render 直抓: {type(remote_exc).__name__}")
            else:
                remote_worker_summary = "not-configured"
        if local_skip_transcript_api:
            if callable(status_cb):
                status_cb(f"本地节点跳过 Transcript API，直接尝试 yt-dlp/Whisper: {video_id}")
            raise RequestBlocked("本地节点跳过 Transcript API")
        if callable(status_cb):
            status_cb(f"尝试 YouTube Transcript API: {video_id}")
        transcript = api.fetch(video_id, languages=langs)
        content = "\n".join([entry.text for entry in transcript])
        if is_html_like_text(content):
            raise RequestBlocked("返回 HTML 页面源码")
        header = f"[{transcript.language_code} | {'自动' if transcript.is_generated else '人工'}] 片段数: {len(transcript)}"
        return header + "\n\n" + content
    except NoTranscriptFound:
        if callable(status_cb):
            status_cb(f"Transcript API 未命中，尝试列出可用字幕: {video_id}")
        transcript_list = api.list(video_id)
        chosen = None
        for t in transcript_list:
            chosen = t
            if not t.is_generated:
                break
        if chosen is None:
            raise
        transcript = chosen.fetch()
        content = "\n".join([entry.text for entry in transcript])
        if is_html_like_text(content):
            raise RequestBlocked("返回 HTML 页面源码")
        header = f"[{transcript.language_code} | {'自动' if transcript.is_generated else '人工'}] 片段数: {len(transcript)}（已忽略语言优先级：{','.join(langs)}）"
        return header + "\n\n" + content
    except (TranscriptsDisabled, PoTokenRequired, RequestBlocked, IpBlocked, requests.exceptions.RequestException) as e:
        proxy_url = str(getattr(api, "_effective_proxy", "") or "")
        timeout_seconds = float(getattr(api, "_timeout_seconds", 60.0) or 60.0)
        retries = int(getattr(api, "_retries", 2) or 2)
        cookies_from_browser = str(getattr(api, "_cookies_from_browser", "") or "")
        cookie_resolve_error = ""
        try:
            cookies_file = resolve_cookie_file(
                cookies_file=str(getattr(api, "_cookies_file", "") or ""),
                cookies_content=str(getattr(api, "_cookies_content", "") or ""),
                cookies_content_b64=str(getattr(api, "_cookies_content_b64", "") or ""),
            )
        except Exception as resolve_exc:
            cookie_resolve_error = f"{type(resolve_exc).__name__}:{resolve_exc}"
            cookies_file = ""
        cookie_debug_summary = build_cookie_runtime_diagnostics(
            api,
            cookies_file=str(getattr(api, "_cookies_file", "") or ""),
            cookies_from_browser=cookies_from_browser,
            resolved_cookie_file=cookies_file,
            resolve_error=cookie_resolve_error,
        )
        if not prefer_remote_first and not str(os.environ.get("LOCAL_FETCH_NODE_MODE", "") or "").strip():
            remote_worker_url = str(os.environ.get("REMOTE_TRANSCRIBE_URL", "") or "").strip()
            if remote_worker_url:
                try:
                    remote_text = try_fetch_transcript_via_remote_worker(
                        video_id=video_id,
                        video_url=video_url,
                        languages=langs,
                        api=api,
                    )
                    if remote_text and not is_html_like_text(remote_text):
                        return remote_text
                except Exception as remote_exc:
                    remote_worker_summary = f"{type(remote_exc).__name__}: {remote_exc}"
            else:
                remote_worker_summary = "not-configured"
        cookie_debug_summary = cookie_debug_summary + f"\n远程抓取诊断: {remote_worker_summary}"
        try:
            if callable(status_cb):
                status_cb(f"尝试 yt-dlp 字幕抓取: {video_id}")
            label, text = fetch_subtitles_with_ytdlp(
                video_url,
                preferred_langs=langs,
                proxy_url=proxy_url,
                timeout_seconds=timeout_seconds,
                retries=retries,
                cookies_file=cookies_file,
                cookies_from_browser=cookies_from_browser,
            )
            if text and text.strip() and not is_html_like_text(text):
                header = f"[yt-dlp | {label}]"
                return header + "\n\n" + text
        except Exception:
            pass

        if asr_enabled:
            if prefer_remote_first and remote_worker_url and not local_fetch_node_mode and disable_render_asr_fallback:
                raise RuntimeError(
                    "未获取到可用字幕，且 Render 已启用低内存保护：已禁止在 Render 本机执行音频下载与 Whisper 转写。"
                    " 请确认本地抓取节点在线并优先处理该任务；如确实需要允许 Render 兜底，可将 "
                    "`REMOTE_TRANSCRIBE_DISABLE_RENDER_ASR_FALLBACK=0`。"
                    f"\n远程抓取诊断: {remote_worker_summary}"
                )
            if callable(status_cb):
                status_cb(f"尝试音频下载与 Whisper 转写: {video_id}")
            # 传递强制CPU标志到内部函数
            setattr(transcribe_video_audio_with_ytdlp, "_force_cpu", asr_force_cpu)
            label, text = transcribe_video_audio_with_ytdlp(
                video_url=video_url,
                proxy_url=proxy_url,
                timeout_seconds=timeout_seconds,
                retries=retries,
                cookies_file=cookies_file,
                cookies_from_browser=cookies_from_browser,
                model_name=asr_model,
                language=asr_language,
                status_callback=status_cb,
                fast_mode=asr_fast_mode,
                cookie_debug_summary=cookie_debug_summary,
            )
            return f"[asr | {label}]\n\n{text}"

        if cookie_debug_summary:
            raise RuntimeError(f"{e}\nCookies 运行时诊断: {cookie_debug_summary}")
        raise e


def list_available_transcripts(api: YouTubeTranscriptApi, video_id: str) -> str:
    transcript_list = api.list(video_id)
    rows = []
    for t in transcript_list:
        kind = "自动" if t.is_generated else "人工"
        translatable = "可翻译" if t.is_translatable else "不可翻译"
        rows.append(f"{t.language_code}\t{kind}\t{translatable}\t{t.language}")
    if not rows:
        return "未检测到任何字幕轨道。"
    return "language_code\t类型\t翻译\t语言\n" + "\n".join(rows)


def format_error(e: Exception) -> str:
    raw_msg = strip_ansi(str(e))
    if isinstance(e, (InvalidVideoId, ValueError)):
        msg = "无法解析视频 ID，请检查链接或直接粘贴 11 位 ID。"
    elif isinstance(e, TranscriptsDisabled):
        msg = "该视频字幕被关闭或不可用。可换一个有字幕的视频，或后续接入“音频转写兜底”。"
    elif isinstance(e, NoTranscriptFound):
        msg = "没有匹配你设置语言优先级的字幕。可先点“检测字幕”查看可用语言代码。"
    elif isinstance(e, (VideoUnavailable, VideoUnplayable)):
        msg = "视频不可用/不可播放（可能地区限制、删除或需要登录）。"
    elif isinstance(e, AgeRestricted):
        msg = "年龄限制视频，可能需要登录或不同的抓取方式。"
    elif isinstance(e, (IpBlocked, RequestBlocked)):
        msg = "请求被 YouTube 限制/风控。建议更换网络、降低频率或配置代理。"
    elif isinstance(e, PoTokenRequired):
        msg = "YouTube 需要 PoToken，当前方式可能抓不到字幕。"
    elif isinstance(e, requests.exceptions.ProxyError):
        msg = "代理不可用，请检查代理地址格式与连通性。"
    elif isinstance(e, requests.exceptions.ConnectTimeout):
        msg = "连接超时，请增大超时或使用代理/可访问网络。"
    elif isinstance(e, requests.exceptions.ChunkedEncodingError):
        msg = "连接中途断流（常见于代理/链路不稳定）。建议增大超时、提高重试次数，或更换更稳定的代理。"
    elif isinstance(e, requests.exceptions.ConnectionError):
        msg = "网络连接失败（可能无法访问 YouTube）。可尝试代理或更换网络。"
    elif isinstance(e, RuntimeError) and "yt-dlp" in str(e):
        msg = str(e)
    elif isinstance(e, RuntimeError) and "openai-whisper" in str(e):
        msg = str(e)
    elif isinstance(e, RuntimeError) and "音频转写失败" in str(e):
        msg = str(e)
    elif isinstance(e, KeyError) and "Linear" in raw_msg:
        msg = "Whisper 模型加载失败 (KeyError: Linear)。可能是 PyTorch 版本不兼容或模型文件损坏。建议尝试删除缓存的模型文件或重装 openai-whisper。"
    elif isinstance(e, FileNotFoundError) and ("ffmpeg" in raw_msg.lower() or "winerror 2" in raw_msg.lower()):
        msg = "未安装 ffmpeg（或未加入 PATH），无法进行音频转写。请先安装 ffmpeg 后重试。"
    elif ("could not copy" in raw_msg.lower() and "cookie" in raw_msg.lower()) or ("cookie database" in raw_msg.lower()) or ("could not find" in raw_msg.lower() and "cookies" in raw_msg.lower()):
        msg = "无法读取浏览器 Cookies（数据库被占用/无权限，或未找到对应浏览器的 Cookie 数据库）。建议关闭浏览器后重试，或切换可用浏览器/关闭自动读取 Cookies。"
    elif "HTTP Error 429" in raw_msg or "too many 429" in raw_msg.lower() or "429" in raw_msg:
        msg = "触发 YouTube 429 限流/风控。建议降低频率并等待一段时间，或使用更稳定的代理；也可启用“自动读取浏览器 Cookies”。"
    else:
        msg = str(e) or e.__class__.__name__

    msg = strip_ansi(msg)
    return f"失败原因：{msg}\n\n异常类型：{e.__class__.__name__}\n\n原始错误：{strip_ansi(repr(e))}"


def get_transcript_from_input(url_or_id: str, languages_csv: str) -> tuple[str, str, str]:
    video_url = normalize_video_url(url_or_id)
    video_id = extract_video_id(video_url)
    languages = [s.strip() for s in (languages_csv or "").split(",") if s.strip()]
    return video_id, video_url, ",".join(languages or ["zh-Hans", "zh", "en"])





def search_channels(
    keyword: str,
    limit: int = 3,
    proxy_url: str = "",
    timeout_seconds: float = 10.0
) -> dict:
    """
    搜索 YouTube 频道
    返回:
    {
        "youtube": [
            {"id": "...", "name": "...", "url": "...", "avatar": "...", "desc": "...", "platform": "youtube"},
            ...
        ]
    }
    """
    import requests
    from concurrent.futures import ThreadPoolExecutor
    
    results = {"youtube": []}
    
    def _search_youtube():
        try:
            # YouTube Search via yt-dlp
            from yt_dlp import YoutubeDL
            
            opts = {
                "quiet": True,
                "no_warnings": True,
                "extract_flat": True,
                "nocheckcertificate": True,
                "ignoreerrors": True,
                "socket_timeout": timeout_seconds,
            }
            if proxy_url:
                opts["proxy"] = proxy_url
            
            # 搜索稍微多一点，以便去重
            search_query = f"ytsearch{limit*2}:{keyword}"
            
            with YoutubeDL(opts) as ydl:
                info = ydl.extract_info(search_query, download=False)
                if not info: return
                
                entries = info.get("entries", [])
                seen_channels = set()
                
                for e in entries:
                    if not e: continue
                    c_id = e.get("channel_id") or e.get("uploader_id")
                    if not c_id: continue
                    
                    if c_id in seen_channels: continue
                    seen_channels.add(c_id)
                    
                    c_name = e.get("channel") or e.get("uploader")
                    c_url = e.get("channel_url") or f"https://www.youtube.com/channel/{c_id}"
                    
                    # 尝试获取头像 (fix: flat 模式无头像，需单独获取)
                    avatar_url = ""
                    try:
                        # 复用 get_channel_info 获取头像
                        # 注意：这会增加耗时，但为了头像显示是必要的
                        _, _, _, avatar_url = get_channel_info(
                           c_url, proxy_url, timeout_seconds=5.0
                        )
                    except:
                        pass
                    
                    results["youtube"].append({
                        "id": c_id,
                        "name": c_name,
                        "url": c_url,
                        "avatar": avatar_url,
                        "desc": "",
                        "platform": "youtube"
                    })
                    
                    if len(results["youtube"]) >= limit:
                        break
                        
        except Exception as e:
            # print(f"YouTube search failed: {e}")
            pass

    # 并发执行
    _search_youtube()
    return results


def get_channel_info(
    channel_url: str,
    proxy_url: str,
    timeout_seconds: float = 20.0,
    retries: int = 2,
    cookies_file: str = "",
    cookies_from_browser: str = "",
) -> tuple[str, str, str, str, str]:
    """
    获取频道信息
    返回: (channel_id, channel_name, canonical_url, avatar_url, platform)
    """
    from yt_dlp import YoutubeDL
    import re

    last_err: Exception | None = None
    original_input = channel_url.strip()
    url_candidates = []

    if original_input.startswith("http"):
        url_candidates.append(original_input)
    elif original_input.startswith("@"):
        url_candidates.append(f"https://www.youtube.com/{original_input}")
    else:
        # 1. 假设是 channel ID
        if re.fullmatch(r"UC[a-zA-Z0-9_-]{22}", original_input):
            url_candidates.append(f"https://www.youtube.com/channel/{original_input}")
        # 2. 尝试 ytsearch
        url_candidates.append(f"ytsearch1:{original_input}")

    for attempt in range(max(1, int(retries))):
        for candidate_url in url_candidates:
            for cookiefile, cfb in CookieManager.get_sources(cookies_file, cookies_from_browser, False):
                opts = {
                    "quiet": True,
                    "no_warnings": True,
                    "extract_flat": True,
                    "playlistend": 1, 
                    "socket_timeout": float(timeout_seconds),
                    "nocheckcertificate": True,
                    "ignoreerrors": True, 
                }
                if proxy_url:
                    opts["proxy"] = proxy_url
                if cookiefile:
                    opts["cookiefile"] = cookiefile
                elif cfb:
                    opts["cookiesfrombrowser"] = (cfb,)

                with YoutubeDL(opts) as ydl:
                    try:
                        info = ydl.extract_info(candidate_url, download=False)
                        if not info:
                            continue
                        
                        if "entries" in info:
                            entries = info.get("entries", [])
                            if not entries:
                                continue
                            info = entries[0]
                        
                        c_id = info.get("channel_id") or info.get("uploader_id")
                        c_name = info.get("channel") or info.get("uploader") or info.get("title")
                        c_url = info.get("channel_url") or info.get("uploader_url")
                        
                        c_avatar = ""
                        platform = "youtube"
                        
                        thumbnails = info.get("thumbnails")
                        if thumbnails:
                            for t in thumbnails:
                                if t.get("id") == "avatar_uncropped":
                                    c_avatar = t.get("url")
                                    break

                            if not c_avatar:
                                for t in thumbnails:
                                    u = t.get("url", "")
                                    if "ggpht.com" in u:
                                        c_avatar = u
                                        break
                            
                            if not c_avatar:
                                for t in thumbnails:
                                    w = t.get("width")
                                    h = t.get("height")
                                    if w and h and w == h:
                                        if "ytimg.com" in t.get("url", ""):
                                            continue
                                        c_avatar = t.get("url")
                                        break
                            
                            if not c_avatar and len(thumbnails) > 0:
                                candidate = thumbnails[-1].get("url")
                                if "ytimg.com" in candidate:
                                    for t in reversed(thumbnails):
                                        if "ytimg.com" not in t.get("url", ""):
                                            c_avatar = t.get("url")
                                            break
                                else:
                                    c_avatar = candidate

                        if c_id and c_name:
                            if not c_url:
                                c_url = f"https://www.youtube.com/channel/{c_id}"
                            
                            return str(c_id), str(c_name), str(c_url), str(c_avatar or ""), platform
                            
                    except Exception as e:
                        last_err = e
                        continue
    
    if last_err:
        raise RuntimeError(strip_ansi(str(last_err)))
    raise RuntimeError(f"无法找到频道信息: {original_input}")




def get_channel_recent_videos(
    channel_url: str,
    limit: int = 5,
    proxy_url: str = "",
    timeout_seconds: float = 20.0,
    retries: int = 2,
    cookies_file: str = "",
    cookies_from_browser: str = "",
    filter_longest: bool = False,
    min_duration_seconds: int = 0,
    only_streams: bool = False,
) -> list[dict]:
    """
    获取频道最新视频，支持混合扫描 /videos 和 /streams 以确保不错过长直播回放。
    
    filter_longest: 如果为 True，将采用"混合提取"策略：
      1. 快速扫描 /videos 和 /streams 获取候选列表
      2. 对前 15 个候选视频进行详细抓取(获取准确时长和日期)
      3. 在最近发布的 6 个视频中选择时长最长的一个
    
    min_duration_seconds: 最小时长限制（秒），用于过滤短视频。
    
    only_streams: 如果为 True，仅扫描 /streams 页面（针对直播回放为主的频道）。
    
    返回: list of {id, title, url, upload_date, duration}
    """
    try:
        from yt_dlp import YoutubeDL
    except ImportError:
        raise RuntimeError("未安装 yt-dlp")

    from datetime import datetime

    def _build_video_item(raw_item: dict, fallback_url: str = "") -> dict:
        """标准化视频元数据，便于统一排序与展示。"""
        timestamp = raw_item.get("timestamp") or raw_item.get("release_timestamp") or 0
        upload_date = raw_item.get("upload_date")
        if not upload_date and timestamp:
            try:
                upload_date = datetime.fromtimestamp(int(timestamp)).strftime("%Y%m%d")
            except Exception:
                upload_date = ""
        return {
            "id": raw_item.get("id"),
            "title": raw_item.get("title", "无标题"),
            "url": raw_item.get("url") or raw_item.get("webpage_url") or fallback_url,
            "upload_date": upload_date,
            "duration": raw_item.get("duration") or 0,
            "timestamp": int(timestamp) if timestamp else 0,
        }

    def _video_sort_key(item: dict) -> tuple[int, str, int]:
        """优先按精确时间排序，缺失时回退到 upload_date 和时长。"""
        upload_date = str(item.get("upload_date") or "")
        return (
            int(item.get("timestamp") or 0),
            upload_date,
            int(item.get("duration") or 0),
        )

    base_url = channel_url.strip().rstrip("/")
    # YouTube 逻辑
    # 初始化 targets 列表
    targets = []
    
    if only_streams:
        # 强制仅扫描直播回放页面
        targets.append(base_url + "/streams")
    # 如果明确指定了 tab，就只扫那个；否则扫 videos 和 streams
    elif any(x in base_url for x in ["/videos", "/shorts", "/streams", "/live", "/featured"]):
        targets.append(base_url)
    else:
        targets.append(base_url + "/videos")
        targets.append(base_url + "/streams")

    candidates_map = {}  # id -> item

    for attempt in range(max(1, int(retries))):
        for cookiefile, cfb in CookieManager.get_sources(cookies_file, cookies_from_browser, False):
            # 第一步：快速扫描 (extract_flat=True)
            for target_url in targets:
                opts = {
                    "quiet": True,
                    "no_warnings": True,
                    "extract_flat": True,
                    "playlistend": 15,  # 扩大扫描范围到前 15 个，确保即便有一堆 Shorts 也能扫到正片
                    "socket_timeout": float(timeout_seconds),
                    "nocheckcertificate": True,
                    "ignoreerrors": True,
                }
                if proxy_url:
                    opts["proxy"] = proxy_url
                if cookiefile:
                    opts["cookiefile"] = cookiefile
                elif cfb:
                    opts["cookiesfrombrowser"] = (cfb,)

                with YoutubeDL(opts) as ydl:
                    try:
                        info = ydl.extract_info(target_url, download=False)
                        if not info:
                            continue
                        
                        entries = info.get("entries", [])
                        for e in entries:
                            if not e: continue
                            v_id = e.get("id")
                            if not v_id: continue
                            
                            # 记录基础信息
                            if v_id not in candidates_map:
                                candidates_map[v_id] = _build_video_item(
                                    e,
                                    fallback_url=f"https://www.youtube.com/watch?v={v_id}",
                                )
                    except Exception:
                        pass
            
            if candidates_map:
                break
        if candidates_map:
            break

    if not candidates_map:
        return []

    # 转为列表
    all_candidates = list(candidates_map.values())

    # 如果不需要过滤长视频，仍然必须先按发布时间排序，不能依赖抓取顺序。
    if not filter_longest:
        all_candidates.sort(key=_video_sort_key, reverse=True)
        return all_candidates[: max(1, int(limit or 1))]

    # === filter_longest 混合模式逻辑 ===
    # 优化：如果 all_candidates 中已经有 duration 和 upload_date，则不需要 fetch_detail
    # 只有当 duration 为 0/None 或者需要精确过滤时才 fetch_detail
    
    # 预过滤：如果 flat 模式已经拿到了时长，且明显小于阈值，直接排除
    # 这样可以避免对 Shorts 发起无意义的 detail 请求，大幅提升速度
    threshold_duration = min_duration_seconds if min_duration_seconds > 0 else 180
    
    check_list = []
    for item in all_candidates:
        flat_dur = item.get("duration")
        # 如果时长已知且小于阈值，直接跳过 (Shorts 过滤)
        if flat_dur and flat_dur > 0 and flat_dur < threshold_duration:
            continue
        check_list.append(item)

    detailed_results = []
    
    # 需要 fetch 的列表
    need_fetch = []
    # 已经可以直接用的列表
    ready_results = []
    
    for item in check_list:
        # 如果是 YouTube，flat 模式通常有 duration
        # 但 upload_date 有时是 None 或者 'NA'
        # 如果 duration > 0 且有 upload_date，直接用
        
        has_dur = item.get("duration") and item.get("duration") > 0
        has_date = item.get("upload_date")
        
        # 如果有了 duration 和 date，直接复用，不重新抓取
        if has_dur and has_date:
            ready_results.append(item)
        else:
            need_fetch.append(item)
            
    # 如果所有都需要 fetch，或者需要更精确的信息
    # 为了保险起见，如果 need_fetch 很多，我们只取前 10 个去 fetch
    if need_fetch:
        need_fetch = need_fetch[:10]

    def fetch_detail_item(item, proxy_url_val):
        try:
            # 缩短单个视频的超时时间，因为如果太慢通常是网络问题
            local_opts = {
                "quiet": True,
                "no_warnings": True,
                "extract_flat": False,
                "socket_timeout": 5, # 再次缩短超时到 5s
                "nocheckcertificate": True,
                "ignoreerrors": True,
            }
            if proxy_url_val:
                local_opts["proxy"] = proxy_url_val
                
            with YoutubeDL(local_opts) as ydl:
                info = ydl.extract_info(item["url"], download=False)
                if info:
                    return _build_video_item(info, fallback_url=item["url"])
        except Exception:
            pass
        # 失败则返回原始数据
        return item

    if need_fetch:
        from concurrent.futures import ThreadPoolExecutor, as_completed

        # 限制并发数，避免瞬间太多请求
        # 提高内部并发数到 10，加快单个频道的处理速度
        with ThreadPoolExecutor(max_workers=10) as executor:
            futures = [executor.submit(fetch_detail_item, item, proxy_url) for item in need_fetch]
            for future in as_completed(futures):
                res = future.result()
                detailed_results.append(res)
    
    # 合并结果
    detailed_results.extend(ready_results)
            
    # 过滤掉获取不到日期的（极少情况）
    valid_results = [x for x in detailed_results if x.get("upload_date") or x.get("timestamp")]

    # 基础时长过滤：默认 180s (3分钟)，如果指定了 min_duration_seconds 则使用指定值
    # 这一步更严格地排除掉 Shorts/剪辑片段/预告片等干扰项，只保留正片
    threshold_duration = min_duration_seconds if min_duration_seconds > 0 else 180
    valid_results = [x for x in valid_results if x.get("duration", 0) >= threshold_duration]

    # 按日期降序排序 (最新的在前)
    # 注意：yt-dlp 返回的日期格式是 YYYYMMDD，字符串比较即可
    valid_results.sort(key=_video_sort_key, reverse=True)

    if not valid_results:
        # 如果过滤完没视频了，尝试放宽标准找一个最长的，避免完全空
        # 比如有些正片可能只有 2 分钟，或者全是短片时至少给一个
        detailed_results.sort(
            key=lambda x: (_video_sort_key(x), int(x.get("duration") or 0)),
            reverse=True,
        )
        if detailed_results:
             return [detailed_results[0]]
        return []

    # 获取今天的日期字符串
    from datetime import datetime
    today_str = datetime.now().strftime("%Y%m%d")

    # 逻辑回退到"简单而健壮"的版本，但增加 scan 深度。
    # 用户反馈之前的版本更好，说明复杂的"今天/昨天"逻辑可能有误判（特别是跨时区时）。
    # 最朴素的需求：
    # 1. 必须是"最近"发布的。
    # 2. 尽量是"正片"（长视频）。
    
    # 我们已经有了 valid_results，它是按 upload_date 降序排列的。
    # 并且已经过滤掉了 < 3分钟的短片。
    
    # 直接取前 3 个。
    # 为什么这样是对的？
    # 1. 如果今天发了 3 个视频，都是长视频 -> 取前3个 -> 正确。
    # 2. 如果今天发了 1 个长视频，昨天发了 2 个 -> 取前3个 -> 正确（包含了最新的）。
    # 3. 如果今天没发，昨天发了 1 个 -> 取前3个 -> 正确（包含了最新的）。
    
    # 唯一的问题：如果今天发了 5 个视频，前 3 个是短评(5min)，第 4 个是正片(30min)。
    # 按时间排序会取前 3 个短评，漏掉正片。
    # 但用户同时也抱怨"不是最新时间的"。
    # 这说明"时间"的优先级 > "时长"。
    # 如果为了找长视频而去取旧视频，用户会不满意。
    
    # 折中方案：
    # 在最近的 N 个（比如 5 个）候选视频中，优先展示。
    # 如果我们直接返回前 3 个，这是最符合"Timeline"逻辑的。
    # 用户说"还不如之前的"，之前的逻辑其实是"取最近10个 -> 加权打分 -> 取1个"。
    # 现在用户要"取3个"。
    
    # 让我们尝试最直观的逻辑：
    # 返回【最近发布的】且【时长合格】的视频，最多 3 个。
    # 不做任何额外的日期分组或时长重排序，完全信任时间轴。
    
    return valid_results[: max(1, int(limit or 1))]


def check_network(proxy_url: str, timeout: float = 5.0) -> str | None:
    try:
        proxies = {"http": proxy_url, "https": proxy_url} if proxy_url else None
        
        def _do_check(verify: bool) -> None:
            s = requests.Session()
            s.trust_env = False
            if proxies:
                s.proxies.update(proxies)
            # 先试 Google
            try:
                r = s.get("https://www.google.com", timeout=timeout, stream=True, verify=verify)
                r.close()
                return
            except Exception:
                pass
            # 再试 YouTube
            r = s.get("https://www.youtube.com", timeout=timeout, stream=True, verify=verify)
            r.close()

        try:
            _do_check(verify=True)
            return None
        except Exception as e1:
            # 如果第一次失败，尝试关闭 SSL 验证（针对某些中间人代理）
            try:
                _do_check(verify=False)
                return None
            except Exception:
                # 如果还是失败，返回第一次的错误
                return f"{e1.__class__.__name__}: {e1}"
    except Exception as e:
        return f"{e.__class__.__name__}: {e}"
